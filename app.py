# =============================================================================
# app.py — CRM Asset Management  v15.0 — Universal Export Hub
# Architecture : Full Modal (@st.dialog) + Split-Screen CRM
# Session-state keys are namespaced per feature to prevent cross-tab conflicts
# =============================================================================

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from datetime import date, timedelta
import io
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from streamlit_agraph import agraph, Node, Edge, Config

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import database as db
import pdf_generator as pdf_gen

# ---------------------------------------------------------------------------
# CONSTANTES
# ---------------------------------------------------------------------------
MARINE  = "#001c4b"
CIEL    = "#019ee1"
ORANGE  = "#f07d00"
BLANC   = "#ffffff"
GRIS    = "#e8e8e8"
GTXT    = "#444444"
B_MID   = "#1a5e8a"
B_PAL   = "#4a8fbd"
B_DEP   = "#003f7a"
PALETTE = [B_MID, MARINE, B_PAL, B_DEP, "#2c7fb8",
           "#004f8c", "#6baed6", "#08519c", "#9ecae1", "#003060"]

TYPES_CLIENT      = ["IFA", "Wholesale", "Instit", "Family Office"]
REGIONS           = db.REGIONS_REFERENTIEL
FONDS             = db.FONDS_REFERENTIEL
STATUTS           = ["Prospect","Initial Pitch","Due Diligence",
                     "Soft Commit","Funded","Paused","Lost","Redeemed"]
STATUTS_ACTIFS    = ["Prospect","Initial Pitch","Due Diligence","Soft Commit"]
RAISONS_PERTE     = ["Pricing","Track Record","Macro","Competitor","Autre"]
TYPES_INTERACTION = ["Call","Meeting","Email","Roadshow","Conference","Autre"]

COUNTRIES_LIST = [
    "", "United Arab Emirates", "Saudi Arabia", "Qatar", "Kuwait", "Bahrain", "Oman",
    "United Kingdom", "France", "Germany", "Switzerland", "Luxembourg", "Netherlands",
    "Italy", "Spain", "Belgium", "Austria", "Sweden", "Norway", "Denmark", "Finland",
    "Singapore", "Japan", "Hong Kong", "China", "South Korea", "Australia", "India",
    "United States", "Canada", "Brazil", "Mexico", "South Africa", "Egypt",
]

STATUT_COLORS = {
    "Funded":        B_MID,    "Soft Commit":   "#2c7fb8",
    "Due Diligence": "#004f8c", "Initial Pitch": B_PAL,
    "Prospect":      "#9ecae1", "Lost":          "#aaaaaa",
    "Paused":        "#c0c0c0", "Redeemed":      "#b8b8d0",
}

fmt_m = db.format_finance

# ---------------------------------------------------------------------------
# PPTX — Account Review Generator (v14.0)
# ---------------------------------------------------------------------------
def _rgb(hex_str):
    """Convert #RRGGBB to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def generate_account_review_pptx(client_data: dict, group_summary: dict,
                                   contacts_df, activities_df) -> bytes:
    """
    Generates a 2-slide PPTX Account Review.
    Slide 1: Title — Revue de Compte Stratégique
    Slide 2: Synthèse — AUM, KYC, Fonds, Next Actions
    Returns bytes ready for st.download_button.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    MARINE_RGB  = _rgb("#001c4b")
    CIEL_RGB    = _rgb("#019ee1")
    ORANGE_RGB  = _rgb("#f07d00")
    BLANC_RGB   = _rgb("#ffffff")
    LIGHT_RGB   = _rgb("#f2f6fa")
    GREY_RGB    = _rgb("#666666")

    blank_layout = prs.slide_layouts[6]  # Completely blank

    def _add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_width_pt=0):
        from pptx.util import Pt as _Pt
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if line_rgb:
            shape.line.color.rgb = line_rgb
            shape.line.width = Pt(line_width_pt)
        else:
            shape.line.fill.background()
        return shape

    def _add_text(slide, text, x, y, w, h, font_size, bold=False, color_rgb=None,
                  align=PP_ALIGN.LEFT, italic=False, wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h))
        txBox.word_wrap = wrap
        tf = txBox.text_frame
        tf.word_wrap = wrap
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color_rgb:
            run.font.color.rgb = color_rgb
        return txBox

    # ── SLIDE 1 — TITLE ──────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)

    # Dark header bar
    _add_rect(slide1, 0, 0, 13.33, 2.8, MARINE_RGB)
    # Accent stripe
    _add_rect(slide1, 0, 2.8, 13.33, 0.08, CIEL_RGB)

    # Company name
    _add_text(slide1,
              "REVUE DE COMPTE STRATÉGIQUE",
              0.6, 0.35, 12.0, 0.7,
              font_size=13, bold=False, color_rgb=CIEL_RGB,
              align=PP_ALIGN.LEFT)
    _add_text(slide1,
              client_data.get("nom_client", "—"),
              0.6, 0.95, 12.0, 1.2,
              font_size=40, bold=True, color_rgb=BLANC_RGB,
              align=PP_ALIGN.LEFT)

    # Date
    _add_text(slide1,
              date.today().strftime("%d %B %Y").upper(),
              0.6, 2.2, 6.0, 0.5,
              font_size=11, bold=False, color_rgb=_rgb("#7ab8d8"),
              align=PP_ALIGN.LEFT)

    # KYC badge (bottom-right of header)
    kyc = client_data.get("kyc_status", "En cours")
    kyc_colors = {"Validé": _rgb("#22a062"), "En cours": ORANGE_RGB, "Bloqué": _rgb("#c0392b")}
    kyc_rgb = kyc_colors.get(kyc, GREY_RGB)
    _add_rect(slide1, 10.0, 1.8, 2.7, 0.7, kyc_rgb)
    _add_text(slide1, "KYC : {}".format(kyc),
              10.0, 1.9, 2.7, 0.5,
              font_size=14, bold=True, color_rgb=BLANC_RGB,
              align=PP_ALIGN.CENTER)

    # Info grid below header
    infos = [
        ("TYPE", client_data.get("type_client","—")),
        ("RÉGION", client_data.get("region","—")),
        ("TIER", client_data.get("tier","—")),
        ("COUNTRY", client_data.get("country","—") or "—"),
    ]
    for idx, (lbl, val) in enumerate(infos):
        xi = 0.5 + idx * 3.2
        _add_rect(slide1, xi, 3.1, 3.0, 1.2, LIGHT_RGB)
        _add_text(slide1, lbl, xi+0.15, 3.2, 2.7, 0.4,
                  font_size=9, bold=False, color_rgb=GREY_RGB)
        _add_text(slide1, val, xi+0.15, 3.55, 2.7, 0.6,
                  font_size=16, bold=True, color_rgb=MARINE_RGB)

    # Product interests
    prods = [p.strip() for p in str(client_data.get("product_interests","")).split(",") if p.strip()]
    if prods:
        _add_text(slide1, "Product Interests : " + "  ·  ".join(prods),
                  0.5, 4.5, 12.3, 0.5,
                  font_size=11, bold=False, color_rgb=GREY_RGB)

    # Footer
    _add_rect(slide1, 0, 7.1, 13.33, 0.4, _rgb("#f0f4f8"))
    _add_text(slide1, "Document a usage interne — Confidentiel — Asset Management CRM",
              0.5, 7.15, 12.0, 0.3,
              font_size=8, bold=False, color_rgb=GREY_RGB, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 — SYNTHÈSE ───────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)

    # Thin top bar
    _add_rect(slide2, 0, 0, 13.33, 0.55, MARINE_RGB)
    _add_text(slide2,
              "SYNTHÈSE  —  {}".format(client_data.get("nom_client","").upper()),
              0.4, 0.1, 10.0, 0.38,
              font_size=11, bold=True, color_rgb=BLANC_RGB)
    _add_text(slide2, date.today().strftime("%d/%m/%Y"),
              11.0, 0.1, 2.0, 0.38,
              font_size=11, bold=False, color_rgb=_rgb("#7ab8d8"),
              align=PP_ALIGN.RIGHT)

    # ── KPI cards row ────────────────────────────────────────────────────────
    kpi_data = [
        ("AUM Consolidé",  db.format_finance(group_summary["aum_consolide"]), MARINE_RGB),
        ("AUM Direct",     db.format_finance(group_summary["aum_direct"]),    CIEL_RGB),
        ("Fonds Investis", str(len(group_summary["fonds_investis"])),          _rgb("#1a5e8a")),
        ("Deals Actifs",   str(len(group_summary["next_actions"])),            ORANGE_RGB),
    ]
    for idx, (lbl, val, clr) in enumerate(kpi_data):
        xi = 0.3 + idx * 3.2
        _add_rect(slide2, xi, 0.75, 3.0, 1.3, MARINE_RGB)
        _add_rect(slide2, xi, 0.75, 3.0, 0.06, clr)  # color accent top
        _add_text(slide2, lbl, xi+0.15, 0.85, 2.7, 0.4,
                  font_size=9, bold=False, color_rgb=_rgb("#7ab8d8"))
        _add_text(slide2, val, xi+0.15, 1.2, 2.7, 0.65,
                  font_size=22, bold=True, color_rgb=BLANC_RGB)

    # ── Fonds Investis ───────────────────────────────────────────────────────
    _add_text(slide2, "FONDS INVESTIS", 0.3, 2.25, 5.5, 0.35,
              font_size=9, bold=True, color_rgb=CIEL_RGB)
    _add_rect(slide2, 0.3, 2.6, 5.8, 0.04, CIEL_RGB)
    if group_summary["fonds_investis"]:
        for fi, fonds_nm in enumerate(group_summary["fonds_investis"]):
            _add_rect(slide2, 0.3, 2.72 + fi*0.42, 5.8, 0.38, LIGHT_RGB)
            _add_text(slide2, fonds_nm, 0.5, 2.74 + fi*0.42, 5.4, 0.36,
                      font_size=12, bold=False, color_rgb=MARINE_RGB)
    else:
        _add_text(slide2, "Aucun fonds financé", 0.3, 2.72, 5.8, 0.4,
                  font_size=11, bold=False, color_rgb=GREY_RGB)

    # ── Prochaines Actions ───────────────────────────────────────────────────
    _add_text(slide2, "PROCHAINES ACTIONS", 6.8, 2.25, 6.2, 0.35,
              font_size=9, bold=True, color_rgb=ORANGE_RGB)
    _add_rect(slide2, 6.8, 2.6, 6.2, 0.04, ORANGE_RGB)
    if group_summary["next_actions"]:
        for ai, act in enumerate(group_summary["next_actions"][:4]):
            row_y = 2.72 + ai * 0.55
            _add_rect(slide2, 6.8, row_y, 6.2, 0.5, LIGHT_RGB)
            label = "{} — {}  |  {}  |  {}".format(
                act["fonds"], act["statut"],
                db.format_finance(act["aum_pipeline"]),
                act["nad"])
            _add_text(slide2, label, 7.0, row_y + 0.07, 5.8, 0.4,
                      font_size=10, bold=False, color_rgb=MARINE_RGB)
    else:
        _add_text(slide2, "Aucune action planifiée", 6.8, 2.72, 6.2, 0.4,
                  font_size=11, bold=False, color_rgb=GREY_RGB)

    # ── Dernières Activités ──────────────────────────────────────────────────
    _add_text(slide2, "DERNIÈRES ACTIVITÉS", 0.3, 5.3, 6.5, 0.35,
              font_size=9, bold=True, color_rgb=MARINE_RGB)
    _add_rect(slide2, 0.3, 5.65, 12.7, 0.03, _rgb("#e8e8e8"))
    TYPE_COLORS_PPTX = {
        "Call": CIEL_RGB, "Meeting": _rgb("#1a5e8a"),
        "Email": _rgb("#4a8fbd"), "Autre": GREY_RGB,
    }
    if not activities_df.empty:
        for ri, (_, act_row) in enumerate(activities_df.head(2).iterrows()):
            row_y = 5.75 + ri * 0.5
            a_typ = str(act_row.get("type_interaction",""))
            a_col = TYPE_COLORS_PPTX.get(a_typ, GREY_RGB)
            _add_rect(slide2, 0.3, row_y, 0.06, 0.38, a_col)
            _add_text(slide2,
                      "[{}]  {}  —  {}".format(
                          a_typ,
                          str(act_row.get("date","")),
                          str(act_row.get("notes",""))[:80]),
                      0.5, row_y + 0.02, 12.3, 0.36,
                      font_size=10, bold=False, color_rgb=MARINE_RGB)
    else:
        _add_text(slide2, "Aucune activité enregistrée.", 0.3, 5.75, 12.7, 0.4,
                  font_size=11, bold=False, color_rgb=GREY_RGB)

    # Footer
    _add_rect(slide2, 0, 7.1, 13.33, 0.4, _rgb("#f0f4f8"))
    _add_text(slide2, "Document a usage interne — Confidentiel — Asset Management CRM",
              0.5, 7.15, 12.0, 0.3,
              font_size=8, bold=False, color_rgb=GREY_RGB, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX — Global Portfolio Report Generator (v15.0)
# ---------------------------------------------------------------------------
def generate_global_pptx(kpis: dict, pipeline_df, mode_comex: bool = False) -> bytes:
    """
    2-slide global PPTX report:
      Slide 1 : Title & KPI Overview
      Slide 2 : Top Deals table (client names masked if mode_comex)
    Returns bytes ready for st.download_button.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    MARINE_RGB  = _rgb("#001c4b")
    CIEL_RGB    = _rgb("#019ee1")
    ORANGE_RGB  = _rgb("#f07d00")
    BLANC_RGB   = _rgb("#ffffff")
    LIGHT_RGB   = _rgb("#f2f6fa")
    GREY_RGB    = _rgb("#666666")

    blank_layout = prs.slide_layouts[6]

    def _rect(slide, x, y, w, h, fill_rgb):
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
        sh.line.fill.background()
        return sh

    def _text(slide, text, x, y, w, h, fs, bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(fs); r.font.bold = bold
        if color: r.font.color.rgb = color
        return tb

    # ── SLIDE 1 — TITRE & KPIs ───────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    _rect(s1, 0, 0, 13.33, 2.4, MARINE_RGB)
    _rect(s1, 0, 2.4, 13.33, 0.07, CIEL_RGB)

    label_comex = " [MODE COMEX — ANONYMISÉ]" if mode_comex else ""
    _text(s1, "RAPPORT GLOBAL — PORTFOLIO PIPELINE" + label_comex,
          0.6, 0.3, 12.0, 0.6, 11, color=CIEL_RGB)
    _text(s1, "Asset Management CRM",
          0.6, 0.85, 12.0, 0.9, 34, bold=True, color=BLANC_RGB)
    _text(s1, date.today().strftime("%d %B %Y").upper(),
          0.6, 1.8, 6.0, 0.5, 11, color=_rgb("#7ab8d8"))

    # KPI row
    kpi_items = [
        ("AUM Financé Total",  fmt_m(kpis.get("total_funded", 0)),    CIEL_RGB),
        ("Pipeline Actif",     fmt_m(kpis.get("pipeline_actif", 0)),   _rgb("#1a5e8a")),
        ("Pipeline Pondéré",   fmt_m(kpis.get("weighted_pipeline",0)), ORANGE_RGB),
        ("Taux Conversion",    "{:.1f}%".format(kpis.get("taux_conversion",0)), _rgb("#22a062")),
    ]
    for idx, (lbl, val, clr) in enumerate(kpi_items):
        xi = 0.4 + idx * 3.2
        _rect(s1, xi, 2.65, 3.0, 1.4, MARINE_RGB)
        _rect(s1, xi, 2.65, 3.0, 0.07, clr)
        _text(s1, lbl, xi+0.15, 2.78, 2.7, 0.4, 9, color=_rgb("#7ab8d8"))
        _text(s1, val, xi+0.15, 3.1,  2.7, 0.75, 22, bold=True, color=BLANC_RGB)

    # Statut pills
    statut_rep = kpis.get("statut_repartition", {})
    statut_items = [(s, statut_rep.get(s, 0)) for s in
                    ["Funded","Soft Commit","Due Diligence","Initial Pitch","Prospect"]
                    if statut_rep.get(s, 0) > 0]
    _text(s1, "RÉPARTITION PAR STATUT", 0.4, 4.25, 12.5, 0.35, 8, bold=True, color=GREY_RGB)
    for si, (stat, cnt) in enumerate(statut_items):
        xi = 0.4 + si * 2.55
        _rect(s1, xi, 4.65, 2.35, 0.9, LIGHT_RGB)
        _text(s1, stat[:12], xi+0.1, 4.72, 2.15, 0.35, 9, color=GREY_RGB)
        _text(s1, str(cnt), xi+0.1, 5.0, 2.15, 0.45, 18, bold=True, color=MARINE_RGB)

    # Footer
    _rect(s1, 0, 7.1, 13.33, 0.4, LIGHT_RGB)
    _text(s1, "Document a usage interne — Confidentiel — Asset Management CRM",
          0.5, 7.15, 12.0, 0.3, 8, color=GREY_RGB, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 — TOP DEALS ──────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _rect(s2, 0, 0, 13.33, 0.55, MARINE_RGB)
    _text(s2, "TOP DEALS — AUM FINANCÉ", 0.4, 0.1, 10.0, 0.38, 11, bold=True, color=BLANC_RGB)
    _text(s2, date.today().strftime("%d/%m/%Y"), 11.0, 0.1, 2.0, 0.38,
          11, color=_rgb("#7ab8d8"), align=PP_ALIGN.RIGHT)

    if mode_comex:
        _text(s2, "Noms de clients masqués — Mode Comex",
              0.4, 0.65, 12.5, 0.35, 9, color=ORANGE_RGB)

    # Table header
    headers = ["Rang", "Client", "Fonds", "Type", "Région", "AUM Financé", "Commercial"]
    col_widths = [0.6, 2.8, 2.0, 1.4, 1.2, 1.5, 1.8]
    col_x = [0.3]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    header_y = 1.05
    for hi, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        _rect(s2, cx, header_y, cw, 0.4, MARINE_RGB)
        _text(s2, hdr, cx+0.05, header_y+0.06, cw-0.1, 0.28,
              8, bold=True, color=BLANC_RGB)

    # Table rows
    df_funded = pipeline_df[pipeline_df["statut"] == "Funded"].sort_values(
        "funded_aum", ascending=False).head(8) if not pipeline_df.empty else pd.DataFrame()

    for ri, (_, row) in enumerate(df_funded.iterrows()):
        ry = header_y + 0.4 + ri * 0.52
        bg = LIGHT_RGB if ri % 2 == 0 else BLANC_RGB
        client_name = ("Client {:02d}".format(ri+1) if mode_comex
                       else str(row.get("nom_client","—"))[:22])
        row_vals = [
            str(ri + 1),
            client_name,
            str(row.get("fonds","—"))[:16],
            str(row.get("type_client","—")),
            str(row.get("region","—")),
            fmt_m(float(row.get("funded_aum", 0))),
            str(row.get("sales_owner","—"))[:14],
        ]
        for vi, (val, cx, cw) in enumerate(zip(row_vals, col_x, col_widths)):
            _rect(s2, cx, ry, cw, 0.48, bg)
            clr = CIEL_RGB if vi == 5 else MARINE_RGB
            fw  = True  if vi == 5 else False
            _text(s2, val, cx+0.05, ry+0.1, cw-0.1, 0.3, 9, bold=fw, color=clr)

    if df_funded.empty:
        _text(s2, "Aucun deal Funded enregistré.", 0.4, 1.6, 12.5, 0.4, 11, color=GREY_RGB)

    _rect(s2, 0, 7.1, 13.33, 0.4, LIGHT_RGB)
    _text(s2, "Document a usage interne — Confidentiel — Asset Management CRM",
          0.5, 7.15, 12.0, 0.3, 8, color=GREY_RGB, align=PP_ALIGN.CENTER)

    # ── SLIDE 3+ — ANALYSES GRAPHIQUES VIA PLOTLY/KALEIDO ────────────────────
    # Génère les figures Plotly du Dashboard et les insère en PNG haute définition.
    def _add_plotly_slide(prs, fig, title_text, subtitle_text="", legend_items=None):
        """
        Convertit une figure Plotly en PNG via kaleido et l'insère dans une slide.
        - Purifie le graphique AVANT to_image (supprime titre, légende, marges).
        - Restitue le titre et la légende NATIVEMENT via des textbox PowerPoint.
        """
        s = prs.slides.add_slide(blank_layout)
        # ── Header institutionnel ──
        _rect(s, 0, 0, 13.33, 0.55, MARINE_RGB)
        _text(s, title_text, 0.4, 0.1, 10.0, 0.38, 11, bold=True, color=BLANC_RGB)
        _text(s, date.today().strftime("%d/%m/%Y"), 11.0, 0.1, 2.0, 0.38,
              11, color=_rgb("#7ab8d8"), align=PP_ALIGN.RIGHT)
        # ── Sous-titre / légende textuelle ──
        _img_top = 0.62
        if subtitle_text:
            _text(s, subtitle_text, 0.4, 0.58, 12.5, 0.28, 8, color=GREY_RGB)
            _img_top = 0.90
        # ── Légende native PPT (si fournie) ──
        if legend_items:
            _leg_x = 0.35
            for _leg_lbl, _leg_hex in legend_items[:8]:
                _rect(s, _leg_x, _img_top, 0.14, 0.14, _rgb(_leg_hex))
                _text(s, _leg_lbl, _leg_x + 0.18, _img_top - 0.01, 1.4, 0.18,
                      7, color=GREY_RGB)
                _leg_x += min(1.65, 13.0 / max(len(legend_items), 1))
            _img_top += 0.22
        # ── Footer ──
        _rect(s, 0, 7.1, 13.33, 0.4, LIGHT_RGB)
        _text(s, "Document a usage interne — Confidentiel — Asset Management CRM",
              0.5, 7.15, 12.0, 0.3, 8, color=GREY_RGB, align=PP_ALIGN.CENTER)
        # ── Image Plotly — purifiée avant export ──
        try:
            import copy as _copy
            _fig_clean = _copy.deepcopy(fig)
            _fig_clean.update_layout(
                title="",
                margin=dict(l=10, r=10, t=10, b=10),
                showlegend=False,
            )
            _img_h_in = 7.1 - _img_top - 0.05
            _img_h_px = int(_img_h_in / 7.5 * 750)
            img_bytes = _fig_clean.to_image(
                format="png", engine="kaleido", width=950, height=max(_img_h_px, 280))
            img_stream = io.BytesIO(img_bytes)
            s.shapes.add_picture(
                img_stream,
                Inches(0.3), Inches(_img_top),
                Inches(12.7), Inches(_img_h_in))
        except Exception as _e:
            _text(s, "Graphique indisponible : {}".format(str(_e)[:80]),
                  0.4, 3.5, 12.5, 0.5, 10, color=GREY_RGB)

    # ── Slide 3 : Funnel Pipeline ─────────────────────────────────────────────
    _statut_funnel_order_p = ["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit"]
    _df_funnel_p = pipeline_df[pipeline_df["statut"].isin(_statut_funnel_order_p)].copy()
    _df_funnel_p["_aum_p"] = _df_funnel_p.apply(
        lambda r: float(r["revised_aum"]) if float(r.get("revised_aum", 0) or 0) > 0
                  else float(r.get("target_aum_initial", 0) or 0), axis=1)
    _funnel_agg_p = (
        _df_funnel_p.groupby("statut")["_aum_p"].sum()
        .reindex(_statut_funnel_order_p, fill_value=0.0)
        .reset_index()
    )
    _funnel_agg_p.columns = ["statut", "aum"]
    _funnel_agg_p = _funnel_agg_p[_funnel_agg_p["aum"] > 0]
    if not _funnel_agg_p.empty:
        _fc = {"Prospect": "#9ecae1", "Initial Pitch": "#4a8fbd",
               "Due Diligence": "#004f8c", "Soft Commit": "#1a5e8a"}
        _fig_funnel_p = go.Figure(go.Funnel(
            y=_funnel_agg_p["statut"].tolist(),
            x=_funnel_agg_p["aum"].tolist(),
            textinfo="value+percent initial",
            text=["{:.1f} M\u20ac".format(v / 1_000_000) for v in _funnel_agg_p["aum"]],
            texttemplate="%{text}",
            marker_color=[_fc.get(s, "#1a5e8a") for s in _funnel_agg_p["statut"]],
            connector=dict(line=dict(color="#e8e8e8", width=1))))
        _fig_funnel_p.update_layout(
            height=420, paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font_color="#001c4b", font_size=12,
            margin=dict(l=20, r=20, t=30, b=20))
        _add_plotly_slide(prs, _fig_funnel_p, "FUNNEL PIPELINE — AUM PAR STATUT",
                          "AUM Pipeline = AUM Revise si > 0, sinon AUM Cible")

    # ── Slide 4 : AUM par Fonds et par Statut (Grouped Bar) ──────────────────
    _df_grp_p = pipeline_df[pipeline_df["statut"].isin(_statut_funnel_order_p)].copy()
    _df_grp_p["_aum_p"] = _df_grp_p.apply(
        lambda r: float(r["revised_aum"]) if float(r.get("revised_aum", 0) or 0) > 0
                  else float(r.get("target_aum_initial", 0) or 0), axis=1)
    _df_grp_p = _df_grp_p[_df_grp_p["_aum_p"] > 0]
    if not _df_grp_p.empty:
        _grp_pivot_p = _df_grp_p.groupby(["fonds", "statut"])["_aum_p"].sum().reset_index()
        _fonds_list_p = sorted(_grp_pivot_p["fonds"].unique().tolist())
        _statut_list_p = [s for s in _statut_funnel_order_p if s in _grp_pivot_p["statut"].unique()]
        _gc = {"Prospect": "#9ecae1", "Initial Pitch": "#4a8fbd",
               "Due Diligence": "#004f8c", "Soft Commit": "#1a5e8a"}
        _fig_grp_p = go.Figure()
        for _st in _statut_list_p:
            _df_st = _grp_pivot_p[_grp_pivot_p["statut"] == _st]
            _st_map = dict(zip(_df_st["fonds"], _df_st["_aum_p"]))
            _y_vals = [_st_map.get(f, 0.0) for f in _fonds_list_p]
            _fig_grp_p.add_trace(go.Bar(
                name=_st, x=_fonds_list_p, y=_y_vals,
                marker_color=_gc.get(_st, "#1a5e8a"),
                marker_line_color="#ffffff", marker_line_width=0.5))
        _fig_grp_p.update_layout(
            barmode="group", height=420,
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font_color="#001c4b", bargap=0.18, bargroupgap=0.05,
            legend_bgcolor="#ffffff", legend_bordercolor="#e8e8e8", legend_borderwidth=1,
            xaxis_showgrid=False, yaxis_showgrid=True, yaxis_gridcolor="#e8e8e8",
            margin=dict(l=20, r=20, t=30, b=20))
        _add_plotly_slide(prs, _fig_grp_p, "AUM PIPELINE PAR FONDS ET PAR STATUT",
                          "Deals actifs uniquement — horizon max")

    # ── Slide 5 : Money in Motion (fig_cf) ───────────────────────────────────
    _df_cf_p = db.get_expected_cashflows()
    if not _df_cf_p.empty and _df_cf_p["aum_pondere"].sum() > 0:
        _months_p = sorted(_df_cf_p["mois"].unique().tolist())
        _PALETTE_P = ["#1a5e8a", "#001c4b", "#4a8fbd", "#003f7a", "#2c7fb8",
                      "#004f8c", "#6baed6", "#08519c", "#9ecae1", "#003060"]
        _fig_cf_p = go.Figure()
        for _i, _fonds in enumerate(sorted(_df_cf_p["fonds"].unique())):
            _df_f = _df_cf_p[_df_cf_p["fonds"] == _fonds]
            _mm = dict(zip(_df_f["mois"], _df_f["aum_pondere"]))
            _y_vals = [_mm.get(m, 0) for m in _months_p]
            _fig_cf_p.add_trace(go.Bar(
                name=_fonds, x=_months_p, y=_y_vals,
                marker_color=_PALETTE_P[_i % len(_PALETTE_P)],
                marker_line_color="#ffffff", marker_line_width=0.5))
        # Trace totaux invisibles pour affichage
        _totaux = [sum(_mm.get(m, 0) for _df_f in
                       [_df_cf_p[_df_cf_p["fonds"] == f] for f in _df_cf_p["fonds"].unique()]
                       for _mm2 in [dict(zip(_df_f["mois"], _df_f["aum_pondere"]))])
                   for m in _months_p]
        _tot_by_month = {}
        for _, _r2 in _df_cf_p.iterrows():
            _tot_by_month[_r2["mois"]] = _tot_by_month.get(_r2["mois"], 0) + _r2["aum_pondere"]
        _totaux_vals = [_tot_by_month.get(m, 0) for m in _months_p]
        _fig_cf_p.add_trace(go.Scatter(
            x=_months_p, y=_totaux_vals,
            mode="text",
            text=["{:.1f}M".format(v / 1_000_000) if v > 0 else "" for v in _totaux_vals],
            textposition="top center",
            textfont=dict(size=9, color="#001c4b"),
            showlegend=False, hoverinfo="skip"))
        _fig_cf_p.update_layout(
            barmode="stack", height=420,
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff", font_color="#001c4b",
            legend_bgcolor="#ffffff", legend_bordercolor="#e8e8e8", legend_borderwidth=1,
            xaxis_showgrid=False, yaxis_showgrid=True, yaxis_gridcolor="#e8e8e8",
            yaxis_tickformat=".2s", yaxis_title="AUM Pondere (EUR)",
            margin=dict(l=20, r=20, t=30, b=20))
        _add_plotly_slide(prs, _fig_cf_p, "MONEY IN MOTION — PROJECTED INFLOWS",
                          "AUM pondere = AUM Pipeline x Probabilite de closing")

    # ── Slide 6 : Whitespace Heatmap ─────────────────────────────────────
    _df_ws_p = db.get_whitespace_matrix()
    if not _df_ws_p.empty:
        _ws_clients_p = _df_ws_p.index.tolist()
        if len(_ws_clients_p) > 20:
            _aum_totals_p = _df_ws_p.fillna(0).sum(axis=1).sort_values(ascending=False)
            _ws_clients_p = _aum_totals_p.head(20).index.tolist()
        _ws_fonds_p = _df_ws_p.columns.tolist()
        _ws_sub_p = _df_ws_p.loc[_ws_clients_p, _ws_fonds_p]
        _z_p = [[0.0 if (v is None or (isinstance(v, float) and np.isnan(v)))
                 else float(v) for v in row] for row in _ws_sub_p.values.tolist()]
        _txt_p = [[("{:.1f}M".format(v / 1e6) if v > 0 else "")
                   for v in row] for row in _z_p]
        _fig_ws_p = go.Figure(go.Heatmap(
            z=_z_p, x=_ws_fonds_p, y=_ws_clients_p,
            text=_txt_p, texttemplate="%{text}",
            textfont=dict(size=8, color="#ffffff"),
            colorscale=[[0, "#f0f4f8"], [0.001, "#c6dff0"], [0.15, "#6baed6"],
                        [0.4, "#2171b5"], [0.7, "#1a5e8a"], [1.0, "#001c4b"]],
            zmin=0, showscale=True, xgap=2, ygap=2,
            colorbar=dict(title=dict(text="AUM (EUR)", font=dict(size=9)),
                          tickfont=dict(size=8), thickness=10, len=0.8)))
        _fig_ws_p.update_layout(
            height=max(350, 26 * len(_ws_clients_p) + 100),
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff", font_color="#001c4b",
            xaxis=dict(side="top", tickangle=-30), yaxis=dict(autorange="reversed"),
            margin=dict(l=10, r=60, t=60, b=10))
        _add_plotly_slide(prs, _fig_ws_p, "WHITESPACE ANALYSIS — CROSS-SELL ENGINE",
                          "Bleu fonce = AUM Finance eleve. Gris clair = opportunite cross-sell.")

    # ── Slide 7 : Choropleth Map ─────────────────────────────────────────
    _df_cty_p = db.get_aum_by_country()
    if not _df_cty_p.empty and _df_cty_p["total_aum"].sum() > 0:
        _ISO_P = {
            "United Arab Emirates": "ARE", "Saudi Arabia": "SAU", "Qatar": "QAT",
            "Kuwait": "KWT", "Bahrain": "BHR", "Oman": "OMN",
            "United Kingdom": "GBR", "France": "FRA", "Germany": "DEU",
            "Switzerland": "CHE", "Luxembourg": "LUX", "Netherlands": "NLD",
            "Italy": "ITA", "Spain": "ESP", "Belgium": "BEL", "Austria": "AUT",
            "Sweden": "SWE", "Norway": "NOR", "Denmark": "DNK", "Finland": "FIN",
            "Singapore": "SGP", "Japan": "JPN", "Hong Kong": "HKG",
            "China": "CHN", "South Korea": "KOR", "Australia": "AUS", "India": "IND",
            "United States": "USA", "Canada": "CAN", "Brazil": "BRA",
            "Mexico": "MEX", "South Africa": "ZAF", "Egypt": "EGY",
        }
        _df_map_p = _df_cty_p.copy()
        _df_map_p["iso"] = _df_map_p["country"].map(_ISO_P)
        _df_map_p = _df_map_p.dropna(subset=["iso"])
        if not _df_map_p.empty:
            _fig_map_p = go.Figure(go.Choropleth(
                locations=_df_map_p["iso"], z=_df_map_p["total_aum"],
                colorscale=[[0, "#f0f4f8"], [0.2, "#c6dff0"], [0.4, "#6baed6"],
                            [0.6, "#2171b5"], [0.8, "#1a5e8a"], [1.0, "#001c4b"]],
                marker_line_color="#ffffff", marker_line_width=0.5,
                showscale=True,
                colorbar=dict(title=dict(text="AUM", font=dict(size=9)),
                              tickfont=dict(size=8), thickness=10, len=0.6)))
            _fig_map_p.update_layout(
                geo=dict(showframe=False, showcoastlines=True,
                         coastlinecolor="#e8e8e8", projection_type="natural earth",
                         bgcolor="#ffffff", landcolor="#f8f9fa",
                         showcountries=True, countrycolor="#e0e0e0"),
                height=420, paper_bgcolor="#ffffff", font_color="#001c4b",
                margin=dict(l=0, r=0, t=10, b=10))
            _add_plotly_slide(prs, _fig_map_p, "GLOBAL ROADSHOW MAP — AUM PAR PAYS",
                              "AUM total (Funded + Pipeline actif) par pays")

    # ── Slide 8 : AUM Time Machine ───────────────────────────────────────
    _df_hist_p = db.get_historical_aum(days_back=365)
    if not _df_hist_p.empty:
        _fig_hist_p = go.Figure()
        _fig_hist_p.add_trace(go.Scatter(
            x=_df_hist_p["date"], y=_df_hist_p["funded_aum"],
            name="Funded AUM", mode="lines",
            line=dict(color="#001c4b", width=2.5),
            fill="tozeroy", fillcolor="rgba(0,28,75,0.08)"))
        _fig_hist_p.add_trace(go.Scatter(
            x=_df_hist_p["date"], y=_df_hist_p["pipeline_aum"],
            name="Active Pipeline", mode="lines",
            line=dict(color="#019ee1", width=2, dash="dot")))
        _fig_hist_p.update_layout(
            height=360, paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font_color="#001c4b",
            legend=dict(orientation="h", y=1.05, x=0.5, xanchor="center"),
            xaxis=dict(showgrid=False, tickformat="%b %Y"),
            yaxis=dict(showgrid=True, gridcolor="#e8e8e8", tickformat=".2s"),
            margin=dict(l=20, r=20, t=30, b=20))
        _add_plotly_slide(prs, _fig_hist_p, "AUM TIME MACHINE — EVOLUTION HISTORIQUE",
                          "Funded AUM et Pipeline Actif sur 12 mois")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# TIMEFRAME HELPER
# ---------------------------------------------------------------------------
from datetime import date as _date_cls

TIMEFRAMES = ["Max", "YTD", "1M Rolling", "3M Rolling", "6M Rolling", "1Y Rolling", "3Y Rolling"]

def _timeframe_cutoff(timeframe: str):
    today = _date_cls.today()
    if timeframe == "YTD":         return _date_cls(today.year, 1, 1)
    if timeframe == "1M Rolling":  return today - timedelta(days=30)
    if timeframe == "3M Rolling":  return today - timedelta(days=91)
    if timeframe == "6M Rolling":  return today - timedelta(days=182)
    if timeframe == "1Y Rolling":  return today - timedelta(days=365)
    if timeframe == "3Y Rolling":  return today - timedelta(days=1095)
    return None

# ---------------------------------------------------------------------------
# CONFIG
# ---------------------------------------------------------------------------
st.set_page_config(page_title="CRM - Asset Management",
                   layout="wide", initial_sidebar_state="expanded")

# ---------------------------------------------------------------------------
# CSS (conservé intégralement)
# ---------------------------------------------------------------------------
st.markdown("""
<style>
.stApp, .main .block-container {
    background-color: #ffffff;
    color: #001c4b;
    font-family: 'Segoe UI', Arial, sans-serif;
}
[data-testid="stSidebar"] { background-color: #001c4b; }
[data-testid="stSidebar"] * { color: #ffffff !important; }
[data-testid="stSidebar"] h1,
[data-testid="stSidebar"] h2,
[data-testid="stSidebar"] h3 { color: #019ee1 !important; }
.crm-header {
    background: #001c4b; padding: 16px 24px;
    margin-bottom: 16px; border-bottom: 3px solid #019ee1;
}
.crm-header h1 { color:#ffffff !important; margin:0; font-size:1.4rem; font-weight:700; }
.crm-header p  { color:#7ab8d8; margin:3px 0 0 0; font-size:0.80rem; }

a.clink {
    display: block; text-decoration: none !important;
    color: inherit; cursor: pointer; outline: none;
}
a.clink:hover, a.clink:focus, a.clink:visited { text-decoration: none !important; color: inherit; }

.kpi-grid {
    display: grid;
    grid-template-columns: repeat(5, 1fr);
    gap: 8px; margin-bottom: 4px;
}
.kpi-card {
    background: #001c4b; padding: 14px 10px 13px 10px;
    text-align: center; border: 1px solid #1a5e8a;
    border-bottom: 2px solid #019ee1;
    transition: background 0.15s, border-color 0.15s, box-shadow 0.15s;
    position: relative;
}
.kpi-card-clickable::after {
    content: "▸"; position: absolute; bottom: 5px; right: 7px;
    font-size: 0.54rem; color: #4a8fbd; opacity: 0.7;
    transition: color 0.15s, opacity 0.15s;
}
a.clink:hover .kpi-card-clickable {
    background: #0b3060; border-color: #f07d00 !important;
    box-shadow: 0 0 0 2px #f07d0050;
}
a.clink:hover .kpi-card-clickable::after { color: #f07d00; opacity: 1; }
.kpi-card-static { border-bottom: 2px solid #1a5e8a; }

.kpi-label { font-size:0.64rem; color:#7ab8d8; text-transform:uppercase; letter-spacing:0.8px; margin-bottom:5px; font-weight:600; }
.kpi-value { font-size:1.4rem; font-weight:800; color:#ffffff; }
.kpi-sub   { font-size:0.61rem; color:#c8dde8; margin-top:3px; }

.statut-grid { display: grid; gap: 6px; margin-bottom: 8px; }
.statut-pill {
    padding: 8px 6px; text-align: center;
    transition: filter 0.12s, box-shadow 0.12s; position: relative;
}
a.clink:hover .statut-pill { filter: brightness(1.10); box-shadow: 0 0 0 2px #f07d0055; }

.alert-overdue {
    background: #fef6f0; border-left: 3px solid #f07d00;
    padding: 8px 12px; margin: 3px 0;
    font-size: 0.77rem; color: #001c4b;
    transition: background 0.12s, box-shadow 0.12s; position: relative;
}
a.clink:hover .alert-overdue { background: #fdebd5; box-shadow: inset 3px 0 0 #f07d00; }

.activity-row {
    border-left: 2px solid #019ee1; padding: 7px 11px; margin: 3px 0;
    background: #f9fbfd; font-size: 0.78rem;
    transition: background 0.12s, border-left-color 0.12s; display: block;
}
a.clink:hover .activity-row { background: #e4f1f9; border-left-color: #f07d00; }

.badge-retard {
    display:inline-block; background:#f07d00; color:#ffffff;
    padding:1px 7px; font-size:0.66rem; font-weight:700; letter-spacing:0.4px;
}
.perimetre-badge {
    display:inline-block; background:#019ee122; border:1px solid #019ee155;
    padding:2px 8px; font-size:0.70rem; color:#ffffff; font-weight:600; margin:2px;
}
.detail-panel {
    background:#f2f6fa; border-left:3px solid #019ee1;
    padding:14px 16px 11px 16px; margin-top:10px;
}
.sales-card { background:#f4f8fc; border:1px solid #001c4b18; padding:13px; border-top:3px solid #019ee1; }
.sales-card-name { font-size:0.87rem; font-weight:700; color:#001c4b; margin-bottom:8px; padding-bottom:5px; border-bottom:1px solid #e8e8e8; }
.sales-metric { font-size:0.69rem; color:#666; margin-bottom:2px; }
.sales-metric-val { font-size:0.93rem; font-weight:700; color:#001c4b; }
.sales-metric-acc { font-size:0.93rem; font-weight:700; color:#019ee1; }
.section-title { font-size:0.90rem; font-weight:700; color:#001c4b; border-bottom:2px solid #001c4b22; padding-bottom:4px; margin:14px 0 9px 0; }
.pipeline-hint { background:#019ee108; border-left:2px solid #001c4b; padding:6px 11px; font-size:0.77rem; color:#001c4b; margin-bottom:8px; }
.sidebar-kpi { background:#ffffff14; padding:8px; margin-bottom:5px; }

.stTabs [data-baseweb="tab-list"] { background:#f0f4f8; border-bottom:2px solid #001c4b20; gap:0; }
.stTabs [data-baseweb="tab"] { color:#001c4b; font-weight:600; font-size:0.81rem; padding:7px 16px; background:#f0f4f8; border-right:1px solid #d0d8e0; }
.stTabs [aria-selected="true"] { background:#001c4b !important; color:#ffffff !important; }

.stButton > button {
    background:#019ee1; color:#ffffff; border:none; font-weight:600;
    padding:6px 15px; font-size:0.80rem; transition:background 0.12s;
}
.stButton > button:hover { background:#f07d00 !important; color:#ffffff !important; }
[data-testid="stDownloadButton"] > button { background:#019ee1 !important; color:#ffffff !important; border:none !important; font-weight:600 !important; }
[data-testid="stDownloadButton"] > button:hover { background:#f07d00 !important; color:#ffffff !important; }
[data-testid="stFormSubmitButton"] > button { background:#001c4b !important; color:#ffffff !important; font-weight:700 !important; }
[data-testid="stFormSubmitButton"] > button:hover { background:#019ee1 !important; }
.stSelectbox label, .stTextInput label, .stNumberInput label,
.stDateInput label, .stTextArea label, .stRadio label { color:#001c4b !important; font-weight:600; font-size:0.78rem; }
h1,h2,h3,h4 { color:#001c4b !important; }
hr { border-color:#001c4b10; }
code { background:#001c4b08; color:#001c4b; }

.statut-badge {
    display: inline-block; padding: 2px 8px; font-size: 0.68rem;
    font-weight: 700; border-radius: 0; letter-spacing: 0.3px;
}
.statut-Funded        { background:#1a5e8a22; color:#1a5e8a; border:1px solid #1a5e8a55; }
.statut-Soft-Commit   { background:#2c7fb822; color:#2c7fb8; border:1px solid #2c7fb855; }
.statut-Due-Diligence { background:#004f8c22; color:#004f8c; border:1px solid #004f8c55; }
.statut-Initial-Pitch { background:#4a8fbd22; color:#4a8fbd; border:1px solid #4a8fbd55; }
.statut-Prospect      { background:#9ecae122; color:#2c7fb8; border:1px solid #9ecae155; }
.statut-Lost          { background:#88888822; color:#555;    border:1px solid #aaaaaa55; }
.statut-Paused        { background:#c0c0c022; color:#666;    border:1px solid #c0c0c055; }
.statut-Redeemed      { background:#b8b8d022; color:#555;    border:1px solid #b8b8d055; }


/* ── Quick Action Center — force Marine primary buttons ── */
[data-testid="stButton"][key^="qac_"] > button,
button[kind="primary"] {
    background: #001c4b !important;
    color: #ffffff !important;
    border: none !important;
    font-weight: 700 !important;
}
button[kind="primary"]:hover {
    background: #019ee1 !important;
    color: #ffffff !important;
}

/* ── Structure card in CRM split-screen ── */
.struct-block {
    background:#f4f8fc; border-left:3px solid #001c4b30;
    padding:10px 14px; margin:6px 0; font-size:0.78rem;
}
.struct-label { font-size:0.65rem; color:#888; text-transform:uppercase; letter-spacing:0.5px; margin-bottom:3px; }
.struct-val   { font-weight:700; color:#001c4b; }
</style>
""", unsafe_allow_html=True)


# ---------------------------------------------------------------------------
# INIT DB
# ---------------------------------------------------------------------------
@st.cache_resource
def _init():
    db.init_db()
    return True
_init()


# ---------------------------------------------------------------------------
# VISUAL HELPERS
# ---------------------------------------------------------------------------
def statut_badge(statut):
    css = "statut-" + statut.replace(" ", "-")
    return '<span class="statut-badge {}">{}</span>'.format(css, statut)

def statut_dot(statut):
    dot_colors = {
        "Funded":"#1a5e8a","Soft Commit":"#2c7fb8","Due Diligence":"#004f8c",
        "Initial Pitch":"#4a8fbd","Prospect":"#9ecae1","Lost":"#aaaaaa",
        "Paused":"#c0c0c0","Redeemed":"#b8b8d0",
    }
    c = dot_colors.get(statut, "#888")
    return ('<span style="display:inline-block;width:8px;height:8px;border-radius:50%;'
            'background:{};margin-right:5px;vertical-align:middle;"></span>').format(c)

def _kyc_dot(kyc):
    colors = {"Validé": "#22a062", "En cours": ORANGE, "Bloqué": "#c0392b"}
    c = colors.get(kyc, "#aaa")
    return ('<span style="display:inline-block;width:10px;height:10px;border-radius:50%;'
            'background:{c};vertical-align:middle;margin-right:6px;" '
            'title="KYC : {kyc}"></span>').format(c=c, kyc=kyc)

def _tier_badge(tier):
    tier_colors = {"Tier 1": MARINE, "Tier 2": B_MID, "Tier 3": B_PAL}
    c = tier_colors.get(tier, GRIS)
    return ('<span style="background:{c};color:#fff;font-size:0.62rem;font-weight:700;'
            'padding:2px 8px;letter-spacing:0.3px;vertical-align:middle;">{t}</span>'
            ).format(c=c, t=tier)

def _content_funded(fonds_filter=None):
    df = db.get_funded_deals_detail(fonds_filter)
    if df.empty: st.info("Aucun deal Funded."); return
    df2 = df.copy()
    for col in ["AUM_Finance","AUM_Cible"]:
        if col in df2.columns: df2[col] = df2[col].apply(fmt_m)
    st.dataframe(df2, use_container_width=True, hide_index=True, height=min(380, 46 + len(df2) * 36))

def _content_pipeline(fonds_filter=None):
    df = db.get_active_deals_detail(fonds_filter)
    if df.empty: st.info("Aucun deal actif."); return
    df2 = df.copy()
    # Show Smart AUM as single "AUM Attendu" column
    if "AUM_Pipeline" in df2.columns:
        df2["AUM Attendu"] = df2["AUM_Pipeline"].apply(fmt_m)
        drop_cols = [c for c in ["AUM_Revise","AUM_Cible","AUM_Pipeline"] if c in df2.columns]
        df2 = df2.drop(columns=drop_cols)
    if "Prochaine_Action" in df2.columns:
        df2["Prochaine_Action"] = df2["Prochaine_Action"].apply(
            lambda d: d.isoformat() if isinstance(d, date) else str(d or "—"))
    st.dataframe(df2, use_container_width=True, hide_index=True, height=min(380, 46 + len(df2) * 36))

def _content_lost(fonds_filter=None):
    df = db.get_lost_deals_detail(fonds_filter)
    if df.empty: st.info("Aucun deal Lost/Paused."); return
    df2 = df.copy()
    if "AUM_Cible" in df2.columns: df2["AUM_Cible"] = df2["AUM_Cible"].apply(fmt_m)
    st.dataframe(df2, use_container_width=True, hide_index=True, height=min(380, 46 + len(df2) * 36))

def _content_overdue():
    df = db.get_overdue_actions()
    if df.empty: st.info("Aucune action en retard."); return
    today = date.today()
    for _, row in df.iterrows():
        pid  = int(row.get("id", 0)) if "id" in df.columns else None
        deal = db.get_overdue_deal_full(pid) if pid else None
        nad  = row.get("next_action_date")
        days = (today - nad).days if isinstance(nad, date) else 0
        nad_str = nad.isoformat() if isinstance(nad, date) else "—"
        st.markdown(
            "<div style='border-left:3px solid #f07d00;padding:8px 14px;"
            "margin:8px 0;background:#fef6f0;'>"
            "<div style='font-size:0.84rem;font-weight:700;color:#001c4b;'>"
            "{} &nbsp;<span style='background:#f07d00;color:#fff;padding:1px 7px;"
            "font-size:0.66rem;font-weight:700;'>RETARD +{}j</span></div>".format(
                row.get("nom_client",""), days), unsafe_allow_html=True)
        if deal:
            c1, c2, c3 = st.columns(3)
            with c1:
                st.markdown("**Fonds**"); st.write(deal.get("fonds","—"))
                st.markdown("**Statut**"); st.write(deal.get("statut","—"))
                st.markdown("**Commercial**"); st.write(deal.get("sales_owner","—"))
            with c2:
                st.markdown("**AUM Cible**"); st.write(fmt_m(deal.get("target_aum_initial",0)))
                st.markdown("**AUM Revise**"); st.write(fmt_m(deal.get("revised_aum",0)))
                st.markdown("**AUM Finance**"); st.write(fmt_m(deal.get("funded_aum",0)))
            with c3:
                st.markdown("**Region**"); st.write(deal.get("region","—"))
                st.markdown("**Prochaine Action**"); st.write(nad_str)
                st.markdown("**Derniere Activite**")
                st.write(deal.get("derniere_activite","") or "—")
        st.markdown("</div>", unsafe_allow_html=True)

def _content_statut(statut_nom, fonds_filter=None):
    df = db.get_pipeline_by_statut(statut_nom, fonds_filter)
    if df.empty: st.info("Aucun deal en statut {}.".format(statut_nom)); return
    today = date.today()
    for _, row in df.iterrows():
        nad = row.get("next_action_date")
        if isinstance(nad, date):
            delta = (nad - today).days
            nad_str = nad.isoformat()
            timing_html = (
                "<span class='badge-retard'>RETARD +{}j</span>".format(abs(delta)) if delta < 0
                else "<span style='color:#019ee1;font-weight:700;'>Aujourd'hui</span>" if delta == 0
                else "<span style='color:#444;'>Dans {}j</span>".format(delta))
        else:
            nad_str = "—"; timing_html = "—"
        act = str(row.get("derniere_activite",""))
        st.markdown(
            "<div style='border-left:3px solid {ciel};padding:8px 14px;"
            "margin:6px 0;background:#f4f8fc;'>"
            "<div style='font-size:0.85rem;font-weight:700;color:{marine};'>{client}</div>"
            "<div style='font-size:0.76rem;color:#555;margin:2px 0;'>"
            "{fonds} &middot; {type} &middot; {region}</div>"
            "<div style='display:grid;grid-template-columns:1fr 1fr 1fr;gap:8px;margin-top:6px;'>"
            "<div><div style='font-size:0.67rem;color:#888;'>AUM Revise</div>"
            "<div style='font-weight:700;color:{marine};font-size:0.82rem;'>{aum}</div></div>"
            "<div><div style='font-size:0.67rem;color:#888;'>Prochaine Action</div>"
            "<div style='font-size:0.78rem;'>{nad}&nbsp;{timing}</div></div>"
            "<div><div style='font-size:0.67rem;color:#888;'>Commercial</div>"
            "<div style='font-size:0.78rem;color:#444;'>{owner}</div></div>"
            "</div>"
            "{act_block}</div>".format(
                ciel=CIEL, marine=MARINE, client=row.get("nom_client",""),
                fonds=row.get("fonds",""), type=row.get("type_client",""),
                region=row.get("region",""),
                aum=fmt_m(float(row.get("aum_pipeline", 0) or 0) if row.get("aum_pipeline") is not None
                        else (float(row.get("revised_aum",0) or 0) if float(row.get("revised_aum",0) or 0) > 0
                              else float(row.get("target_aum_initial",0) or 0))),
                nad=nad_str, timing=timing_html, owner=row.get("sales_owner",""),
                act_block=(
                    "<div style='margin-top:5px;font-size:0.73rem;color:#666;"
                    "border-top:1px solid #e8e8e8;padding-top:4px;'>"
                    "<b>Derniere activite :</b> {}</div>".format(act)
                ) if act else ""), unsafe_allow_html=True)


# ---------------------------------------------------------------------------
# STATE GUARD — prevents cross-tab state pollution
# Each tab uses its own namespace prefix in session_state
# ---------------------------------------------------------------------------
_filtre_effectif = None  # set in sidebar


# ============================================================================
# DIALOG FACTORY — all modals in one place
# ============================================================================

@st.dialog("Ajouter un compte client", width="large")
def dialog_add_client():
    with st.form("dlg_form_client", clear_on_submit=True):
        c1, c2 = st.columns(2)
        with c1:
            nom_client  = st.text_input("Nom du Client")
            type_client = st.selectbox("Type Client", TYPES_CLIENT)
            region      = st.selectbox("Région", REGIONS)
        with c2:
            is_parent = st.checkbox("Ce compte est une Maison Mère", value=False)
            all_clients_opts = db.get_all_clients()
            parent_options   = ["(Aucune — client autonome)"] + all_clients_opts["nom_client"].tolist()
            parent_sel       = st.selectbox("Maison Mère (optionnel)", parent_options,
                                             disabled=is_parent)
            tier_sel    = st.selectbox("Tier", db.TIERS_REFERENTIEL, index=1)
            kyc_sel     = st.selectbox("Statut KYC", db.KYC_STATUTS, index=1)
            country_sel = st.selectbox("Country", COUNTRIES_LIST)
        interests_sel = st.multiselect("Product Interests", db.PRODUCT_INTERESTS)
        if st.form_submit_button("Enregistrer le compte", use_container_width=True):
            if not nom_client.strip():
                st.error("Nom obligatoire.")
            else:
                try:
                    _parent_id = None
                    if not is_parent and parent_sel != "(Aucune — client autonome)":
                        _pr = all_clients_opts[all_clients_opts["nom_client"] == parent_sel]
                        if not _pr.empty:
                            _parent_id = int(_pr.iloc[0]["id"])
                    db.add_client(nom_client.strip(), type_client, region,
                                  country=country_sel, parent_id=_parent_id,
                                  tier=tier_sel, kyc_status=kyc_sel,
                                  product_interests=",".join(interests_sel))
                    st.success("Client {} ajouté.".format(nom_client))
                    st.rerun()
                except Exception as e:
                    st.warning("Ce client existe déjà." if "UNIQUE" in str(e)
                               else "Erreur : {}".format(e))


@st.dialog("Modifier le compte", width="large")
def dialog_edit_client(client_data: dict):
    sel_id   = int(client_data["id"])
    sel_nom  = str(client_data.get("nom_client",""))
    sel_type = str(client_data.get("type_client",""))
    sel_reg  = str(client_data.get("region",""))
    sel_tier = str(client_data.get("tier","Tier 2"))
    sel_kyc  = str(client_data.get("kyc_status","En cours"))
    sel_par  = str(client_data.get("parent_nom",""))
    cur_country = str(client_data.get("country",""))
    cur_prod    = [p.strip() for p in str(client_data.get("product_interests","")).split(",") if p.strip()]

    all_clients_edit = db.get_all_clients()
    parent_opts_edit = ["(Aucune — client autonome)"] + [
        n for n in all_clients_edit["nom_client"].tolist() if n != sel_nom]
    cur_parent_nom = sel_par or "(Aucune — client autonome)"

    with st.form("dlg_edit_company_{}".format(sel_id), clear_on_submit=False):
        ecp1, ecp2 = st.columns(2)
        with ecp1:
            ec_nom    = st.text_input("Nom du compte", value=sel_nom)
            ec_type   = st.selectbox("Type Client", TYPES_CLIENT,
                index=TYPES_CLIENT.index(sel_type) if sel_type in TYPES_CLIENT else 0)
            ec_region = st.selectbox("Région", REGIONS,
                index=REGIONS.index(sel_reg) if sel_reg in REGIONS else 0)
            ec_country = st.selectbox("Country", COUNTRIES_LIST,
                index=COUNTRIES_LIST.index(cur_country) if cur_country in COUNTRIES_LIST else 0)
        with ecp2:
            ec_tier   = st.selectbox("Tier", db.TIERS_REFERENTIEL,
                index=db.TIERS_REFERENTIEL.index(sel_tier) if sel_tier in db.TIERS_REFERENTIEL else 1)
            ec_kyc    = st.selectbox("Statut KYC", db.KYC_STATUTS,
                index=db.KYC_STATUTS.index(sel_kyc) if sel_kyc in db.KYC_STATUTS else 1)
            ec_parent = st.selectbox("Maison Mère", parent_opts_edit,
                index=parent_opts_edit.index(cur_parent_nom)
                      if cur_parent_nom in parent_opts_edit else 0)
            ec_is_parent = st.checkbox("Ce compte est une Maison Mère", value=(sel_par == ""))
        ec_interests = st.multiselect("Product Interests", db.PRODUCT_INTERESTS, default=cur_prod)
        esave, ecancel = st.columns(2)
        with esave:
            if st.form_submit_button("Enregistrer les modifications", use_container_width=True):
                _new_parent_id = None
                if not ec_is_parent and ec_parent != "(Aucune — client autonome)":
                    _pr = all_clients_edit[all_clients_edit["nom_client"] == ec_parent]
                    if not _pr.empty:
                        _new_parent_id = int(_pr.iloc[0]["id"])
                ok, err = db.update_client(
                    sel_id, ec_nom.strip(), ec_type, ec_region,
                    country=ec_country, parent_id=_new_parent_id,
                    tier=ec_tier, kyc_status=ec_kyc,
                    product_interests=",".join(ec_interests))
                if ok:
                    st.success("Compte mis à jour.")
                    st.rerun()
                else:
                    st.error("Erreur : {}".format(err))
        with ecancel:
            if st.form_submit_button("Annuler", use_container_width=True):
                st.rerun()


@st.dialog("Ajouter un contact", width="large")
def dialog_add_contact(client_id: int, client_nom: str = ""):
    st.caption("Compte : **{}**".format(client_nom) if client_nom else "")
    with st.form("dlg_add_contact_{}".format(client_id), clear_on_submit=True):
        cc1, cc2 = st.columns(2)
        with cc1:
            qc_prenom = st.text_input("Prénom")
            qc_nom    = st.text_input("Nom")
            qc_role   = st.selectbox("Rôle", [""] + db.ROLES_CONTACT)
        with cc2:
            qc_email   = st.text_input("Email")
            qc_tel     = st.text_input("Téléphone")
            qc_li      = st.text_input("LinkedIn (URL)")
        qc_primary = st.checkbox("Contact principal")
        if st.form_submit_button("Enregistrer le Contact", use_container_width=True):
            if not qc_nom.strip():
                st.error("Nom obligatoire.")
            else:
                db.add_contact(client_id, qc_prenom, qc_nom, qc_role,
                               qc_email, qc_tel, qc_li, qc_primary)
                st.success("Contact ajouté.")
                st.rerun()


@st.dialog("Modifier le contact", width="large")
def dialog_edit_contact(contact_data: dict):
    ct_id = int(contact_data.get("id", 0))
    with st.form("dlg_edit_ct_{}".format(ct_id), clear_on_submit=True):
        ec1, ec2 = st.columns(2)
        with ec1:
            ec_prenom = st.text_input("Prénom", value=str(contact_data.get("prenom","")))
            ec_nom    = st.text_input("Nom",    value=str(contact_data.get("nom","")))
            ec_role   = st.selectbox("Rôle", [""] + db.ROLES_CONTACT,
                index=([""] + db.ROLES_CONTACT).index(str(contact_data.get("role","")))
                      if str(contact_data.get("role","")) in ([""] + db.ROLES_CONTACT) else 0)
        with ec2:
            ec_email = st.text_input("Email",     value=str(contact_data.get("email","")))
            ec_tel   = st.text_input("Téléphone", value=str(contact_data.get("telephone","")))
            ec_li    = st.text_input("LinkedIn",  value=str(contact_data.get("linkedin","")))
        ec_primary = st.checkbox("Contact principal", value=bool(contact_data.get("is_primary", 0)))
        ecs1, ecs2 = st.columns(2)
        with ecs1:
            if st.form_submit_button("Enregistrer les modifications", use_container_width=True):
                ok, err = db.update_contact(ct_id, ec_prenom, ec_nom,
                                            ec_role, ec_email, ec_tel, ec_li, ec_primary)
                if ok:
                    st.success("Contact mis à jour.")
                    st.rerun()
                else:
                    st.error(err)
        with ecs2:
            if st.form_submit_button("Annuler", use_container_width=True):
                st.rerun()


@st.dialog("Ajouter un deal", width="large")
def dialog_add_deal(preselect_client_id: int = None):
    _clients = db.get_client_options()
    _st_team = db.get_sales_team()
    _marchés = sorted(_st_team["marche"].unique().tolist()) if not _st_team.empty else ["Global"]
    if not _clients:
        st.info("Ajoutez d'abord un client.")
        return
    with st.form("dlg_add_deal", clear_on_submit=True):
        ca, cb = st.columns(2)
        with ca:
            # Pre-select client if called from CRM tab
            client_names = list(_clients.keys())
            default_idx = 0
            if preselect_client_id:
                _matching = [i for i, cid in enumerate(_clients.values())
                             if cid == preselect_client_id]
                if _matching:
                    default_idx = _matching[0]
            client_sel  = st.selectbox("Client", client_names, index=default_idx)
            fonds_sel   = st.selectbox("Fonds", FONDS)
            statut_sel  = st.selectbox("Statut", STATUTS)
            marche_sel  = st.selectbox("Marché", ["Tous"] + _marchés)
            if marche_sel == "Tous":
                _owners = _st_team["nom"].tolist() if not _st_team.empty else ["Non assigne"]
            else:
                _owners = _st_team[_st_team["marche"]==marche_sel]["nom"].tolist()
            if not _owners: _owners = ["Non assigne"]
            owner_sel = st.selectbox("Commercial", _owners)
        with cb:
            target_aum_m  = st.number_input("AUM Cible (en M€)", min_value=0.0, step=1.0, format="%.2f")
            revised_aum_m = st.number_input("AUM Révisé (en M€)", min_value=0.0, step=1.0, format="%.2f")
            funded_aum_m  = st.number_input("AUM Financé (en M€)", min_value=0.0, step=1.0, format="%.2f",
                                             help="Laisser à 0 si Funded → utilise l'AUM Révisé automatiquement")
            closing_prob = st.slider("Probabilité de closing (%)", 0, 100, 50, 5)
        raison_perte, concurrent = "", ""
        if statut_sel in ("Lost","Paused"):
            cc, cd = st.columns(2)
            with cc: raison_perte = st.selectbox("Raison", RAISONS_PERTE)
            with cd: concurrent   = st.text_input("Concurrent")
        next_action = st.date_input("Prochaine Action", value=date.today() + timedelta(days=14))
        if st.form_submit_button("Enregistrer le Deal", use_container_width=True):
            if statut_sel in ("Lost","Paused") and not raison_perte:
                st.error("Raison obligatoire.")
            else:
                db.add_pipeline_entry(
                    _clients[client_sel], fonds_sel, statut_sel,
                    target_aum_m * 1_000_000,
                    revised_aum_m * 1_000_000,
                    funded_aum_m * 1_000_000,
                    raison_perte, concurrent, next_action.isoformat(),
                    owner_sel, closing_prob)
                st.success("Deal {} / {} enregistré.".format(fonds_sel, client_sel))
                st.rerun()


@st.dialog("Enregistrer une activité", width="large")
def dialog_add_activity(preselect_client_id: int = None):
    _clients = db.get_client_options()
    if not _clients:
        st.info("Ajoutez d'abord un client.")
        return
    with st.form("dlg_add_activity", clear_on_submit=True):
        client_names = list(_clients.keys())
        default_idx  = 0
        if preselect_client_id:
            _matching = [i for i, cid in enumerate(_clients.values())
                         if cid == preselect_client_id]
            if _matching:
                default_idx = _matching[0]
        dlg_client = st.selectbox("Client", client_names, index=default_idx)
        dlg_type   = st.selectbox("Type d'interaction", TYPES_INTERACTION)
        dlg_date   = st.date_input("Date", value=date.today())
        dlg_notes  = st.text_area("Notes", height=100)
        if st.form_submit_button("Enregistrer", use_container_width=True):
            db.add_activity(_clients[dlg_client], dlg_date.isoformat(), dlg_notes, dlg_type)
            st.success("Activité enregistrée.")
            st.rerun()


@st.dialog("Modifier le Deal", width="large")
def dialog_edit_pipeline(pipeline_id: int, row_data: dict):
    """Modale d'édition pipeline — sans scroll, état isolé via clé 'crm_pipeline_dialog_id'."""
    client_name    = str(row_data.get("nom_client",""))
    current_statut = str(row_data.get("statut","Prospect"))

    st.markdown(
        '<div style="font-size:0.84rem;font-weight:700;color:{marine};">'
        'Client : <span style="color:{ciel};">{name}</span>'
        '&nbsp;{badge}&nbsp;<span style="font-size:0.68rem;color:#888;">ID #{pid}</span>'
        '</div>'.format(
            marine=MARINE, ciel=B_MID, name=client_name,
            badge=statut_badge(current_statut), pid=pipeline_id),
        unsafe_allow_html=True)

    _st_team    = db.get_sales_team()
    _marchés_ed = sorted(_st_team["marche"].unique().tolist()) if not _st_team.empty else ["Global"]

    with st.form(key="dlg_edit_pipeline_{}".format(pipeline_id)):
        r1c1, r1c2 = st.columns(2)
        with r1c1:
            fi = FONDS.index(row_data["fonds"]) if row_data["fonds"] in FONDS else 0
            new_fonds  = st.selectbox("Fonds", FONDS, index=fi)
            cur_owner  = str(row_data.get("sales_owner") or "Non assigne")
            cur_marche = "Tous"
            if not _st_team.empty:
                rows_m = _st_team[_st_team["nom"] == cur_owner]
                if not rows_m.empty:
                    cur_marche = rows_m.iloc[0]["marche"]
            marche_edit = st.selectbox(
                "Marché commercial", ["Tous"] + _marchés_ed,
                index=(["Tous"] + _marchés_ed).index(cur_marche)
                      if cur_marche in (["Tous"] + _marchés_ed) else 0,
                key="dlg_marche_{}".format(pipeline_id))
            if marche_edit == "Tous":
                _owners_edit = _st_team["nom"].tolist() if not _st_team.empty else ["Non assigne"]
            else:
                _owners_edit = _st_team[_st_team["marche"]==marche_edit]["nom"].tolist()
            if not _owners_edit: _owners_edit = ["Non assigne"]
            oi = _owners_edit.index(cur_owner) if cur_owner in _owners_edit else 0
            new_sales = st.selectbox("Commercial", _owners_edit, index=oi,
                                     key="dlg_owner_{}".format(pipeline_id))
        with r1c2:
            si = STATUTS.index(current_statut) if current_statut in STATUTS else 0
            new_statut = st.selectbox("Statut", STATUTS, index=si)
            cur_prob   = float(row_data.get("closing_probability") or 50)
            new_prob   = st.slider("Probabilité de closing (%)", 0, 100, int(cur_prob), 5,
                                   key="dlg_prob_{}".format(pipeline_id))
        r2c1, r2c2, r2c3 = st.columns(3)
        with r2c1:
            new_target_m = st.number_input(
                "AUM Cible (en M€)",
                value=round(float(row_data.get("target_aum_initial",0.0)) / 1_000_000, 2),
                min_value=0.0, step=1.0, format="%.2f")
        with r2c2:
            new_revised_m = st.number_input(
                "AUM Révisé (en M€)",
                value=round(float(row_data.get("revised_aum",0.0)) / 1_000_000, 2),
                min_value=0.0, step=1.0, format="%.2f")
        with r2c3:
            new_funded_m = st.number_input(
                "AUM Financé (en M€)",
                value=round(float(row_data.get("funded_aum",0.0)) / 1_000_000, 2),
                min_value=0.0, step=1.0, format="%.2f",
                help="Laisser à 0 si Funded → utilise l'AUM Révisé automatiquement")
        r3c1, r3c2, r3c3 = st.columns(3)
        with r3c1:
            ropts  = [""] + RAISONS_PERTE
            cur_r  = str(row_data.get("raison_perte") or "")
            ri     = ropts.index(cur_r) if cur_r in ropts else 0
            lbl_r  = "Raison (obligatoire)" if new_statut in ("Lost","Paused") else "Raison"
            new_raison = st.selectbox(lbl_r, ropts, index=ri)
        with r3c2:
            new_conc = st.text_input("Concurrent", value=str(row_data.get("concurrent_choisi") or ""))
        with r3c3:
            nad = row_data.get("next_action_date")
            if not isinstance(nad, date): nad = date.today() + timedelta(days=14)
            new_nad = st.date_input("Prochaine Action", value=nad)

        col_save, col_cancel = st.columns([3, 1])
        with col_save:
            sub    = st.form_submit_button("Sauvegarder les modifications", use_container_width=True)
        with col_cancel:
            cancel = st.form_submit_button("Annuler", use_container_width=True)

    # Audit log inside dialog
    df_audit = db.get_audit_log(pipeline_id)
    if not df_audit.empty:
        with st.expander("Historique modifications", expanded=False):
            st.dataframe(df_audit, use_container_width=True, hide_index=True,
                         height=min(200, 46 + len(df_audit) * 36))

    # Delete deal — inside dialog
    st.markdown("---")
    if st.button("Supprimer ce deal (irréversible)", type="tertiary",
                 key="dlg_del_btn_{}".format(pipeline_id)):
        st.session_state["crm_pipeline_del_confirm"] = pipeline_id

    if st.session_state.get("crm_pipeline_del_confirm") == pipeline_id:
        st.error("Suppression définitive du deal #{} et de son historique.".format(pipeline_id))
        dd1, dd2, _ = st.columns([1, 1, 4])
        with dd1:
            if st.button("Confirmer", key="dlg_del_yes_{}".format(pipeline_id)):
                ok, err = db.delete_pipeline_row(pipeline_id)
                if ok:
                    st.session_state.pop("crm_pipeline_del_confirm", None)
                    st.session_state.pop("crm_pipeline_dialog_id", None)
                    st.rerun()
                else:
                    st.error(err)
        with dd2:
            if st.button("Annuler", key="dlg_del_no_{}".format(pipeline_id)):
                st.session_state.pop("crm_pipeline_del_confirm", None)
                st.rerun()

    # Form submission
    if cancel:
        st.session_state.pop("crm_pipeline_dialog_id", None)
        st.rerun()

    if sub:
        new_target  = new_target_m  * 1_000_000
        new_revised = new_revised_m * 1_000_000
        new_funded  = new_funded_m  * 1_000_000
        if new_statut == "Funded":
            kyc_check = str(row_data.get("kyc_status", "En cours"))
            if kyc_check in ("Bloqué", "En cours"):
                st.error("KYC incomplet ({}). Validez le KYC avant de passer au statut Funded.".format(kyc_check))
                return
        ok, msg = db.update_pipeline_row({
            "id": pipeline_id, "fonds": new_fonds, "statut": new_statut,
            "target_aum_initial": new_target, "revised_aum": new_revised,
            "funded_aum": new_funded, "raison_perte": new_raison,
            "concurrent_choisi": new_conc, "next_action_date": new_nad,
            "sales_owner": new_sales, "closing_probability": new_prob})
        if ok:
            st.session_state.pop("crm_pipeline_dialog_id", None)
            st.success("Deal mis à jour.")
            st.rerun()
        else:
            st.error(msg)


@st.dialog("Modifier l'activité", width="large")
def dialog_edit_activity(act_id: int, act_data: dict):
    with st.form("dlg_edit_act_{}".format(act_id), clear_on_submit=True):
        ea1, ea2 = st.columns(2)
        with ea1:
            ea_type = st.selectbox("Type", TYPES_INTERACTION,
                index=TYPES_INTERACTION.index(act_data.get("type","Call"))
                      if act_data.get("type") in TYPES_INTERACTION else 0)
            try:
                ea_date_val = date.fromisoformat(act_data.get("date",""))
            except Exception:
                ea_date_val = date.today()
            ea_date = st.date_input("Date", value=ea_date_val)
        with ea2:
            ea_notes = st.text_area("Notes", value=act_data.get("notes",""), height=100)
        col_save, col_cancel = st.columns(2)
        with col_save:
            if st.form_submit_button("Enregistrer les modifications", use_container_width=True):
                ok, err = db.update_activity(act_id, ea_date.isoformat(), ea_notes, ea_type)
                if ok:
                    st.rerun()
                else:
                    st.error(err)
        with col_cancel:
            if st.form_submit_button("Annuler", use_container_width=True):
                st.rerun()




@st.dialog("Meeting Brief — Analyse Compte", width="large")
def dialog_meeting_brief(client_data: dict, group_summary: dict,
                          contacts_df, activities_df, whitespace_df):
    """One-minute briefing for CRM director / relationship manager."""
    sel_nom = client_data.get("nom_client", "")
    st.markdown(
        '<div style="background:#001c4b;padding:14px 18px;margin-bottom:12px;">'
        '<div style="font-size:0.72rem;color:#7ab8d8;text-transform:uppercase;letter-spacing:1px;">'
        'Meeting Brief — Confidentiel</div>'
        '<div style="font-size:1.2rem;font-weight:800;color:#ffffff;">{}</div>'
        '</div>'.format(sel_nom), unsafe_allow_html=True)

    bc1, bc2 = st.columns(2)
    with bc1:
        st.markdown("**AUM Consolidé Groupe**")
        st.markdown(
            '<span style="font-size:1.6rem;font-weight:800;color:#001c4b;">{}</span>'.format(
                fmt_m(group_summary["aum_consolide"])), unsafe_allow_html=True)
        if group_summary["fonds_investis"]:
            st.markdown("**Fonds investis :** " + " · ".join(group_summary["fonds_investis"]))
        else:
            st.info("Aucun deal Funded à ce jour.")
    with bc2:
        st.markdown("**Contacts clés**")
        if not contacts_df.empty:
            for _, ct in contacts_df.head(3).iterrows():
                st.markdown("- **{}** {} — *{}*".format(
                    ct.get("prenom",""), ct.get("nom",""), ct.get("role","—")))
        else:
            st.caption("Aucun contact enregistré.")

    st.divider()
    st.markdown("**2 Dernières Activités**")
    if not activities_df.empty:
        for _, act in activities_df.head(2).iterrows():
            st.markdown("- `{}` **{}** — {}".format(
                str(act.get("date","")),
                str(act.get("type_interaction","")),
                str(act.get("notes",""))[:120]))
    else:
        st.caption("Aucune activité.")

    st.divider()
    st.markdown("**Opportunités Whitespace (Cross-Sell)**")
    if not whitespace_df.empty and sel_nom in whitespace_df.index:
        row = whitespace_df.loc[sel_nom]
        gaps = [f for f in row.index if pd.isna(row[f])]
        invested = [f for f in row.index if not pd.isna(row[f])]
        if gaps:
            st.markdown("Fonds **non investis** (opportunités) : " +
                        ", ".join("**{}**".format(g) for g in gaps))
        if invested:
            st.markdown("Fonds investis : " + ", ".join(invested))
    else:
        st.caption("Données whitespace indisponibles.")

    if group_summary["next_actions"]:
        st.divider()
        st.markdown("**Prochaines Actions Pipeline**")
        for act in group_summary["next_actions"][:3]:
            st.markdown("- **{}** — {} — {} — `{}`".format(
                act["fonds"], act["statut"],
                fmt_m(act["aum_pipeline"]), act["nad"]))


@st.dialog("Gérer les Commerciaux", width="large")
def dialog_manage_sales():
    st_team = db.get_sales_team()
    MARCHES_OPTIONS = ["GCC", "EMEA", "APAC", "Americas", "Nordics", "Global"]

    st.markdown("**Equipe commerciale actuelle**")
    if st_team.empty:
        st.info("Aucun commercial enregistré.")
    else:
        # data_editor : modification inline du nom et du marché
        edited_df = st.data_editor(
            st_team[["id", "nom", "marche"]].copy(),
            column_config={
                "id":     st.column_config.NumberColumn("ID", disabled=True, width="small"),
                "nom":    st.column_config.TextColumn("Nom", width="medium"),
                "marche": st.column_config.SelectboxColumn(
                    "Marché", options=MARCHES_OPTIONS, width="medium"),
            },
            hide_index=True,
            use_container_width=True,
            num_rows="fixed",
            key="sales_editor"
        )
        col_save, col_spacer = st.columns([1, 3])
        with col_save:
            if st.button("Enregistrer les modifications", key="dlg_sales_save",
                         use_container_width=True, type="primary"):
                errors = []
                for _, edit_row in edited_df.iterrows():
                    orig = st_team[st_team["id"] == edit_row["id"]]
                    if orig.empty:
                        continue
                    orig_nom    = str(orig.iloc[0]["nom"])
                    orig_marche = str(orig.iloc[0]["marche"])
                    new_nom     = str(edit_row["nom"]).strip()
                    new_marche  = str(edit_row["marche"]).strip()
                    if new_nom and (new_nom != orig_nom or new_marche != orig_marche):
                        ok, err = db.update_sales_member(
                            int(edit_row["id"]), new_nom, new_marche)
                        if not ok:
                            errors.append("ID {}: {}".format(edit_row["id"], err))
                if errors:
                    st.error("Erreurs : " + " | ".join(errors))
                else:
                    st.success("Modifications enregistrées.")
                    st.rerun()

        # Suppression individuelle
        st.markdown("**Supprimer un commercial**")
        sales_names = st_team["nom"].tolist()
        del_nom = st.selectbox(
            "Sélectionner le commercial à supprimer",
            options=[""] + sales_names,
            key="dlg_sales_del_sel"
        )
        if del_nom:
            st.warning(
                "La suppression de **{}** remettra ses deals "
                "pipeline à 'Non assigne'.".format(del_nom))
            if st.button("Confirmer la suppression", key="dlg_sales_del_confirm",
                         type="secondary", use_container_width=True):
                _del_id = int(
                    st_team[st_team["nom"] == del_nom].iloc[0]["id"])
                ok, err = db.delete_sales_member(_del_id)
                if ok:
                    st.success("Commercial supprimé.")
                    st.rerun()
                else:
                    st.error("Erreur : {}".format(err))

    st.divider()
    st.markdown("**Ajouter un commercial**")
    with st.form("dlg_add_sales", clear_on_submit=True):
        nc1, nc2 = st.columns(2)
        with nc1: new_sales_nom    = st.text_input("Nom")
        with nc2: new_sales_marche = st.selectbox("Marché", MARCHES_OPTIONS)
        if st.form_submit_button("Ajouter", use_container_width=True):
            if new_sales_nom.strip():
                ok = db.add_sales_member(new_sales_nom.strip(), new_sales_marche)
                st.success("Commercial ajouté." if ok else "Nom déjà existant.")
                st.rerun()


# ---------------------------------------------------------------------------
# SIDEBAR
# ---------------------------------------------------------------------------
with st.sidebar:
    st.markdown(
        '<div style="text-align:center;padding:6px 0 14px 0;">'
        '<div style="font-size:0.91rem;font-weight:800;color:#ffffff;">CRM Asset Management</div>'
        '<div style="font-size:0.65rem;color:#e8e8e8;margin-top:3px;">'
        + date.today().strftime("%d %B %Y") + '</div></div>',
        unsafe_allow_html=True)
    st.divider()

    kpis_global = db.get_kpis()
    st.markdown('<div style="color:#4a8fbd;font-size:0.67rem;font-weight:700;'
                'text-transform:uppercase;letter-spacing:1px;margin-bottom:5px;">Apercu Global</div>',
                unsafe_allow_html=True)
    for lbl, val in [("AUM Finance Total", fmt_m(kpis_global["total_funded"])),
                     ("Pipeline Actif",    fmt_m(kpis_global["pipeline_actif"]))]:
        st.markdown('<div class="sidebar-kpi">'
                    '<div style="font-size:0.62rem;color:#c8dde8;">{}</div>'
                    '<div style="font-size:1.06rem;font-weight:800;color:#ffffff;">{}</div>'
                    '</div>'.format(lbl, val), unsafe_allow_html=True)

    st.divider()
    st.markdown('<div style="color:#4a8fbd;font-size:0.67rem;font-weight:700;'
                'text-transform:uppercase;letter-spacing:1px;margin-bottom:5px;">Backup</div>',
                unsafe_allow_html=True)
    try:
        excel_bytes = db.get_excel_backup()
        st.download_button("Exporter Backup Excel", data=excel_bytes,
                           file_name="backup_crm_{}.xlsx".format(date.today().isoformat()),
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                           use_container_width=True)
    except Exception as e_xl:
        st.error("Backup Excel : {}".format(e_xl))

    st.divider()
    st.markdown('<div style="color:#f07d00;font-size:0.67rem;font-weight:700;'
                'text-transform:uppercase;letter-spacing:1px;margin-bottom:5px;">Export Hub</div>',
                unsafe_allow_html=True)

    _dyn_sidebar   = db.get_dynamic_filters()
    _fonds_dyn_sb  = _dyn_sidebar.get("fonds", FONDS)
    fonds_perimetre = st.multiselect(
        "Périmètre de l'export",
        options=_fonds_dyn_sb,
        default=_fonds_dyn_sb,
        key="fonds_perimetre_select"
    )
    _filtre_effectif = (fonds_perimetre
                        if (fonds_perimetre and len(fonds_perimetre) < len(_fonds_dyn_sb))
                        else None)
    mode_comex = st.toggle("Mode Comex — Anonymisation", value=False, key="hub_mode_comex")

    if fonds_perimetre and len(fonds_perimetre) < len(_fonds_dyn_sb):
        st.markdown("<div>{}</div>".format(" ".join(
            '<span class="perimetre-badge">{}</span>'.format(f) for f in fonds_perimetre)),
            unsafe_allow_html=True)
    elif not fonds_perimetre:
        st.warning("Sélectionnez au moins un fonds.")

    hub_format = st.radio(
        "Format du rapport",
        ["PDF", "PPTX", "Email"],
        horizontal=True,
        key="hub_format_radio"
    )

    # PDF-only options
    if hub_format == "PDF":
        st.markdown('<div style="color:#8ab8d8;font-size:0.62rem;font-weight:600;'
                    'text-transform:uppercase;letter-spacing:0.7px;margin:6px 0 3px 0;">Sections PDF</div>',
                    unsafe_allow_html=True)
        include_top10    = st.checkbox("Top 10 Inflows",     value=True,  key="hub_top10")
        include_outflows = st.checkbox("Top 10 Outflows",    value=False, key="hub_outflows",
                                       disabled=not include_top10)
        include_perf_pdf = st.checkbox("Performance NAV",    value=True,  key="hub_perf_nav")
    else:
        include_top10 = True; include_outflows = False; include_perf_pdf = False

    perf_data_pdf   = st.session_state.get("perf_data",   None)
    nav_base100_pdf = st.session_state.get("nav_base100", None)
    if perf_data_pdf is not None and hub_format == "PDF":
        st.caption("NAV chargée : {} fonds.".format(len(perf_data_pdf)))

    # ── Single "Générer le Rapport" button ────────────────────────────────────
    _btn_disabled = not fonds_perimetre
    if _btn_disabled:
        st.button("Générer le Rapport", key="hub_gen_btn_disabled",
                  disabled=True, use_container_width=True)
    elif st.button("Générer le Rapport", key="hub_gen_btn", use_container_width=True):
        with st.spinner("Génération en cours…"):
            try:
                pipeline_hub = db.get_pipeline_with_clients(fonds_filter=_filtre_effectif)
                kpis_hub     = db.get_kpis(fonds_filter=_filtre_effectif)

                if hub_format == "PDF":
                    aum_region_pdf = db.get_aum_by_region(fonds_filter=_filtre_effectif)
                    pf_pdf = perf_data_pdf
                    nb_pdf = nav_base100_pdf
                    if pf_pdf is not None and _filtre_effectif and "Fonds" in pf_pdf.columns:
                        pf_pdf = pf_pdf[pf_pdf["Fonds"].isin(_filtre_effectif)]
                    if nb_pdf is not None and _filtre_effectif and hasattr(nb_pdf, "columns"):
                        cols_k = [c for c in nb_pdf.columns if c in _filtre_effectif]
                        nb_pdf = nb_pdf[cols_k] if cols_k else None
                    _include_perf = (include_perf_pdf and pf_pdf is not None
                                     and hasattr(pf_pdf, "empty") and not pf_pdf.empty)
                    _bi_pngs = []
                    try:
                        _export_cfg = dict(format="png", engine="kaleido",
                                           width=900, height=500)
                        _export_layout = dict(
                            title="", showlegend=True,
                            margin=dict(l=40, r=40, t=20, b=30),
                            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff")
                        import copy as _copy_pdf
                        _df_cf_pdf = db.get_expected_cashflows()
                        if not _df_cf_pdf.empty and _df_cf_pdf["aum_pondere"].sum() > 0:
                            _months_pdf = sorted(_df_cf_pdf["mois"].unique().tolist())
                            _fig_cf_pdf = go.Figure()
                            for _i, _f in enumerate(sorted(_df_cf_pdf["fonds"].unique())):
                                _df_f = _df_cf_pdf[_df_cf_pdf["fonds"] == _f]
                                _mm = dict(zip(_df_f["mois"], _df_f["aum_pondere"]))
                                _fig_cf_pdf.add_trace(go.Bar(
                                    name=_f, x=_months_pdf,
                                    y=[_mm.get(m, 0) for m in _months_pdf],
                                    marker_color=PALETTE[_i % len(PALETTE)]))
                            _fig_cf_pdf.update_layout(barmode="stack", **_export_layout)
                            _bi_pngs.append(("Money in Motion — Cashflows",
                                io.BytesIO(_fig_cf_pdf.to_image(**_export_cfg))))

                        _df_ws_pdf = db.get_whitespace_matrix()
                        if not _df_ws_pdf.empty:
                            _ws_c = _df_ws_pdf.index.tolist()[:20]
                            _ws_f = _df_ws_pdf.columns.tolist()
                            _z = [[0.0 if (v is None or (isinstance(v, float) and np.isnan(v)))
                                   else float(v) for v in r]
                                  for r in _df_ws_pdf.loc[_ws_c, _ws_f].values.tolist()]
                            _fig_ws_pdf = go.Figure(go.Heatmap(
                                z=_z, x=_ws_f, y=_ws_c,
                                colorscale=[[0,"#f0f4f8"],[0.5,"#2171b5"],[1,"#001c4b"]],
                                xgap=2, ygap=2, showscale=True))
                            _fig_ws_pdf.update_layout(
                                xaxis=dict(side="top"), yaxis=dict(autorange="reversed"),
                                **_export_layout)
                            _bi_pngs.append(("Whitespace Analysis — Cross-Sell",
                                io.BytesIO(_fig_ws_pdf.to_image(**_export_cfg))))

                        _df_hist_pdf = db.get_historical_aum(days_back=365)
                        if not _df_hist_pdf.empty:
                            _fig_h_pdf = go.Figure()
                            _fig_h_pdf.add_trace(go.Scatter(
                                x=_df_hist_pdf["date"], y=_df_hist_pdf["funded_aum"],
                                name="Funded AUM", line=dict(color="#001c4b", width=2),
                                fill="tozeroy", fillcolor="rgba(0,28,75,0.08)"))
                            _fig_h_pdf.add_trace(go.Scatter(
                                x=_df_hist_pdf["date"], y=_df_hist_pdf["pipeline_aum"],
                                name="Active Pipeline",
                                line=dict(color="#019ee1", width=2, dash="dot")))
                            _fig_h_pdf.update_layout(**_export_layout)
                            _bi_pngs.append(("AUM Time Machine",
                                io.BytesIO(_fig_h_pdf.to_image(**_export_cfg))))
                    except Exception:
                        pass
                    pdf_bytes = pdf_gen.generate_pdf(
                        pipeline_df=pipeline_hub, kpis=kpis_hub,
                        aum_by_region=aum_region_pdf, mode_comex=mode_comex,
                        perf_data=pf_pdf, nav_base100_df=nb_pdf,
                        fonds_perimetre=fonds_perimetre,
                        include_top10=include_top10, include_outflows=include_outflows,
                        include_perf=_include_perf,
                        bi_chart_pngs=_bi_pngs if _bi_pngs else None)
                    fname_pdf = "report{}_{}.pdf".format(
                        "_comex" if mode_comex else "", date.today().isoformat())
                    st.download_button("Télécharger le PDF", data=pdf_bytes,
                                       file_name=fname_pdf, mime="application/pdf",
                                       key="hub_dl_pdf", use_container_width=True)
                    st.success("PDF généré.")

                elif hub_format == "PPTX":
                    pptx_bytes = generate_global_pptx(kpis_hub, pipeline_hub,
                                                       mode_comex=mode_comex)
                    fname_pptx = "rapport_global{}_{}.pptx".format(
                        "_comex" if mode_comex else "", date.today().isoformat())
                    st.download_button(
                        "Télécharger le PPTX", data=pptx_bytes,
                        file_name=fname_pptx,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="hub_dl_pptx", use_container_width=True)
                    st.success("PPTX généré.")

                elif hub_format == "Email":
                    import urllib.parse
                    # ── Données brutes ────────────────────────────────────────
                    _aum_f   = fmt_m(kpis_hub.get("total_funded", 0))
                    _pip_a   = fmt_m(kpis_hub.get("pipeline_actif", 0))
                    _wp      = fmt_m(kpis_hub.get("weighted_pipeline", 0))
                    _taux    = "{:.1f}%".format(kpis_hub.get("taux_conversion", 0))
                    _nb_f    = kpis_hub.get("nb_funded", 0)
                    _nb_l    = kpis_hub.get("nb_lost", 0)
                    _nb_a    = kpis_hub.get("nb_deals_actifs", 0)
                    _perims  = ", ".join(fonds_perimetre) if fonds_perimetre else "Tous les fonds"
                    # Date en français
                    _mois_fr = {1:"Janvier",2:"Fevrier",3:"Mars",4:"Avril",5:"Mai",
                                6:"Juin",7:"Juillet",8:"Aout",9:"Septembre",
                                10:"Octobre",11:"Novembre",12:"Decembre"}
                    _today_e = date.today()
                    _date_str = "{} {} {}".format(
                        _today_e.day, _mois_fr[_today_e.month], _today_e.year)
                    _comex_note = " [MODE COMEX — ANONYMISE]" if mode_comex else ""

                    # ── Section 2 : Alertes Opérationnelles ──────────────────
                    _df_overdue = db.get_overdue_actions()
                    _today = date.today()
                    if _df_overdue.empty:
                        _alertes_txt = "Aucune action commerciale en retard n'est recensee a ce jour."
                        _alertes_nb  = 0
                    else:
                        _alertes_nb  = len(_df_overdue)
                        _alertes_lines = []
                        # Limiter aux 5 deals en retard les plus importants
                        for _, _ov in _df_overdue.head(5).iterrows():
                            _nad_ov = _ov.get("next_action_date")
                            _days_l = (_today - _nad_ov).days if isinstance(_nad_ov, date) else 0
                            _nad_s  = _nad_ov.strftime("%d/%m/%Y") if isinstance(_nad_ov, date) else "—"
                            _alertes_lines.append(
                                "  - {client} / {fonds} ({statut}) : action prevue le {nad}"
                                ", soit un retard de {d} jour(s). Commercial : {owner}.".format(
                                    client=str(_ov.get("nom_client","—")),
                                    fonds=str(_ov.get("fonds","—")),
                                    statut=str(_ov.get("statut","—")),
                                    nad=_nad_s, d=_days_l,
                                    owner=str(_ov.get("sales_owner","—"))))
                        if _alertes_nb > 5:
                            _alertes_lines.append(
                                "  ...et {} autre(s) action(s) necessitant"
                                " l'attention des equipes.".format(_alertes_nb - 5))
                        _alertes_txt = "\n".join(_alertes_lines)

                    # ── Section 3 : Top Opportunités (Due Diligence / Soft Commit) ─
                    _df_top_opp = pipeline_hub[
                        pipeline_hub["statut"].isin(["Due Diligence", "Soft Commit"])
                    ].copy()
                    if not _df_top_opp.empty:
                        _df_top_opp["_aum_p"] = _df_top_opp.apply(
                            lambda r: float(r["revised_aum"]) if float(r.get("revised_aum",0) or 0) > 0
                                      else float(r.get("target_aum_initial",0) or 0), axis=1)
                        _df_top_opp = _df_top_opp.sort_values("_aum_p", ascending=False).head(3)
                    _top_opp_lines = []
                    for _rank, (_, _op) in enumerate(_df_top_opp.iterrows(), 1):
                        _client_nm = ("Opportunite {:02d}".format(_rank)
                                      if mode_comex else str(_op.get("nom_client","—")))
                        _top_opp_lines.append(
                            "  {rank}. {client} — {fonds} ({statut}) — AUM : {aum}"
                            " — Commercial : {owner}".format(
                                rank=_rank, client=_client_nm,
                                fonds=str(_op.get("fonds","—")),
                                statut=str(_op.get("statut","—")),
                                aum=fmt_m(float(_op.get("_aum_p",0))),
                                owner=str(_op.get("sales_owner","—"))))
                    if not _top_opp_lines:
                        _top_opp_txt = "Aucun deal en phase Due Diligence ou Soft Commit a ce jour."
                    else:
                        _top_opp_txt = "\n".join(_top_opp_lines)

                    # ── Rendu Markdown pour prévisualisation ──────────────────
                    _md_preview = (
                        "**NOTE DE SYNTHESE — PORTFOLIO COMMERCIAL**\n\n"
                        "**Date :** {date}{comex}  \n"
                        "**Perimetre :** {perim}\n\n"
                        "---\n\n"
                        "**I. SYNTHESE GLOBALE**\n\n"
                        "A la date du {date}, les encours finances s'elevent a **{aum_f}**, "
                        "repartis sur {nb_f} mandat(s) consolide(s). "
                        "Le pipeline commercial actif atteint **{pip_a}** ({nb_a} deal(s) en cours), "
                        "soit un pipeline pondere de **{wp}** apres application des probabilites de closing. "
                        "Le taux de conversion historique s'etablit a **{taux}** "
                        "({nb_f} deals finances, {nb_l} deals perdus).\n\n"
                        "---\n\n"
                        "**II. ALERTES OPERATIONNELLES** ({nb_ov} action(s) en retard)\n\n"
                        "{alertes}\n\n"
                        "---\n\n"
                        "**III. TOP OPPORTUNITES — Due Diligence et Soft Commit**\n\n"
                        "{top_opp}\n\n"
                        "---\n\n"
                        "*Document genere automatiquement par le CRM Asset Management — "
                        "Usage interne strictement confidentiel.*"
                    ).format(
                        date=_date_str,
                        comex=_comex_note,
                        perim=_perims,
                        aum_f=_aum_f,
                        pip_a=_pip_a,
                        wp=_wp,
                        taux=_taux,
                        nb_f=_nb_f,
                        nb_a=_nb_a,
                        nb_l=_nb_l,
                        nb_ov=_alertes_nb,
                        alertes=_alertes_txt if _alertes_nb > 0 else "Aucune action en retard.",
                        top_opp=_top_opp_txt,
                    )
                    st.markdown(_md_preview)

                    # ── Corps email pour mailto ───────────────────────────────
                    _body  = "Madame, Monsieur,\n\n"
                    _body += "Veuillez trouver ci-apres la note de synthese du portfolio commercial "
                    _body += "etablie a la date du {}{}.\n\n".format(_date_str, _comex_note)
                    _body += "I. SYNTHESE GLOBALE\n"
                    _body += "-" * 40 + "\n"
                    _body += "Encours finances (AUM) : {}\n".format(_aum_f)
                    _body += "Pipeline actif         : {} ({} deal(s))\n".format(_pip_a, _nb_a)
                    _body += "Pipeline pondere       : {}\n".format(_wp)
                    _body += "Taux de conversion     : {} ({} finances / {} perdus)\n".format(
                        _taux, _nb_f, _nb_l)
                    _body += "Perimetre              : {}\n\n".format(_perims)
                    _body += "II. ALERTES OPERATIONNELLES ({} action(s) en retard)\n".format(_alertes_nb)
                    _body += "-" * 40 + "\n"
                    _body += _alertes_txt + "\n\n"
                    _body += "III. TOP OPPORTUNITES (Due Diligence / Soft Commit)\n"
                    _body += "-" * 40 + "\n"
                    _body += _top_opp_txt + "\n\n"
                    _body += "--\n"
                    _body += "CRM Asset Management — Note generee automatiquement — Usage interne confidentiel."

                    _subject = "Note de Synthese Portfolio — {}{}".format(
                        date.today().strftime("%d/%m/%Y"),
                        " [COMEX]" if mode_comex else "")
                    mailto_url = "mailto:?subject={}&body={}".format(
                        urllib.parse.quote(_subject),
                        urllib.parse.quote(_body))
                    st.link_button("Ouvrir dans Outlook / Gmail",
                                   url=mailto_url,
                                   use_container_width=True)

            except Exception as e:
                st.error("Erreur génération : {}".format(e))

    st.divider()
    st.caption("Version 15.0 — Asset Management CRM — Universal Export Hub")


# ---------------------------------------------------------------------------
# EN-TETE + TABS
# ---------------------------------------------------------------------------
st.markdown(
    '<div class="crm-header"><h1>CRM &amp; Reporting — Asset Management</h1>'
    '<p>Pipeline commercial &middot; Suivi des mandats &middot; '
    'Reporting executif &middot; Performance NAV</p></div>',
    unsafe_allow_html=True)

tab_crm, tab_pipeline, tab_dash, tab_sales, tab_activites, tab_settings, tab_perf = st.tabs([
    "CRM Directory", "Pipeline Management", "Executive Dashboard",
    "Sales Tracking", "Activities", "Settings & Admin", "Performance & NAV"])


# ============================================================================
# ONGLET 1 — CRM DIRECTORY  (Split-Screen)
# ============================================================================
with tab_crm:
    st.markdown('<div class="section-title">CRM Directory — Clients &amp; Contacts</div>',
                unsafe_allow_html=True)

    # ── Toolbar ──────────────────────────────────────────────────────────────
    tb1, tb2, tb_spacer = st.columns([1, 1, 6])
    with tb1:
        if st.button("Nouveau compte", key="crm_tb_new_compte", use_container_width=True):
            dialog_add_client()
    with tb2:
        if st.button("Mailing List", key="crm_tb_mailing", use_container_width=True, type="secondary"):
            st.session_state["crm_show_mailing"] = not st.session_state.get("crm_show_mailing", False)

    # ── Mailing List Generator (toggled inline, no expander) ─────────────────
    if st.session_state.get("crm_show_mailing", False):
        st.markdown('<div class="pipeline-hint">Filtrer les contacts pour générer une liste de diffusion ciblée. Seuls les contacts avec un email sont inclus.</div>', unsafe_allow_html=True)
        mg1, mg2, mg3, mg4 = st.columns(4)
        with mg1: ml_regions   = st.multiselect("Region", db.REGIONS_REFERENTIEL, key="ml_regions")
        with mg2: ml_countries = st.multiselect("Country", COUNTRIES_LIST[1:], key="ml_countries")
        with mg3: ml_tiers     = st.multiselect("Tier", db.TIERS_REFERENTIEL, key="ml_tiers")
        with mg4: ml_interests = st.multiselect("Product Interests", db.PRODUCT_INTERESTS, key="ml_interests")
        df_ml = db.get_mailing_list(
            regions=ml_regions or None, countries=ml_countries or None,
            tiers=ml_tiers or None, product_interests=ml_interests or None)
        if df_ml.empty:
            st.info("Aucun contact ne correspond aux filtres.")
        else:
            display_cols = ["first_name","last_name","company","role","email"]
            st.markdown('<div class="pipeline-hint"><b>{}</b> contact(s) trouvé(s)</div>'.format(len(df_ml)),
                        unsafe_allow_html=True)
            st.dataframe(df_ml[display_cols].rename(columns={
                "first_name":"Prénom","last_name":"Nom","company":"Compte",
                "role":"Rôle","email":"Email"}),
                use_container_width=True, hide_index=True,
                height=min(280, 46 + len(df_ml) * 36))
            st.code("; ".join(df_ml["email"].tolist()), language=None)
            st.download_button("Exporter CSV", key="crm_ml_dl",
                data=df_ml[display_cols].to_csv(index=False).encode("utf-8"),
                file_name="mailing_list_{}.csv".format(date.today().isoformat()),
                mime="text/csv")
        st.markdown("---")

    # ── Search ───────────────────────────────────────────────────────────────
    df_hier = db.get_client_hierarchy()
    if df_hier.empty:
        st.markdown(
            '<div style="background:#001c4b04;border:2px dashed #001c4b18;padding:40px;'
            'text-align:center;margin-top:12px;">'
            '<div style="font-size:0.92rem;font-weight:700;color:{m};">Aucun client enregistré</div>'
            '<div style="color:#888;font-size:0.79rem;margin-top:4px;">'
            'Cliquez sur <b>Nouveau compte</b> ci-dessus pour commencer.</div>'
            '</div>'.format(m=MARINE), unsafe_allow_html=True)
    else:
        noms_clients = sorted(df_hier["nom_client"].tolist())
        crm_search   = st.selectbox(
            "Rechercher un compte client…",
            options=noms_clients,
            index=None,
            placeholder="Tapez le nom du client pour rechercher…",
            key="crm_search_box")

        if crm_search is None:
            nb_total   = len(df_hier)
            nb_valide  = int((df_hier["kyc_status"] == "Validé").sum())
            nb_encours = int((df_hier["kyc_status"] == "En cours").sum())
            nb_bloque  = int((df_hier["kyc_status"] == "Bloqué").sum())
            st.markdown(
                '<div style="display:grid;grid-template-columns:repeat(4,1fr);gap:8px;margin:16px 0;">'
                '<div class="kpi-card kpi-card-static"><div class="kpi-label">Comptes</div>'
                '<div class="kpi-value">{t}</div><div class="kpi-sub">clients enregistrés</div></div>'
                '<div class="kpi-card kpi-card-static"><div class="kpi-label">KYC Validé</div>'
                '<div class="kpi-value" style="color:#22a062;">{v}</div><div class="kpi-sub">&nbsp;</div></div>'
                '<div class="kpi-card kpi-card-static"><div class="kpi-label">KYC En cours</div>'
                '<div class="kpi-value" style="color:{o};">{e}</div><div class="kpi-sub">&nbsp;</div></div>'
                '<div class="kpi-card kpi-card-static"><div class="kpi-label">KYC Bloqué</div>'
                '<div class="kpi-value" style="color:#c0392b;">{b}</div><div class="kpi-sub">&nbsp;</div></div>'
                '</div>'.format(t=nb_total, v=nb_valide, e=nb_encours, b=nb_bloque, o=ORANGE),
                unsafe_allow_html=True)
            st.markdown(
                '<div style="color:#888;font-size:0.78rem;text-align:center;padding:8px 0;">'
                'Sélectionnez un compte dans la liste déroulante pour afficher sa fiche.</div>',
                unsafe_allow_html=True)
        else:
            # ── FICHE CLIENT — SPLIT SCREEN ───────────────────────────────
            sel_row = df_hier[df_hier["nom_client"] == crm_search]
            if sel_row.empty:
                st.warning("Client introuvable.")
            else:
                sel      = sel_row.iloc[0]
                sel_id   = int(sel["id"])
                sel_nom  = str(sel["nom_client"])
                sel_tier = str(sel.get("tier", "Tier 2"))
                sel_kyc  = str(sel.get("kyc_status", "En cours"))
                sel_type = str(sel.get("type_client", ""))
                sel_reg  = str(sel.get("region", ""))
                sel_cty  = str(sel.get("country", ""))
                sel_par  = str(sel.get("parent_nom", ""))
                sel_prod = str(sel.get("product_interests", ""))
                kyc_color = {"Validé": "#22a062", "En cours": ORANGE,
                             "Bloqué": "#c0392b"}.get(sel_kyc, "#aaa")

                # Fetch filiales
                conn_tmp = db.get_connection()
                try:
                    _c = conn_tmp.cursor()
                    _c.execute("SELECT nom_client FROM clients WHERE parent_id = ? ORDER BY nom_client", (sel_id,))
                    filiales = [r[0] for r in _c.fetchall()]
                finally:
                    conn_tmp.close()

                # ── SPLIT SCREEN ──────────────────────────────────────────
                col_gauche, col_droite = st.columns([1.2, 1], gap="large")

                # ════════════════════════════════════════════════════════
                # COLONNE GAUCHE — Entité
                # ════════════════════════════════════════════════════════
                with col_gauche:
                    # Product interest tags
                    prods_html = ""
                    if sel_prod.strip():
                        tags = " ".join(
                            '<span style="background:{c}18;border:1px solid {c}44;color:{m};'
                            'padding:1px 8px;font-size:0.68rem;font-weight:600;margin:2px;">{p}</span>'
                            .format(c=CIEL, m=MARINE, p=p.strip())
                            for p in sel_prod.split(",") if p.strip())
                        prods_html = (
                            '<div style="margin-top:10px;font-size:0.74rem;">'
                            '<span style="color:#888;margin-right:6px;">Product Interests :</span>'
                            + tags + '</div>')

                    st.markdown(
                        '<div class="detail-panel" style="margin-bottom:10px;">'
                        '<div style="display:flex;align-items:center;gap:10px;flex-wrap:wrap;margin-bottom:10px;">'
                        '<span style="font-size:1.05rem;font-weight:800;color:{m};">{nom}</span>'
                        '{tier_b}'
                        '<span style="background:{kycc};color:#fff;font-size:0.68rem;font-weight:700;padding:2px 10px;">'
                        '{kycd} KYC : {kyc}</span>'
                        '</div>'
                        '<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;font-size:0.79rem;">'
                        '<div><div style="color:#888;font-size:0.67rem;text-transform:uppercase;letter-spacing:0.5px;">Type</div>'
                        '<div style="font-weight:600;color:{m};">{type}</div></div>'
                        '<div><div style="color:#888;font-size:0.67rem;text-transform:uppercase;letter-spacing:0.5px;">Région</div>'
                        '<div style="font-weight:600;color:{m};">{reg}</div></div>'
                        '<div><div style="color:#888;font-size:0.67rem;text-transform:uppercase;letter-spacing:0.5px;">Tier</div>'
                        '<div style="font-weight:600;color:{m};">{tier}</div></div>'
                        '<div><div style="color:#888;font-size:0.67rem;text-transform:uppercase;letter-spacing:0.5px;">Country</div>'
                        '<div style="font-weight:600;color:{m};">{cty}</div></div>'
                        '</div>'
                        '{prods_html}'
                        '</div>'.format(
                            m=MARINE, nom=sel_nom, tier_b=_tier_badge(sel_tier), tier=sel_tier,
                            kycc=kyc_color, kycd=_kyc_dot(sel_kyc), kyc=sel_kyc,
                            type=sel_type, reg=sel_reg, cty=sel_cty or "—",
                            prods_html=prods_html),
                        unsafe_allow_html=True)

                    _btn_row1, _btn_row2 = st.columns(2)
                    with _btn_row1:
                        if st.button("Modifier le compte",
                                     key="crm_edit_co_{}".format(sel_id),
                                     type="tertiary", use_container_width=True):
                            dialog_edit_client(sel.to_dict())
                    with _btn_row2:
                        if st.button("Meeting Brief",
                                     key="crm_brief_{}".format(sel_id),
                                     type="tertiary", use_container_width=True):
                            _grp = db.get_client_group_summary(sel_id)
                            _ws  = db.get_whitespace_matrix()
                            dialog_meeting_brief(sel.to_dict(), _grp,
                                                 db.get_contacts(sel_id),
                                                 db.get_activities(client_id=sel_id),
                                                 _ws)

                    # PPTX export → Universal Export Hub in sidebar

                    # ── Corporate Structure block ─────────────────────────
                    _grp_summary = db.get_client_group_summary(sel_id)
                    has_struct = sel_par or filiales
                    if has_struct:
                        st.markdown(
                            '<div style="font-size:0.72rem;font-weight:700;color:{m};'
                            'text-transform:uppercase;letter-spacing:0.6px;'
                            'margin:14px 0 6px 0;border-bottom:1px solid #001c4b15;padding-bottom:3px;">'
                            'Structure du groupe</div>'.format(m=MARINE),
                            unsafe_allow_html=True)
                        if sel_par:
                            st.markdown(
                                '<div class="struct-block">'
                                '<div class="struct-label">Filiale de</div>'
                                '<div class="struct-val">{p}</div>'
                                '</div>'.format(p=sel_par),
                                unsafe_allow_html=True)
                        if filiales:
                            # Enrich filiales with AUM + Tier from group summary
                            _fils_enriched = _grp_summary.get("filiales", [])
                            _fils_dict = {f["nom"]: f for f in _fils_enriched}
                            fils_items = ""
                            for f_nom in filiales:
                                f_data = _fils_dict.get(f_nom, {})
                                f_aum  = fmt_m(f_data.get("aum", 0))
                                f_tier = f_data.get("tier", "—")
                                fils_items += (
                                    '<div style="display:flex;justify-content:space-between;'
                                    'align-items:center;padding:5px 0;'
                                    'border-bottom:1px solid #001c4b08;">'
                                    '<span style="font-weight:600;color:{m};font-size:0.78rem;">{f}</span>'
                                    '<span style="font-size:0.70rem;color:#888;">{tier}</span>'
                                    '<span style="font-size:0.74rem;font-weight:700;color:{ciel};">{aum}</span>'
                                    '</div>').format(m=MARINE, ciel=CIEL,
                                                      f=f_nom, tier=f_tier, aum=f_aum)
                            # Consolidated total
                            aum_grp = fmt_m(_grp_summary.get("aum_consolide", 0))
                            st.markdown(
                                '<div class="struct-block">'
                                '<div style="display:flex;justify-content:space-between;'
                                'margin-bottom:6px;">'
                                '<div class="struct-label">Filiales ({n})</div>'
                                '<div style="font-size:0.72rem;font-weight:700;color:{ciel};">'
                                'Groupe : {aum_grp}</div></div>'
                                '{fils}'
                                '</div>'.format(n=len(filiales), fils=fils_items,
                                              ciel=CIEL, aum_grp=aum_grp),
                                unsafe_allow_html=True)

                    # ── Network Graph (Pappers-style) ────────────────────
                    _net = db.get_client_network(sel_id)
                    _net_nodes = []
                    _net_edges = []
                    _seen_funds = set()
                    if _net.get("root"):
                        _root = _net["root"]
                        _is_root = (_root["id"] == sel_id)
                        _net_nodes.append(Node(
                            id="c_{}".format(_root["id"]),
                            label=_root["nom"],
                            size=30,
                            color=MARINE,
                            font={"color": "#ffffff", "size": 12, "bold": True},
                            shape="dot",
                        ))
                        for sub in _net.get("subsidiaries", []):
                            _net_nodes.append(Node(
                                id="c_{}".format(sub["id"]),
                                label=sub["nom_client"],
                                size=20,
                                color=CIEL,
                                font={"color": MARINE, "size": 10},
                                shape="dot",
                            ))
                            _net_edges.append(Edge(
                                source="c_{}".format(_root["id"]),
                                target="c_{}".format(sub["id"]),
                                color="#cccccc", width=1.5,
                            ))
                        for fl in _net.get("fund_links", []):
                            fund_id = "f_{}".format(fl["fonds"].replace(" ", "_"))
                            if fund_id not in _seen_funds:
                                _seen_funds.add(fund_id)
                                _net_nodes.append(Node(
                                    id=fund_id,
                                    label=fl["fonds"],
                                    size=15,
                                    color=ORANGE,
                                    font={"color": MARINE, "size": 9},
                                    shape="diamond",
                                ))
                            _edge_color = ORANGE if fl["statut"] == "Funded" else "#cccccc"
                            _edge_style = {"strokeDasharray": ""} if fl["statut"] == "Funded" else {}
                            _net_edges.append(Edge(
                                source="c_{}".format(fl["client_id"]),
                                target=fund_id,
                                color=_edge_color,
                                width=2 if fl["statut"] == "Funded" else 1,
                                dashes=fl["statut"] != "Funded",
                            ))

                    if _net_nodes:
                        st.markdown(
                            '<div style="font-size:0.72rem;font-weight:700;color:{m};'
                            'text-transform:uppercase;letter-spacing:0.6px;'
                            'margin:14px 0 6px 0;border-bottom:1px solid #001c4b15;padding-bottom:3px;">'
                            'Network Graph</div>'.format(m=MARINE),
                            unsafe_allow_html=True)
                        _graph_cfg = Config(
                            width="100%", height=300,
                            directed=False,
                            physics=True,
                            hierarchical=False,
                            nodeHighlightBehavior=True,
                            highlightColor=CIEL,
                            collapsible=False,
                        )
                        agraph(nodes=_net_nodes, edges=_net_edges, config=_graph_cfg)
                        st.markdown(
                            '<div style="font-size:0.65rem;color:#999;margin-top:2px;">'
                            '<span style="color:{m};">●</span> Parent &nbsp;'
                            '<span style="color:{c};">●</span> Filiale &nbsp;'
                            '<span style="color:{o};">◆</span> Fonds</div>'.format(
                                m=MARINE, c=CIEL, o=ORANGE),
                            unsafe_allow_html=True)

                    # ── Activités Récentes ────────────────────────────────
                    st.markdown(
                        '<div style="display:flex;align-items:center;justify-content:space-between;'
                        'margin:14px 0 6px 0;border-bottom:1px solid #001c4b15;padding-bottom:3px;">'
                        '<span style="font-size:0.72rem;font-weight:700;color:{m};'
                        'text-transform:uppercase;letter-spacing:0.6px;">Activités récentes</span>'
                        '</div>'.format(m=MARINE),
                        unsafe_allow_html=True)
                    df_act_client = db.get_activities(client_id=sel_id)
                    if df_act_client.empty:
                        st.markdown(
                            '<div style="color:#888;font-size:0.78rem;padding:6px 0;">Aucune activité enregistrée.</div>',
                            unsafe_allow_html=True)
                    else:
                        TYPE_COLORS_CRM = {"Call":CIEL,"Meeting":B_MID,"Email":B_PAL,
                                           "Roadshow":"#2c7fb8","Conference":"#004f8c","Autre":"#888"}
                        for _, act_row in df_act_client.head(3).iterrows():
                            a_typ   = str(act_row.get("type_interaction",""))
                            a_color = TYPE_COLORS_CRM.get(a_typ, "#888")
                            a_notes = str(act_row.get("notes","")) or "—"
                            a_date  = str(act_row.get("date",""))
                            st.markdown(
                                "<div style='border-left:3px solid {c};padding:6px 11px;"
                                "margin:4px 0;background:#f9fbfd;font-size:0.77rem;'>"
                                "<div style='display:flex;justify-content:space-between;'>"
                                "<span style='font-weight:700;color:{c};'>{typ}</span>"
                                "<span style='color:#aaa;font-size:0.68rem;'>{dt}</span></div>"
                                "<div style='color:#444;margin-top:3px;line-height:1.5;'>{notes}</div>"
                                "</div>".format(c=a_color, typ=a_typ, dt=a_date, notes=a_notes[:120]),
                                unsafe_allow_html=True)
                    if st.button("+ Enregistrer une activité", key="crm_add_act_{}".format(sel_id),
                                 type="tertiary"):
                        dialog_add_activity(preselect_client_id=sel_id)

                # ── Client-Specific Pipeline Editor ──────────────────────
                st.markdown(
                    '<div style="font-size:0.72rem;font-weight:700;color:{m};'
                    'text-transform:uppercase;letter-spacing:0.6px;'
                    'margin:18px 0 8px 0;border-bottom:1px solid #001c4b15;padding-bottom:3px;">'
                    'Pipeline — {nom}</div>'.format(m=MARINE, nom=sel_nom),
                    unsafe_allow_html=True)
                _df_client_pipe = db.get_pipeline_with_clients().copy()
                _df_client_pipe = _df_client_pipe[_df_client_pipe["client_id"] == sel_id].reset_index(drop=True)
                if _df_client_pipe.empty:
                    st.markdown(
                        '<div style="color:#888;font-size:0.78rem;padding:6px 0;">'
                        'Aucun deal dans le pipeline pour ce client.</div>',
                        unsafe_allow_html=True)
                    if st.button("+ Nouveau deal", key="crm_pipe_new_{}".format(sel_id),
                                 type="tertiary"):
                        dialog_add_deal(preselect_client_id=sel_id)
                else:
                    _df_cp_edit = _df_client_pipe.copy()
                    _df_cp_edit["next_action_date"] = _df_cp_edit["next_action_date"].apply(
                        lambda d: d.isoformat() if isinstance(d, date) else "")
                    _df_cp_edit["aum_m"] = _df_cp_edit.apply(
                        lambda r: round(float(r["funded_aum"]) / 1e6, 2) if r["statut"] == "Funded"
                        else round((float(r["revised_aum"]) if float(r["revised_aum"]) > 0
                              else float(r["target_aum_initial"])) / 1e6, 2), axis=1)
                    if "closing_probability" in _df_cp_edit.columns:
                        _df_cp_edit["closing_probability"] = _df_cp_edit["closing_probability"].fillna(50)

                    _cp_cols = ["id", "fonds", "statut", "aum_m",
                                "closing_probability", "next_action_date", "sales_owner"]
                    _cp_edited = st.data_editor(
                        _df_cp_edit[_cp_cols],
                        use_container_width=True,
                        hide_index=True,
                        num_rows="fixed",
                        height=min(250, 46 + len(_df_cp_edit) * 36),
                        column_config={
                            "id":                  st.column_config.NumberColumn("ID", width="small", disabled=True),
                            "fonds":               st.column_config.SelectboxColumn("Fonds", options=FONDS),
                            "statut":              st.column_config.SelectboxColumn("Statut", options=STATUTS),
                            "aum_m":               st.column_config.NumberColumn("AUM (M€)", format="%.2f",
                                                       min_value=0.0, step=1.0),
                            "closing_probability": st.column_config.NumberColumn("Proba %", format="%.0f%%",
                                                       width="small", min_value=0, max_value=100, step=5),
                            "next_action_date":    st.column_config.TextColumn("Next Action"),
                            "sales_owner":         st.column_config.TextColumn("Commercial"),
                        },
                        key="crm_client_pipe_editor_{}".format(sel_id))

                    _cp_changes = []
                    _cp_orig = _df_cp_edit[_cp_cols]
                    for idx in range(len(_cp_edited)):
                        if idx >= len(_cp_orig):
                            break
                        ed = _cp_edited.iloc[idx]
                        orig = _cp_orig.iloc[idx]
                        for c in ["statut", "fonds", "aum_m", "closing_probability",
                                   "next_action_date", "sales_owner"]:
                            if str(ed.get(c, "")) != str(orig.get(c, "")):
                                _cp_changes.append(idx)
                                break

                    _cp_btn_c1, _cp_btn_c2 = st.columns([3, 1])
                    with _cp_btn_c1:
                        _cp_label = "Sauvegarder les deals du client"
                        if _cp_changes:
                            _cp_label += " ({})".format(len(_cp_changes))
                        if st.button(_cp_label, key="crm_pipe_save_{}".format(sel_id),
                                     type="primary", use_container_width=True,
                                     disabled=(len(_cp_changes) == 0)):
                            _cp_ok, _cp_err = 0, 0
                            for idx in _cp_changes:
                                ed = _cp_edited.iloc[idx]
                                pid = int(_cp_orig.iloc[idx]["id"])
                                _orig_row = db.get_pipeline_row_by_id(pid)
                                if not _orig_row:
                                    _cp_err += 1
                                    continue
                                _upd = dict(_orig_row)
                                _upd["id"] = pid
                                for c in ["statut", "fonds"]:
                                    _upd[c] = str(ed.get(c, _upd.get(c, "")))
                                try:
                                    _upd["closing_probability"] = float(ed.get("closing_probability", 50))
                                except (ValueError, TypeError):
                                    pass
                                _nad = str(ed.get("next_action_date", "")).strip()
                                if _nad:
                                    _upd["next_action_date"] = _nad
                                _so = str(ed.get("sales_owner", "")).strip()
                                if _so:
                                    _upd["sales_owner"] = _so
                                _new_aum = float(ed.get("aum_m", 0)) * 1e6
                                if _orig_row["statut"] == "Funded" or str(ed.get("statut", "")) == "Funded":
                                    _upd["funded_aum"] = _new_aum
                                else:
                                    _upd["revised_aum"] = _new_aum
                                ok, msg = db.update_pipeline_row(_upd)
                                if ok:
                                    _cp_ok += 1
                                else:
                                    _cp_err += 1
                                    st.warning("Deal #{}: {}".format(pid, msg))
                            if _cp_ok:
                                st.success("{} deal(s) mis à jour.".format(_cp_ok))
                            if _cp_err:
                                st.error("{} erreur(s).".format(_cp_err))
                            if _cp_ok:
                                st.rerun()
                    with _cp_btn_c2:
                        if st.button("+ Deal", key="crm_pipe_add_{}".format(sel_id),
                                     type="tertiary", use_container_width=True):
                            dialog_add_deal(preselect_client_id=sel_id)

                # ════════════════════════════════════════════════════════
                # COLONNE DROITE — Contacts
                # ════════════════════════════════════════════════════════
                with col_droite:
                    df_contacts = db.get_contacts(sel_id)

                    # Header with inline "Ajouter" button
                    h_left, h_right = st.columns([3, 1])
                    with h_left:
                        st.markdown(
                            '<div style="font-size:0.72rem;font-weight:700;color:{m};'
                            'text-transform:uppercase;letter-spacing:0.6px;">'
                            'Contacts associés ({n})</div>'.format(m=MARINE, n=len(df_contacts)),
                            unsafe_allow_html=True)
                    with h_right:
                        if st.button("+ Ajouter", key="crm_add_ct_{}".format(sel_id),
                                     use_container_width=True):
                            dialog_add_contact(sel_id, sel_nom)

                    if df_contacts.empty:
                        st.markdown(
                            '<div style="color:#888;font-size:0.78rem;padding:10px 0;">'
                            'Aucun contact enregistré.</div>',
                            unsafe_allow_html=True)
                    else:
                        for _, ct in df_contacts.iterrows():
                            prenom   = str(ct.get("prenom", ""))
                            ct_nom   = str(ct.get("nom", ""))
                            role     = str(ct.get("role", ""))
                            email    = str(ct.get("email", "")).strip()
                            tel      = str(ct.get("telephone", "")).strip()
                            linkedin = str(ct.get("linkedin", "")).strip()
                            primary  = bool(ct.get("is_primary", 0))
                            ct_id    = int(ct.get("id", 0))

                            email_html = (
                                '<a href="mailto:{e}" style="color:{ciel};text-decoration:none;">{e}</a>'
                                .format(e=email, ciel=CIEL)) if email else '<span style="color:#bbb;">—</span>'
                            li_url = linkedin
                            if li_url and not li_url.startswith("http"):
                                li_url = "https://" + li_url
                            li_html = (
                                '<a href="{url}" target="_blank" style="color:{ciel};text-decoration:none;">LinkedIn</a>'
                                .format(url=li_url, ciel=CIEL)) if linkedin else ""
                            primary_badge = (
                                '<span style="background:{o};color:#fff;font-size:0.60rem;'
                                'font-weight:700;padding:1px 7px;margin-left:6px;">Principal</span>'
                                .format(o=ORANGE)) if primary else ""

                            # Contact card
                            st.markdown(
                                '<div style="background:#f4f8fc;border-left:3px solid {ciel};'
                                'padding:9px 12px;margin:5px 0;">'
                                '<div style="display:flex;align-items:center;gap:4px;margin-bottom:4px;">'
                                '<span style="font-weight:700;font-size:0.84rem;color:{m};">'
                                '{prenom} {nom}</span>{primary}'
                                '<span style="color:#888;font-size:0.70rem;margin-left:auto;">{role}</span>'
                                '</div>'
                                '<div style="font-size:0.74rem;line-height:1.8;">'
                                '<span style="color:#666;margin-right:12px;">Tél : {tel}</span>'
                                '{email}&nbsp;&nbsp;{li}'
                                '</div></div>'.format(
                                    ciel=CIEL, m=MARINE,
                                    prenom=prenom, nom=ct_nom,
                                    primary=primary_badge, role=role,
                                    tel=tel if tel else "—",
                                    email=email_html, li=li_html),
                                unsafe_allow_html=True)

                            # Action buttons — compact [name+role | Modifier | Supprimer]
                            btn_sp, btn_mod, btn_del = st.columns([8, 1, 1])
                            with btn_mod:
                                if st.button("Modifier", key="crm_edit_ct_{}".format(ct_id),
                                             type="tertiary", help="Modifier ce contact"):
                                    dialog_edit_contact(ct.to_dict())
                            with btn_del:
                                if st.button("Supprimer", key="crm_del_ct_{}".format(ct_id),
                                             type="tertiary", help="Supprimer ce contact"):
                                    st.session_state["crm_del_ct_confirm"] = ct_id

                            if st.session_state.get("crm_del_ct_confirm") == ct_id:
                                st.warning("Supprimer {} {} ?".format(prenom, ct_nom))
                                cdc1, cdc2, _ = st.columns([1, 1, 6])
                                with cdc1:
                                    if st.button("Confirmer", key="crm_del_ct_yes_{}".format(ct_id)):
                                        ok, err = db.delete_contact(ct_id)
                                        if ok:
                                            st.session_state.pop("crm_del_ct_confirm", None)
                                            st.rerun()
                                        else:
                                            st.error(err)
                                with cdc2:
                                    if st.button("Annuler", key="crm_del_ct_no_{}".format(ct_id)):
                                        st.session_state.pop("crm_del_ct_confirm", None)
                                        st.rerun()


# ============================================================================
# ONGLET 2 — PIPELINE MANAGEMENT
# ============================================================================
with tab_pipeline:
    st.markdown('<div class="section-title">Pipeline Management</div>',
                unsafe_allow_html=True)

    # ── Quick Action Center — 4 primary buttons, equal width ────────────────
    qac1, qac2, qac3, qac4 = st.columns(4)
    with qac1:
        if st.button("Nouveau Compte", key="qac_new_compte", use_container_width=True, type="primary"):
            dialog_add_client()
    with qac2:
        if st.button("Nouveau Deal", key="qac_new_deal", use_container_width=True, type="primary"):
            dialog_add_deal()
    with qac3:
        if st.button("Nouvelle Activité", key="qac_new_activity", use_container_width=True, type="primary"):
            dialog_add_activity()
    with qac4:
        if st.button("Nouveau Commercial", key="qac_new_sales", use_container_width=True, type="primary"):
            dialog_manage_sales()
    st.markdown("<br>", unsafe_allow_html=True)

    # ── Filtres ───────────────────────────────────────────────────────────────
    tp_filt, tp_spacer = st.columns([1, 7])
    with tp_filt:
        if st.button("Filtres", key="pipe_filt_toggle", use_container_width=True, type="secondary"):
            st.session_state["pipe_show_filters"] = not st.session_state.get("pipe_show_filters", False)

    if st.session_state.get("pipe_show_filters", False):
        _dyn = db.get_dynamic_filters()
        _dyn_statuts  = _dyn.get("statuts", STATUTS)
        _dyn_fonds    = _dyn.get("fonds",   FONDS)
        _dyn_regions  = _dyn.get("regions", REGIONS)
        _default_statuts = [s for s in STATUTS_ACTIFS if s in _dyn_statuts]
        fc1, fc2, fc3, fc4 = st.columns(4)
        with fc1: filt_statuts = st.multiselect("Statuts", _dyn_statuts,
                                                 default=_default_statuts,
                                                 key="pipe_filt_statuts")
        with fc2: filt_fonds   = st.multiselect("Fonds",   _dyn_fonds,
                                                 key="pipe_filt_fonds")
        with fc3: filt_regions = st.multiselect("Régions", _dyn_regions,
                                                 key="pipe_filt_regions")
        with fc4: filt_countries = st.multiselect("Pays", COUNTRIES_LIST[1:],
                                                   key="pipe_filt_countries")
    else:
        filt_statuts   = STATUTS_ACTIFS
        filt_fonds     = []
        filt_regions   = []
        filt_countries = []

    df_pipe_full = db.get_pipeline_with_last_activity()
    df_view      = df_pipe_full.copy()
    if filt_statuts:   df_view = df_view[df_view["statut"].isin(filt_statuts)]
    if filt_fonds:     df_view = df_view[df_view["fonds"].isin(filt_fonds)]
    if filt_regions:   df_view = df_view[df_view["region"].isin(filt_regions)]
    if filt_countries: df_view = df_view[df_view["country"].isin(filt_countries)]

    st.markdown('<div class="pipeline-hint">'
                '<b>{} deal(s)</b> affiché(s) — éditez directement dans le tableau, '
                'cochez 🗑️ pour supprimer, puis cliquez <b>Sauvegarder</b></div>'.format(len(df_view)),
                unsafe_allow_html=True)

    df_display = df_view.copy()
    df_display["next_action_date"] = df_display["next_action_date"].apply(
        lambda d: d.isoformat() if isinstance(d, date) else "")
    df_display["aum_pipeline"] = df_display.apply(
        lambda r: float(r["funded_aum"]) if r["statut"] == "Funded"
        else (float(r["revised_aum"]) if float(r["revised_aum"]) > 0
              else float(r["target_aum_initial"])), axis=1)
    df_display["aum_pipeline_fmt"] = df_display["aum_pipeline"].apply(fmt_m)
    df_display["funded_aum_fmt"]   = df_display["funded_aum"].apply(fmt_m)
    if "closing_probability" in df_display.columns:
        df_display["closing_probability"] = df_display["closing_probability"].fillna(50)

    df_display.insert(0, "\U0001f5d1\ufe0f Delete", False)

    _cols_for_editor = ["\U0001f5d1\ufe0f Delete",
                        "id","nom_client","type_client","region","country",
                        "fonds","statut","aum_pipeline_fmt","funded_aum_fmt",
                        "closing_probability","raison_perte",
                        "next_action_date","sales_owner","derniere_activite"]

    edited_df = st.data_editor(
        df_display[_cols_for_editor],
        use_container_width=True, height=400, hide_index=True, num_rows="fixed",
        column_config={
            "\U0001f5d1\ufe0f Delete":  st.column_config.CheckboxColumn("\U0001f5d1\ufe0f", default=False,
                                            width="small"),
            "id":                  st.column_config.NumberColumn("ID", width="small", disabled=True),
            "nom_client":          st.column_config.TextColumn("Client", disabled=True),
            "type_client":         st.column_config.TextColumn("Type", width="small", disabled=True),
            "region":              st.column_config.TextColumn("Region", width="small", disabled=True),
            "country":             st.column_config.TextColumn("Pays", width="small", disabled=True),
            "fonds":               st.column_config.SelectboxColumn("Fonds", options=FONDS),
            "statut":              st.column_config.SelectboxColumn("Statut", options=STATUTS),
            "aum_pipeline_fmt":    st.column_config.TextColumn("AUM Pipeline", disabled=True),
            "funded_aum_fmt":      st.column_config.TextColumn("AUM Financé", disabled=True),
            "closing_probability": st.column_config.NumberColumn("Proba %", format="%.0f%%",
                                       width="small", min_value=0, max_value=100, step=5),
            "raison_perte":        st.column_config.SelectboxColumn("Raison",
                                       options=[""] + RAISONS_PERTE),
            "next_action_date":    st.column_config.TextColumn("Next Action"),
            "sales_owner":         st.column_config.TextColumn("Commercial"),
            "derniere_activite":   st.column_config.TextColumn("Dernière Activité", disabled=True),
        }, key="pipeline_editor")

    _del_col = "\U0001f5d1\ufe0f Delete"
    _ids_to_delete = []
    _rows_to_update = []
    _orig_df = df_display[_cols_for_editor]

    for idx in range(len(edited_df)):
        if idx >= len(_orig_df):
            break
        row_ed   = edited_df.iloc[idx]
        row_orig = _orig_df.iloc[idx]
        pid      = int(row_orig["id"])

        if row_ed[_del_col]:
            _ids_to_delete.append(pid)
            continue

        changed = False
        for c in ["statut", "fonds", "closing_probability", "raison_perte",
                   "next_action_date", "sales_owner"]:
            if str(row_ed.get(c, "")) != str(row_orig.get(c, "")):
                changed = True
                break
        if changed:
            _rows_to_update.append((pid, row_ed))

    _n_actions = len(_ids_to_delete) + len(_rows_to_update)
    _btn_label = "Sauvegarder les modifications"
    if _n_actions > 0:
        _parts = []
        if _rows_to_update:
            _parts.append("{} modif.".format(len(_rows_to_update)))
        if _ids_to_delete:
            _parts.append("{} suppr.".format(len(_ids_to_delete)))
        _btn_label += " ({})".format(" + ".join(_parts))

    if st.button(_btn_label, key="pipe_save_all", type="primary",
                 use_container_width=True, disabled=(_n_actions == 0)):
        _ok_del, _ok_upd, _err = 0, 0, 0

        if _ids_to_delete:
            _ok_del = db.delete_pipeline_rows(_ids_to_delete)

        for pid, row_ed in _rows_to_update:
            _orig_row = db.get_pipeline_row_by_id(pid)
            if not _orig_row:
                _err += 1
                continue
            _upd = dict(_orig_row)
            _upd["id"] = pid
            for c in ["statut", "fonds"]:
                _upd[c] = str(row_ed.get(c, _upd.get(c, "")))
            try:
                _upd["closing_probability"] = float(row_ed.get("closing_probability",
                                                                _upd.get("closing_probability", 50)))
            except (ValueError, TypeError):
                pass
            _upd["raison_perte"] = str(row_ed.get("raison_perte", "")) or ""
            _nad = str(row_ed.get("next_action_date", "")).strip()
            if _nad:
                _upd["next_action_date"] = _nad
            _so = str(row_ed.get("sales_owner", "")).strip()
            if _so:
                _upd["sales_owner"] = _so
            ok, msg = db.update_pipeline_row(_upd)
            if ok:
                _ok_upd += 1
            else:
                _err += 1
                st.warning("Deal #{} : {}".format(pid, msg))

        _msgs = []
        if _ok_upd:
            _msgs.append("{} mis à jour".format(_ok_upd))
        if _ok_del:
            _msgs.append("{} supprimé(s)".format(_ok_del))
        if _msgs:
            st.success(" · ".join(_msgs))
        if _err:
            st.error("{} erreur(s).".format(_err))
        if _ok_upd or _ok_del:
            st.rerun()

    st.divider()
    df_viz_raw = db.get_pipeline_with_clients()

    # ── VISUALISATION 1 : Funnel Chart — AUM Pipeline par Statut ─────────────
    st.markdown("#### Funnel Pipeline — AUM par Statut")
    _statut_funnel_order = ["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit"]
    _df_funnel = df_viz_raw[df_viz_raw["statut"].isin(_statut_funnel_order)].copy()
    _df_funnel["_aum_p"] = _df_funnel.apply(
        lambda r: float(r["revised_aum"]) if float(r.get("revised_aum", 0) or 0) > 0
                  else float(r.get("target_aum_initial", 0) or 0), axis=1)
    _funnel_agg = (
        _df_funnel.groupby("statut")["_aum_p"].sum()
        .reindex(_statut_funnel_order, fill_value=0.0)
        .reset_index()
    )
    _funnel_agg.columns = ["statut", "aum"]
    _funnel_agg = _funnel_agg[_funnel_agg["aum"] > 0]
    if not _funnel_agg.empty:
        _funnel_colors = {
            "Prospect":      "#9ecae1",
            "Initial Pitch": B_PAL,
            "Due Diligence": "#004f8c",
            "Soft Commit":   B_MID,
        }
        fig_funnel = go.Figure(go.Funnel(
            y=_funnel_agg["statut"].tolist(),
            x=_funnel_agg["aum"].tolist(),
            textinfo="value+percent initial",
            text=[fmt_m(v) for v in _funnel_agg["aum"]],
            texttemplate="%{text}",
            hovertemplate="<b>%{y}</b><br>AUM Pipeline : %{text}<extra></extra>",
            marker_color=[_funnel_colors.get(s, B_MID) for s in _funnel_agg["statut"]],
            connector=dict(line=dict(color=GRIS, width=1)),
        ))
        fig_funnel.update_layout(
            height=320, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
            font_color=MARINE, font_size=11,
            margin=dict(l=10, r=10, t=30, b=10),
            funnelmode="stack")
        st.plotly_chart(fig_funnel, use_container_width=True, config={"displayModeBar": False}, key="chart_funnel")
        st.caption("AUM Pipeline = AUM Revise si > 0, sinon AUM Cible. Deals actifs uniquement.")
    else:
        st.info("Aucun deal actif pour le graphique entonnoir.")

    # ── VISUALISATION 2 : Grouped Bar — AUM par Fonds et par Statut ──────────
    st.markdown("#### AUM Pipeline par Fonds et par Statut")
    _grp_horizon = st.radio(
        "Activite recente",
        ["Tout l'historique", "1 Mois", "1 Trimestre", "6 Mois", "1 An"],
        horizontal=True,
        key="grp_horizon_radio")
    _df_grp = df_viz_raw[df_viz_raw["statut"].isin(_statut_funnel_order)].copy()
    # Filtre EN ARRIÈRE : deals ayant eu une activite dans la periode passee
    if _grp_horizon != "Tout l'historique":
        _today_grp  = date.today()
        _lookback   = {"1 Mois": 30, "1 Trimestre": 91, "6 Mois": 182, "1 An": 365}[_grp_horizon]
        _cutoff_grp = _today_grp - timedelta(days=_lookback)
        # Récupère la date MAX d'activite par pipeline_id via derniere_activite du df complet
        _df_la = db.get_pipeline_with_last_activity()
        # Extraire la date depuis la colonne derniere_activite si disponible
        # Alternativement, utiliser une requête directe sur les activites
        _conn_g = db.get_connection()
        try:
            import pandas as _pd_g
            _df_act_dates = _pd_g.read_sql_query(
                "SELECT p.id AS pipeline_id, MAX(a.date) AS last_act_date"
                " FROM pipeline p"
                " LEFT JOIN activites a ON a.client_id = p.client_id"
                " GROUP BY p.id",
                _conn_g)
        finally:
            _conn_g.close()
        if not _df_act_dates.empty and "id" in _df_grp.columns:
            _df_act_dates["last_act_date"] = pd.to_datetime(
                _df_act_dates["last_act_date"], errors="coerce")
            _df_grp = _df_grp.merge(
                _df_act_dates.rename(columns={"pipeline_id": "id"}),
                on="id", how="left")
            _df_grp = _df_grp[
                _df_grp["last_act_date"].notna() &
                (_df_grp["last_act_date"].dt.date >= _cutoff_grp)
            ]
            _df_grp = _df_grp.drop(columns=["last_act_date"], errors="ignore")
    _df_grp["_aum_p"] = _df_grp.apply(
        lambda r: float(r["revised_aum"]) if float(r.get("revised_aum", 0) or 0) > 0
                  else float(r.get("target_aum_initial", 0) or 0), axis=1)
    _df_grp = _df_grp[_df_grp["_aum_p"] > 0]
    if not _df_grp.empty:
        _grp_pivot = (
            _df_grp.groupby(["fonds", "statut"])["_aum_p"].sum()
            .reset_index()
        )
        _fonds_list  = sorted(_grp_pivot["fonds"].unique().tolist())
        _statut_list = [s for s in _statut_funnel_order
                        if s in _grp_pivot["statut"].unique()]
        _grp_colors  = {
            "Prospect":      "#9ecae1",
            "Initial Pitch": B_PAL,
            "Due Diligence": "#004f8c",
            "Soft Commit":   B_MID,
        }
        fig_grp = go.Figure()
        for _st in _statut_list:
            _df_st = _grp_pivot[_grp_pivot["statut"] == _st]
            _st_map = dict(zip(_df_st["fonds"], _df_st["_aum_p"]))
            _y_vals = [_st_map.get(f, 0.0) for f in _fonds_list]
            _txt    = [fmt_m(v) if v > 0 else "" for v in _y_vals]
            fig_grp.add_trace(go.Bar(
                name=_st,
                x=_fonds_list,
                y=_y_vals,
                text=_txt,
                textposition="outside",
                textfont=dict(size=8, color=MARINE),
                marker_color=_grp_colors.get(_st, B_MID),
                marker_line_color=BLANC, marker_line_width=0.5,
                hovertemplate="<b>{s}</b><br>%{{x}}<br>AUM : %{{customdata}}<extra></extra>".format(
                    s=_st),
                customdata=_txt))
        fig_grp.update_layout(
            barmode="group", height=340,
            paper_bgcolor=BLANC, plot_bgcolor=BLANC,
            font_color=MARINE, bargap=0.18, bargroupgap=0.05,
            legend_bgcolor=BLANC, legend_bordercolor=GRIS, legend_borderwidth=1,
            legend_font_size=10,
            xaxis_showgrid=False, xaxis_tickangle=-15,
            yaxis_showgrid=True, yaxis_gridcolor=GRIS,
            margin=dict(l=10, r=10, t=30, b=10))
        st.plotly_chart(fig_grp, use_container_width=True, config={"displayModeBar": False}, key="chart_grp")
        st.caption("Vue croisee Fonds x Statut — deals actifs uniquement.")
    else:
        st.info("Aucune donnee disponible pour la vue croisee Fonds / Statut.")

    df_lp = db.get_pipeline_with_clients()
    df_lp = df_lp[df_lp["statut"].isin(["Lost","Paused"])].copy()
    if not df_lp.empty:
        st.divider()
        st.markdown("#### Lost / Paused Deals")
        df_lp2 = df_lp[["nom_client","fonds","statut","target_aum_initial",
                         "raison_perte","concurrent_choisi","sales_owner"]].copy()
        df_lp2["target_aum_initial"] = df_lp2["target_aum_initial"].apply(fmt_m)
        st.dataframe(df_lp2, use_container_width=True, hide_index=True)


# ============================================================================
# ONGLET 3 — EXECUTIVE DASHBOARD
# ============================================================================
with tab_dash:
    st.markdown('<div class="section-title">Executive Dashboard</div>',
                unsafe_allow_html=True)

    tf_dash = st.selectbox("Timeframe", TIMEFRAMES, key="tf_dash",
                           label_visibility="collapsed")
    cutoff_dash = _timeframe_cutoff(tf_dash)

    kpis = db.get_kpis()
    # pipeline_actif already uses Smart AUM (CASE WHEN revised_aum > 0 ...) in database.py
    if cutoff_dash is not None:
        import pandas as _pd
        _conn = db.get_connection()
        _cutoff_str = cutoff_dash.isoformat()
        _df_p = _pd.read_sql_query(
            "SELECT p.statut, p.funded_aum, p.revised_aum, p.target_aum_initial"
            " FROM pipeline p WHERE DATE(p.updated_at) >= ?",
            _conn, params=(_cutoff_str,))
        _conn.close()
        if not _df_p.empty:
            kpis["total_funded"]    = float(_df_p[_df_p["statut"]=="Funded"]["funded_aum"].sum())
            _df_actif = _df_p[_df_p["statut"].isin(
                ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])].copy()
            # Smart AUM: revised if >0, else target
            _df_actif["_aum_p"] = _df_actif.apply(
                lambda r: float(r["revised_aum"]) if float(r.get("revised_aum",0) or 0) > 0
                          else float(r.get("target_aum_initial",0) or 0), axis=1)
            kpis["pipeline_actif"] = float(_df_actif["_aum_p"].sum())
            kpis["nb_funded"]       = int((_df_p["statut"]=="Funded").sum())
            kpis["nb_lost"]         = int((_df_p["statut"]=="Lost").sum())
            kpis["nb_paused"]       = int((_df_p["statut"]=="Paused").sum())
            kpis["nb_deals_actifs"] = int(_df_p["statut"].isin(
                ["Prospect","Initial Pitch","Due Diligence","Soft Commit"]).sum())
            nb_fl = kpis["nb_funded"] + kpis["nb_lost"]
            kpis["taux_conversion"] = round(kpis["nb_funded"] / nb_fl * 100, 1) if nb_fl > 0 else 0.0
        else:
            for k in ("total_funded","pipeline_actif","nb_funded","nb_lost",
                      "nb_paused","nb_deals_actifs","taux_conversion"):
                kpis[k] = 0

    nb_lost_paused = kpis["nb_lost"] + kpis.get("nb_paused", 0)

    card_lp = (
        '<div class="kpi-card kpi-card-static">'
        '<div class="kpi-label">Lost / Paused Deals</div>'
        '<div class="kpi-value">{n}</div>'
        '<div class="kpi-sub">{n} deals</div>'
        '</div>'
    ).format(n=nb_lost_paused) if nb_lost_paused > 0 else (
        '<div class="kpi-card kpi-card-static">'
        '<div class="kpi-label">Lost / Paused Deals</div>'
        '<div class="kpi-value">0</div>'
        '<div class="kpi-sub">&nbsp;</div>'
        '</div>'
    )
    st.markdown(
        '<div class="kpi-grid">'
        '<div class="kpi-card kpi-card-static">'
        '<div class="kpi-label">Total Funded AUM</div>'
        '<div class="kpi-value">{aum_f}</div>'
        '<div class="kpi-sub">{nb_f} deal(s) funded</div>'
        '</div>'
        '<div class="kpi-card kpi-card-static">'
        '<div class="kpi-label">Active Pipeline</div>'
        '<div class="kpi-value">{aum_p}</div>'
        '<div class="kpi-sub">{nb_p} active deals</div>'
        '</div>'
        '<div class="kpi-card kpi-card-static">'
        '<div class="kpi-label">Weighted Pipeline</div>'
        '<div class="kpi-value">{wp}</div>'
        '<div class="kpi-sub">probability-weighted</div>'
        '</div>'
        '<div class="kpi-card kpi-card-static">'
        '<div class="kpi-label">Conversion Rate</div>'
        '<div class="kpi-value">{taux:.1f}%</div>'
        '<div class="kpi-sub">{nb_f2} funded / {nb_l} lost</div>'
        '</div>'
        '{card_lp}'
        '</div>'.format(
            aum_f=fmt_m(kpis["total_funded"]), nb_f=kpis["nb_funded"],
            aum_p=fmt_m(kpis["pipeline_actif"]), nb_p=kpis["nb_deals_actifs"],
            wp=fmt_m(kpis.get("weighted_pipeline", 0)),
            taux=kpis["taux_conversion"], nb_f2=kpis["nb_funded"], nb_l=kpis["nb_lost"],
            card_lp=card_lp),
        unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── AUM Time Machine — Historical Line Chart ─────────────────────────
    st.markdown("#### AUM Time Machine")
    _tm_options = {"YTD": None, "1M": 30, "3M": 91, "6M": 182, "1Y": 365}
    _tm_sel = st.selectbox("Période", list(_tm_options.keys()),
                            index=4, key="dash_time_machine_period",
                            label_visibility="collapsed")
    _tm_days = _tm_options[_tm_sel]
    if _tm_days is None:
        _tm_days = (date.today() - date(date.today().year, 1, 1)).days or 30
    _df_hist = db.get_historical_aum(days_back=_tm_days)
    if not _df_hist.empty:
        fig_tm = go.Figure()
        fig_tm.add_trace(go.Scatter(
            x=_df_hist["date"], y=_df_hist["funded_aum"],
            name="Funded AUM", mode="lines",
            line=dict(color=MARINE, width=2.5),
            fill="tozeroy", fillcolor="rgba(0,28,75,0.08)",
            hovertemplate="<b>Funded</b><br>%{x|%d/%m/%Y}<br>%{customdata}<extra></extra>",
            customdata=[fmt_m(v) for v in _df_hist["funded_aum"]]))
        fig_tm.add_trace(go.Scatter(
            x=_df_hist["date"], y=_df_hist["pipeline_aum"],
            name="Active Pipeline", mode="lines",
            line=dict(color=CIEL, width=2, dash="dot"),
            hovertemplate="<b>Pipeline</b><br>%{x|%d/%m/%Y}<br>%{customdata}<extra></extra>",
            customdata=[fmt_m(v) for v in _df_hist["pipeline_aum"]]))
        fig_tm.update_layout(
            height=300, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
            font_color=MARINE,
            legend=dict(orientation="h", y=1.08, x=0.5, xanchor="center",
                        font_size=10, bgcolor=BLANC, bordercolor=GRIS, borderwidth=1),
            xaxis=dict(showgrid=False, tickformat="%b %Y"),
            yaxis=dict(showgrid=True, gridcolor=GRIS, tickformat=".2s",
                       title="AUM (EUR)"),
            margin=dict(l=10, r=10, t=40, b=10))
        st.plotly_chart(fig_tm, use_container_width=True,
                        config={"displayModeBar": False}, key="chart_time_machine")
        st.caption("Evolution historique du Funded AUM et du Pipeline Actif. "
                   "Période : {}.".format(_tm_sel))
    else:
        st.info("Aucune donnée historique disponible.")

    st.markdown("<br>", unsafe_allow_html=True)

    statut_order = [s for s in STATUTS if kpis["statut_repartition"].get(s, 0) > 0]
    if statut_order:
        pills = ""
        for s in statut_order:
            c_hex = STATUT_COLORS.get(s, GRIS)
            count = kpis["statut_repartition"][s]
            pills += (
                '<div class="statut-pill" style="background:{c}16;border:1px solid {c}44;">'
                '<div style="font-size:0.61rem;color:{marine};font-weight:700;text-transform:uppercase;">{s}</div>'
                '<div style="font-size:1.3rem;font-weight:800;color:{c};">{n}</div>'
                '</div>'
            ).format(s=s, c=c_hex, marine=MARINE, n=count)
        st.markdown(
            '<div class="statut-grid" style="grid-template-columns:repeat({n},1fr);">'
            '{pills}</div>'.format(n=len(statut_order), pills=pills),
            unsafe_allow_html=True)

        with st.expander("Détail par statut", expanded=False):
            tabs_statut = st.tabs(["{} ({})".format(s, kpis["statut_repartition"][s])
                                   for s in statut_order])
            for tab_s, s in zip(tabs_statut, statut_order):
                with tab_s:
                    _content_statut(s, _filtre_effectif)

    # ── Alertes opérationnelles — filtrées avec les mêmes variables que le dashboard ──
    df_overdue = db.get_overdue_actions(fonds_filter=_filtre_effectif)
    # Appliquer filtre région si actif (même variable que le reste du dashboard)
    if not df_overdue.empty:
        _pipe_filt_regions_dash = st.session_state.get("pipe_filt_regions", [])
        _pipe_filt_statuts_dash = st.session_state.get("pipe_filt_statuts", [])
        if _pipe_filt_regions_dash and "region" in df_overdue.columns:
            df_overdue = df_overdue[df_overdue.get("region", pd.Series(dtype=str)).isin(
                _pipe_filt_regions_dash)] if "region" in df_overdue.columns else df_overdue
        # Pas de filtre statut sur les alertes (elles sont déjà exclues Lost/Funded/Redeemed)
    if not df_overdue.empty:
        today = date.today()
        alertes_html = ""
        for _, row in df_overdue.head(5).iterrows():
            nad = row.get("next_action_date")
            days_late = (today - nad).days if isinstance(nad, date) else 0
            nad_str   = nad.isoformat() if isinstance(nad, date) else "—"
            owner     = str(row.get("sales_owner","")) or ""
            alertes_html += (
                '<div class="alert-overdue">'
                '<b>{client}</b> — {fonds}'
                ' <span style="color:{ciel};font-weight:600;">({statut})</span>'
                ' — Prevue le <b>{nad}</b>'
                ' &nbsp;<span class="badge-retard">+{days}j</span>'
                '{owner_part}'
                '</div>'
            ).format(
                client=row["nom_client"], fonds=row["fonds"], ciel=CIEL,
                statut=row["statut"], nad=nad_str, days=days_late,
                owner_part=" — {}".format(owner) if owner else "")
        st.markdown(alertes_html, unsafe_allow_html=True)
        with st.expander("{} action(s) en retard — détail".format(len(df_overdue)), expanded=False):
            _content_overdue()

    st.divider()

    gcol1, gcol2, gcol3 = st.columns([1, 1, 1.2], gap="medium")
    with gcol1:
        st.markdown("#### Par Type Client")
        abt = kpis.get("aum_by_type", {})
        if abt:
            fig_type = go.Figure(go.Pie(
                labels=list(abt.keys()), values=list(abt.values()), hole=0.52,
                marker_colors=PALETTE[:len(abt)], marker_line_color=BLANC, marker_line_width=2,
                textinfo="percent", textfont_size=10, textfont_color=BLANC))
            fig_type.add_annotation(text=fmt_m(sum(abt.values())),
                                    x=0.5, y=0.55, font_size=11, font_color=MARINE, showarrow=False)
            fig_type.add_annotation(text="Finance",
                                    x=0.5, y=0.42, font_size=8, font_color=GTXT, showarrow=False)
            fig_type.update_layout(height=300, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
                                   font_color=MARINE, showlegend=True,
                                   legend_x=1.02, legend_y=0.5, legend_font_size=9,
                                   margin=dict(l=0, r=80, t=36, b=10))
            st.plotly_chart(fig_type, use_container_width=True, config={"displayModeBar": False}, key="chart_pie_type")
    with gcol2:
        st.markdown("#### Par Région")
        aum_reg_dash = db.get_aum_by_region()
        if aum_reg_dash:
            fig_reg = go.Figure(go.Pie(
                labels=list(aum_reg_dash.keys()), values=list(aum_reg_dash.values()), hole=0.52,
                marker_colors=PALETTE[:len(aum_reg_dash)], marker_line_color=BLANC, marker_line_width=2,
                textinfo="percent", textfont_size=10, textfont_color=BLANC))
            fig_reg.add_annotation(text=fmt_m(sum(aum_reg_dash.values())),
                                   x=0.5, y=0.55, font_size=11, font_color=MARINE, showarrow=False)
            fig_reg.add_annotation(text="Finance",
                                   x=0.5, y=0.42, font_size=8, font_color=GTXT, showarrow=False)
            fig_reg.update_layout(height=300, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
                                  font_color=MARINE, showlegend=True,
                                  legend_x=1.02, legend_y=0.5, legend_font_size=9,
                                  margin=dict(l=0, r=80, t=36, b=10))
            st.plotly_chart(fig_reg, use_container_width=True, config={"displayModeBar": False}, key="chart_pie_reg")
    with gcol3:
        st.markdown("#### AUM par Fonds")
        abf = kpis.get("aum_by_fonds", {})
        if abf:
            fs = sorted(abf.items(), key=lambda x: x[1], reverse=True)
            fig_fonds = go.Figure(go.Bar(
                x=[v for _, v in fs], y=[k for k, _ in fs], orientation="h",
                marker_color=PALETTE[:len(fs)], marker_line_color=BLANC, marker_line_width=0.5,
                text=[fmt_m(v) for _, v in fs], textposition="outside",
                textfont_size=9, textfont_color=MARINE))
            fig_fonds.update_layout(height=300, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
                                    font_color=MARINE, xaxis_showgrid=True, xaxis_gridcolor=GRIS,
                                    yaxis=dict(autorange="reversed", automargin=True),
                                    margin=dict(l=180, r=20, t=36, b=10))
            st.plotly_chart(fig_fonds, use_container_width=True, config={"displayModeBar": False}, key="chart_bar_fonds")

    # ── Global Roadshow Map — Choropleth ────────────────────────────────────
    st.divider()
    st.markdown("#### Global Roadshow Map — AUM par Pays")
    _df_country_aum = db.get_aum_by_country()
    if not _df_country_aum.empty and _df_country_aum["total_aum"].sum() > 0:
        _COUNTRY_ISO = {
            "United Arab Emirates": "ARE", "Saudi Arabia": "SAU", "Qatar": "QAT",
            "Kuwait": "KWT", "Bahrain": "BHR", "Oman": "OMN",
            "United Kingdom": "GBR", "France": "FRA", "Germany": "DEU",
            "Switzerland": "CHE", "Luxembourg": "LUX", "Netherlands": "NLD",
            "Italy": "ITA", "Spain": "ESP", "Belgium": "BEL", "Austria": "AUT",
            "Sweden": "SWE", "Norway": "NOR", "Denmark": "DNK", "Finland": "FIN",
            "Singapore": "SGP", "Japan": "JPN", "Hong Kong": "HKG",
            "China": "CHN", "South Korea": "KOR", "Australia": "AUS", "India": "IND",
            "United States": "USA", "Canada": "CAN", "Brazil": "BRA",
            "Mexico": "MEX", "South Africa": "ZAF", "Egypt": "EGY",
        }
        _df_map = _df_country_aum.copy()
        _df_map["iso"] = _df_map["country"].map(_COUNTRY_ISO)
        _df_map = _df_map.dropna(subset=["iso"])
        _df_map["hover_text"] = _df_map.apply(
            lambda r: "<b>{}</b><br>Funded: {}<br>Pipeline: {}<br>Total: {}".format(
                r["country"], fmt_m(r["funded_aum"]),
                fmt_m(r["pipeline_aum"]), fmt_m(r["total_aum"])), axis=1)
        if not _df_map.empty:
            fig_map = go.Figure(go.Choropleth(
                locations=_df_map["iso"],
                z=_df_map["total_aum"],
                text=_df_map["hover_text"],
                hovertemplate="%{text}<extra></extra>",
                colorscale=[
                    [0.0, "#f0f4f8"], [0.2, "#c6dff0"], [0.4, "#6baed6"],
                    [0.6, "#2171b5"], [0.8, "#1a5e8a"], [1.0, MARINE],
                ],
                marker_line_color="#ffffff",
                marker_line_width=0.5,
                colorbar=dict(
                    title=dict(text="AUM (EUR)", font=dict(size=9, color=MARINE)),
                    tickfont=dict(size=8, color=MARINE),
                    thickness=12, len=0.6,
                ),
            ))
            fig_map.update_layout(
                geo=dict(
                    showframe=False, showcoastlines=True,
                    coastlinecolor=GRIS, projection_type="natural earth",
                    bgcolor=BLANC, landcolor="#f8f9fa",
                    showlakes=False, showcountries=True, countrycolor="#e0e0e0",
                ),
                height=360, paper_bgcolor=BLANC,
                font_color=MARINE,
                margin=dict(l=0, r=0, t=10, b=10),
            )
            st.plotly_chart(fig_map, use_container_width=True,
                            config={"displayModeBar": False}, key="chart_choropleth")
            st.caption("AUM total (Funded + Pipeline actif) par pays. Couleur plus foncée = AUM plus élevé.")
    else:
        st.info("Aucune donnée géographique disponible pour la carte.")

    st.divider()
    _td_col1, _td_col2 = st.columns([3, 1])
    with _td_col1:
        st.markdown("#### Top Deals — Funded AUM")
    with _td_col2:
        _td_region = st.selectbox("Filtrer par Région", ["Toutes"] + REGIONS,
                                   key="dash_top_deals_region",
                                   label_visibility="collapsed")
    df_funded_top = (db.get_pipeline_with_clients()
                     .query("statut == 'Funded' and funded_aum > 0")
                     .sort_values("funded_aum", ascending=False))
    if _td_region != "Toutes":
        df_funded_top = df_funded_top[df_funded_top["region"] == _td_region]
    df_funded_top = df_funded_top.head(10)
    if not df_funded_top.empty:
        max_f = float(df_funded_top["funded_aum"].max())
        for i, (_, row) in enumerate(df_funded_top.iterrows()):
            val   = float(row["funded_aum"])
            pct   = val / max_f * 100 if max_f > 0 else 0
            bar_c = MARINE if i == 0 else B_MID if i < 3 else B_PAL
            st.markdown(
                '<div style="display:flex;align-items:center;gap:10px;'
                'margin:4px 0;padding:7px 12px;border-bottom:1px solid #e8e8e8;">'
                '<div style="font-size:0.72rem;font-weight:700;color:{ciel};min-width:26px;">No.{rank}</div>'
                '<div style="flex:1;">'
                '<div style="font-size:0.80rem;font-weight:700;color:{marine};">{client}</div>'
                '<div style="font-size:0.67rem;color:#777;">{fonds} &middot; {type} &middot; {region}</div>'
                '<div style="background:{gris};height:3px;margin-top:4px;overflow:hidden;">'
                '<div style="background:{barc};width:{pct:.0f}%;height:100%;"></div></div></div>'
                '<div style="text-align:right;min-width:72px;">'
                '<div style="font-size:0.86rem;font-weight:800;color:{marine};">{aum}</div>'
                '</div></div>'.format(
                    rank=i+1, ciel=CIEL, marine=MARINE, gris=GRIS,
                    client=row["nom_client"], fonds=row["fonds"],
                    type=row["type_client"], region=row["region"],
                    barc=bar_c, pct=pct, aum=fmt_m(val)), unsafe_allow_html=True)

    with st.expander("Détail des deals — Funded / Actifs / Perdus", expanded=False):
        dt1, dt2, dt3 = st.tabs(["Funded", "Pipeline Actif", "Lost & Paused"])
        with dt1: _content_funded(_filtre_effectif)
        with dt2: _content_pipeline(_filtre_effectif)
        with dt3: _content_lost(_filtre_effectif)

    # ── TOP CLIENTS CONSOLIDE — Stacked Bar horizontal par Fonds ─────────────
    st.divider()
    _strat_col1, _strat_col2 = st.columns([3, 1])
    with _strat_col1:
        st.markdown("#### Top Clients — Vue Consolidee (AUM Finance par Fonds)")
    with _strat_col2:
        _strat_region = st.selectbox("Filtrer par Région", ["Toutes"] + REGIONS,
                                      key="dash_strat_region",
                                      label_visibility="collapsed")
    _df_top_clients_raw = (
        db.get_pipeline_with_clients()
        .query("statut == 'Funded' and funded_aum > 0")
        .copy()
    )
    if _strat_region != "Toutes":
        _df_top_clients_raw = _df_top_clients_raw[
            _df_top_clients_raw["region"] == _strat_region]
    if not _df_top_clients_raw.empty:
        _df_tc = (
            _df_top_clients_raw
            .groupby(["nom_client", "fonds"])["funded_aum"]
            .sum()
            .reset_index()
        )
        # Identifier les top N clients par AUM Finance total
        _client_totals = (
            _df_tc.groupby("nom_client")["funded_aum"].sum()
            .sort_values(ascending=False)
            .head(12)
        )
        _top_clients = _client_totals.index.tolist()
        _df_tc = _df_tc[_df_tc["nom_client"].isin(_top_clients)].copy()
        # Ordonner les clients par AUM total decroissant (Y axis = haut en bas)
        _client_order = list(reversed(_top_clients))
        _fonds_present = sorted(_df_tc["fonds"].unique().tolist())
        fig_tc = go.Figure()
        for _fi, _fonds in enumerate(_fonds_present):
            _df_f = _df_tc[_df_tc["fonds"] == _fonds]
            _fmap = dict(zip(_df_f["nom_client"], _df_f["funded_aum"]))
            _x_vals = [_fmap.get(c, 0.0) for c in _client_order]
            _txt    = [fmt_m(v) if v > 0 else "" for v in _x_vals]
            fig_tc.add_trace(go.Bar(
                name=_fonds,
                x=_x_vals,
                y=_client_order,
                orientation="h",
                marker_color=PALETTE[_fi % len(PALETTE)],
                marker_line_color=BLANC, marker_line_width=0.4,
                text=_txt,
                textposition="inside",
                textfont=dict(size=8, color=BLANC),
                hovertemplate=(
                    "<b>%{{y}}</b><br>{fonds}<br>"
                    "AUM Finance : %{{customdata}}<extra></extra>").format(fonds=_fonds),
                customdata=_txt))
        fig_tc.update_layout(
            barmode="stack",
            height=max(300, 36 * len(_client_order) + 80),
            paper_bgcolor=BLANC, plot_bgcolor=BLANC,
            font_color=MARINE,
            legend_bgcolor=BLANC, legend_bordercolor=GRIS,
            legend_borderwidth=1, legend_font_size=10,
            xaxis_showgrid=True, xaxis_gridcolor=GRIS,
            xaxis_title="AUM Finance (EUR)",
            yaxis=dict(automargin=True, tickfont=dict(size=10)),
            margin=dict(l=10, r=20, t=30, b=10))
        st.plotly_chart(fig_tc, use_container_width=True, config={"displayModeBar": False}, key="chart_tc")
        st.caption(
            "AUM Finance consolide par client — top {} clients. "
            "Chaque segment represente la contribution d'un fonds aux encours totaux du client.".format(
                len(_top_clients)))
    else:
        st.markdown(
            '<div style="background:#001c4b04;border:1px dashed #001c4b20;'
            'padding:18px;text-align:center;">'
            '<div style="color:{m};font-size:0.84rem;">Aucun deal Funded enregistre.</div>'
            '</div>'.format(m=MARINE), unsafe_allow_html=True)

    # ── DECISION SUPPORT ──────────────────────────────────────────────────
    st.divider()
    ds_col1, ds_col2 = st.columns([1.1, 1], gap="large")

    with ds_col1:
        st.markdown("#### Money in Motion — Projected Inflows")
        df_cf = db.get_expected_cashflows()
        if df_cf.empty or df_cf["aum_pondere"].sum() == 0:
            st.markdown(
                '<div style="background:#001c4b04;border:1px dashed #001c4b20;'
                'padding:18px;text-align:center;">'
                '<div style="color:{m};font-size:0.84rem;">Aucun deal actif avec prochaine action planifiée.</div>'
                '</div>'.format(m=MARINE), unsafe_allow_html=True)
        else:
            all_months = sorted(df_cf["mois"].unique().tolist())
            fig_cf = go.Figure()
            for i, fonds in enumerate(sorted(df_cf["fonds"].unique())):
                df_f = df_cf[df_cf["fonds"] == fonds]
                month_map = dict(zip(df_f["mois"], df_f["aum_pondere"]))
                y_vals = [month_map.get(m, 0) for m in all_months]
                fig_cf.add_trace(go.Bar(
                    name=fonds, x=all_months, y=y_vals,
                    marker_color=PALETTE[i % len(PALETTE)],
                    marker_line_color=BLANC, marker_line_width=0.5,
                    hovertemplate="<b>{}</b><br>%{{x}}<br>%{{customdata}}<extra></extra>".format(fonds),
                    customdata=[fmt_m(v) for v in y_vals]))
            # ── Trace invisible pour afficher le Total au-dessus de chaque barre ──
            _cf_totaux = {}
            for _, _r in df_cf.iterrows():
                _cf_totaux[_r["mois"]] = _cf_totaux.get(_r["mois"], 0.0) + _r["aum_pondere"]
            _cf_total_vals = [_cf_totaux.get(m, 0.0) for m in all_months]
            fig_cf.add_trace(go.Scatter(
                x=all_months,
                y=_cf_total_vals,
                mode="text",
                text=[fmt_m(v) if v > 0 else "" for v in _cf_total_vals],
                textposition="top center",
                textfont=dict(size=9, color=MARINE),
                showlegend=False,
                hoverinfo="skip",
                name="Total"))
            fig_cf.update_layout(
                barmode="stack", height=280,
                paper_bgcolor=BLANC, plot_bgcolor=BLANC, font_color=MARINE,
                legend_bgcolor=BLANC, legend_bordercolor=GRIS, legend_borderwidth=1,
                legend_font_size=9,
                xaxis_showgrid=False, yaxis_showgrid=True, yaxis_gridcolor=GRIS,
                yaxis_tickformat=".2s", yaxis_title="AUM Pondéré (€)",
                margin=dict(l=10, r=10, t=28, b=10))
            st.plotly_chart(fig_cf, use_container_width=True,
                            config={"displayModeBar": False}, key="chart_cf")
            st.caption("AUM pondéré = AUM Pipeline × Probabilité de closing. Répartition par mois de Next Action.")

    with ds_col2:
        st.markdown("#### Whitespace Analysis — Cross-Sell")
        df_ws = db.get_whitespace_matrix()
        if df_ws.empty:
            st.markdown(
                '<div style="background:#001c4b04;border:1px dashed #001c4b20;'
                'padding:18px;text-align:center;">'
                '<div style="color:{m};font-size:0.84rem;">Aucune donnée disponible.</div>'
                '</div>'.format(m=MARINE), unsafe_allow_html=True)
        else:
            # Top 20 clients par AUM total pour garder la heatmap lisible
            _ws_clients = df_ws.index.tolist()
            if len(_ws_clients) > 20:
                _aum_totals = df_ws.fillna(0).sum(axis=1).sort_values(ascending=False)
                _ws_clients = _aum_totals.head(20).index.tolist()
            _ws_fonds   = df_ws.columns.tolist()
            _ws_sub     = df_ws.loc[_ws_clients, _ws_fonds]

            # Construction de la matrice z : NaN → 0 (opportunité), valeur → AUM
            _z_raw   = _ws_sub.values.tolist()
            _z_clean = [[0.0 if (v is None or (isinstance(v, float) and np.isnan(v)))
                         else float(v) for v in row] for row in _z_raw]

            # Données pipeline actif pour enrichir le hover : AUM Pipeline par client x fonds
            _ws_pipe_all = db.get_pipeline_with_clients()
            _ws_pipe_actif = _ws_pipe_all[
                _ws_pipe_all["statut"].isin(
                    ["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit"])
            ].copy()
            _ws_pipe_actif["_aum_p"] = _ws_pipe_actif.apply(
                lambda r: float(r["revised_aum"]) if float(r.get("revised_aum", 0) or 0) > 0
                          else float(r.get("target_aum_initial", 0) or 0), axis=1)
            _ws_pipe_dict = {}
            for _, _rp in (_ws_pipe_actif.groupby(["nom_client", "fonds"])["_aum_p"]
                           .sum().reset_index().iterrows()):
                _ws_pipe_dict[(_rp["nom_client"], _rp["fonds"])] = _rp["_aum_p"]

            # Hover propre : "Client : X | Fonds : Y | AUM Finance : ... | AUM Pipeline : ..."
            _hover = [[
                "Client : {client}<br>Fonds : {fonds}<br>"
                "AUM Finance : {aum_f}<br>AUM Pipeline : {aum_p}".format(
                    client=_ws_clients[ri],
                    fonds=_ws_fonds[ci],
                    aum_f=(fmt_m(_z_clean[ri][ci]) if _z_clean[ri][ci] > 0 else "Non investi"),
                    aum_p=(fmt_m(_ws_pipe_dict.get((_ws_clients[ri], _ws_fonds[ci]), 0.0))
                           if _ws_pipe_dict.get((_ws_clients[ri], _ws_fonds[ci]), 0.0) > 0
                           else "—"))
                for ci in range(len(_ws_fonds))]
                for ri in range(len(_ws_clients))]

            # Texte affiché dans les cases (AUM si investi, vide si opportunite)
            _text_ws = [[fmt_m(_z_clean[ri][ci]) if _z_clean[ri][ci] > 0 else ""
                         for ci in range(len(_ws_fonds))]
                        for ri in range(len(_ws_clients))]

            fig_ws = go.Figure(go.Heatmap(
                z=_z_clean,
                x=_ws_fonds,
                y=_ws_clients,
                text=_text_ws,
                texttemplate="%{text}",
                textfont=dict(size=8, color="#ffffff"),
                hovertext=_hover,
                hovertemplate="%{hovertext}<extra></extra>",
                # 0 (vide) = gris tres clair, >0 = degradé bleu institutionnel
                colorscale=[
                    [0.0,   "#f0f4f8"],
                    [0.001, "#c6dff0"],
                    [0.15,  "#6baed6"],
                    [0.40,  "#2171b5"],
                    [0.70,  "#1a5e8a"],
                    [1.0,   "#001c4b"],
                ],
                zmin=0,
                showscale=True,
                colorbar=dict(
                    title=dict(text="AUM (EUR)", font=dict(size=9, color=MARINE)),
                    tickfont=dict(size=8, color=MARINE),
                    thickness=10,
                    len=0.8,
                ),
                # Grille matricielle stricte
                xgap=2,
                ygap=2,
            ))
            fig_ws.update_layout(
                height=max(300, 26 * len(_ws_clients) + 100),
                paper_bgcolor=BLANC, plot_bgcolor=BLANC,
                font_color=MARINE, font_size=9,
                xaxis=dict(side="top", tickangle=-30, automargin=True,
                           tickfont=dict(size=9, color=MARINE), showgrid=False),
                yaxis=dict(autorange="reversed", automargin=True,
                           tickfont=dict(size=9, color=MARINE), showgrid=False),
                margin=dict(l=10, r=60, t=60, b=10))
            st.plotly_chart(fig_ws, use_container_width=True,
                            config={"displayModeBar": False}, key="chart_ws")
            st.caption(
                "Bleu fonce = AUM Finance eleve. Gris clair = opportunite cross-sell. "
                "Survol pour detail AUM Finance et AUM Pipeline.")




# ============================================================================
# ONGLET 4 — SALES TRACKING
# ============================================================================
with tab_sales:
    st.markdown('<div class="section-title">Sales Tracking</div>',
                unsafe_allow_html=True)
    if st.button("Gérer l'équipe commerciale", key="btn_manage_sales_tab"):
        dialog_manage_sales()
    df_sm = db.get_sales_metrics()
    df_na = db.get_next_actions_by_sales(days_ahead=30)

    if df_sm.empty:
        st.info("Aucune donnée de pipeline disponible.")
    else:
        # ── Filtres Sales ─────────────────────────────────────────────────────
        _sf1, _sf2, _sf3 = st.columns(3)
        _all_commerciaux = sorted(df_sm["Commercial"].tolist())
        with _sf1:
            _sel_commerciaux = st.multiselect(
                "Commercial",
                options=_all_commerciaux,
                default=[],
                key="sales_search_ms",
                placeholder="Tous les commerciaux"
            )
        # Filtre marché : extrait depuis sales_team
        _st_data = db.get_sales_team()
        _marches_opts = sorted(_st_data["marche"].unique().tolist()) if not _st_data.empty else []
        with _sf2:
            _sel_marches = st.multiselect(
                "Marché",
                options=_marches_opts,
                default=[],
                key="sales_filt_marche",
                placeholder="Tous les marchés"
            )
        with _sf3:
            _sel_countries_sales = st.multiselect(
                "Pays client",
                options=COUNTRIES_LIST[1:],
                default=[],
                key="sales_filt_country",
                placeholder="Tous les pays"
            )

        # Appliquer les filtres
        df_sm_show = df_sm.copy()
        if _sel_commerciaux:
            df_sm_show = df_sm_show[df_sm_show["Commercial"].isin(_sel_commerciaux)]
        if _sel_marches and not _st_data.empty:
            _owners_in_marche = _st_data[_st_data["marche"].isin(_sel_marches)]["nom"].tolist()
            df_sm_show = df_sm_show[df_sm_show["Commercial"].isin(_owners_in_marche)]

        # ── Enrichissement analytique : Produit Phare & Marché Clé ───────────
        _df_enrich = db.get_pipeline_with_clients()
        # Filtre pays sur le pipeline affiché dans les cartes
        if _sel_countries_sales:
            _df_enrich = _df_enrich[_df_enrich["country"].isin(_sel_countries_sales)]
        _produit_phare_map = {}
        _marche_cle_map    = {}
        if not _df_enrich.empty:
            # AUM total par (sales_owner, fonds) — toutes lignes
            _enr = _df_enrich.copy()
            _enr["_aum_tot"] = _enr.apply(
                lambda r: float(r["funded_aum"]) if r["statut"] == "Funded"
                          else (float(r["revised_aum"]) if float(r.get("revised_aum", 0) or 0) > 0
                                else float(r.get("target_aum_initial", 0) or 0)), axis=1)
            _aum_sf = _enr.groupby(["sales_owner", "fonds"])["_aum_tot"].sum().reset_index()
            _reg_sf = (_enr[_enr["region"].str.strip() != ""]
                       .groupby(["sales_owner", "region"])
                       .size().reset_index(name="n"))
            for _ow in _all_commerciaux:
                _sub_a = _aum_sf[_aum_sf["sales_owner"] == _ow]
                if not _sub_a.empty and _sub_a["_aum_tot"].max() > 0:
                    _produit_phare_map[_ow] = _sub_a.loc[_sub_a["_aum_tot"].idxmax(), "fonds"]
                else:
                    _produit_phare_map[_ow] = "—"
                _sub_r = _reg_sf[_reg_sf["sales_owner"] == _ow]
                if not _sub_r.empty:
                    _marche_cle_map[_ow] = _sub_r.loc[_sub_r["n"].idxmax(), "region"]
                else:
                    _marche_cle_map[_ow] = "—"

        n_own  = len(df_sm_show)
        n_cols = min(max(n_own, 1), 3)
        s_cols = st.columns(n_cols, gap="medium")
        for i, (_, row) in enumerate(df_sm_show.iterrows()):
            retard_val  = int(row.get("Retards", 0))
            retard_html = (
                '<span class="badge-retard">RETARD : {}</span>'.format(retard_val)
                if retard_val > 0
                else '<span style="color:{};font-size:0.74rem;font-weight:600;">A jour</span>'.format(CIEL))
            _owner_nm = row["Commercial"]
            _ph = _produit_phare_map.get(_owner_nm, "—")
            _mk = _marche_cle_map.get(_owner_nm, "—")
            with s_cols[i % n_cols]:
                st.markdown(
                    '<div class="sales-card">'
                    '<div class="sales-card-name">{name}</div>'
                    '<div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;">'
                    '<div><div class="sales-metric">Total Deals</div><div class="sales-metric-val">{nb}</div></div>'
                    '<div><div class="sales-metric">Funded</div><div class="sales-metric-val">{funded}</div></div>'
                    '<div><div class="sales-metric">AUM Finance</div><div class="sales-metric-acc">{aum}</div></div>'
                    '<div><div class="sales-metric">Pipeline Actif</div><div class="sales-metric-val">{pipe}</div></div>'
                    '<div><div class="sales-metric">Actifs / Perdus</div><div class="sales-metric-val">{actifs} / {perdus}</div></div>'
                    '<div><div class="sales-metric">Alertes</div><div style="margin-top:2px;">{retard}</div></div>'
                    '<div style="grid-column:1/3;border-top:1px solid #e8e8e8;padding-top:6px;margin-top:4px;">'
                    '<div class="sales-metric">Produit Phare</div>'
                    '<div style="font-size:0.83rem;font-weight:700;color:{ciel};">{ph}</div></div>'
                    '<div style="grid-column:1/3;">'
                    '<div class="sales-metric">Marche Cle</div>'
                    '<div style="font-size:0.83rem;font-weight:700;color:{marine};">{mk}</div></div>'
                    '</div></div><br>'.format(
                        name=_owner_nm, nb=int(row["Nb_Deals"]), funded=int(row["Funded"]),
                        aum=fmt_m(float(row["AUM_Finance"])), pipe=fmt_m(float(row["Pipeline_Actif"])),
                        actifs=int(row["Actifs"]), perdus=int(row["Perdus"]), retard=retard_html,
                        ph=_ph, mk=_mk, ciel=CIEL, marine=MARINE),
                    unsafe_allow_html=True)

        st.divider()
        st.markdown("#### Funded AUM vs Active Pipeline by Sales Rep")
        if df_sm["AUM_Finance"].sum() > 0:
            fig_sales = go.Figure()
            for lbl, col_key, color in [
                ("AUM Finance",    "AUM_Finance",   MARINE),
                ("Pipeline Actif", "Pipeline_Actif", B_MID)]:
                fig_sales.add_trace(go.Bar(name=lbl, x=df_sm["Commercial"].tolist(),
                                           y=df_sm[col_key].tolist(), marker_color=color,
                                           marker_line_color=BLANC, marker_line_width=0.5))
            fig_sales.update_layout(
                height=300, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
                font_color=MARINE, barmode="group", bargap=0.25,
                legend_bgcolor=BLANC, legend_bordercolor=GRIS, legend_borderwidth=1,
                xaxis_showgrid=False, yaxis_showgrid=True, yaxis_gridcolor=GRIS,
                margin=dict(l=10, r=10, t=20, b=10))
            st.plotly_chart(fig_sales, use_container_width=True, config={"displayModeBar": False}, key="chart_sales")

        st.divider()
        st.markdown("#### Strategic Analysis — Fund Breakdown by Market")
        _si_r1, _si_r2 = st.columns([2, 1])
        with _si_r1:
            si_mode = st.radio(
                "Périmètre",
                ["Pipeline Actif (AUM Révisé)", "Funded (AUM Financé)"],
                horizontal=True, key="si_mode")
        with _si_r2:
            _si_region_filter = st.selectbox(
                "Région", ["Toutes"] + REGIONS,
                key="si_region_filter",
                label_visibility="collapsed"
            )
        df_mfb = db.get_market_fonds_breakdown(
            mode="funded" if "Funded" in si_mode else "pipeline")
        # Filtre région : restreindre aux commerciaux couvrant cette région
        if _si_region_filter != "Toutes" and not df_mfb.empty:
            _pipe_si = db.get_pipeline_with_clients()
            _pipe_si = _pipe_si[_pipe_si["region"] == _si_region_filter]
            _owners_in_reg = _pipe_si["sales_owner"].unique().tolist()
            _st_si = db.get_sales_team()
            if not _st_si.empty:
                _marches_in_reg = _st_si[_st_si["nom"].isin(_owners_in_reg)]["marche"].unique().tolist()
                df_mfb = df_mfb[df_mfb["marche"].isin(_marches_in_reg)] if _marches_in_reg else df_mfb
            else:
                df_mfb = df_mfb[df_mfb["marche"].isin(_owners_in_reg)]

        if df_mfb.empty or df_mfb["aum"].sum() == 0:
            st.markdown(
                '<div style="background:#001c4b04;border:1px dashed #001c4b20;'
                'padding:24px;text-align:center;">'
                '<div style="color:{m};font-weight:600;font-size:0.85rem;">'
                'Aucune donnée disponible pour ce périmètre</div>'
                '<div style="color:#888;font-size:0.75rem;margin-top:3px;">'
                'Enregistrez des deals et associez-les à des commerciaux pour activer cette analyse.</div>'
                '</div>'.format(m=MARINE), unsafe_allow_html=True)
        else:
            marches = sorted(df_mfb["marche"].unique().tolist())
            fonds_presents = sorted(df_mfb["fonds"].unique().tolist())
            fig_si = go.Figure()
            for i, fonds in enumerate(fonds_presents):
                df_f   = df_mfb[df_mfb["fonds"] == fonds]
                y_vals = []
                for m in marches:
                    row_m = df_f[df_f["marche"] == m]
                    y_vals.append(float(row_m["aum"].iloc[0]) if not row_m.empty else 0.0)
                text_vals = [fmt_m(v) for v in y_vals]
                fig_si.add_trace(go.Bar(
                    name=fonds, x=marches, y=y_vals,
                    text=text_vals, textposition="inside",
                    textfont=dict(size=9, color=BLANC),
                    marker_color=PALETTE[i % len(PALETTE)],
                    marker_line_color=BLANC, marker_line_width=0.5,
                    hovertemplate=(
                        "<b>{f}</b><br>Marché : %{{x}}<br>"
                        "AUM : %{{customdata}}<extra></extra>").format(f=fonds),
                    customdata=text_vals))
            fig_si.update_layout(
                barmode="stack", height=360,
                paper_bgcolor=BLANC, plot_bgcolor=BLANC, font_color=MARINE,
                legend_bgcolor=BLANC, legend_bordercolor=GRIS, legend_borderwidth=1,
                legend_font_size=10, legend_title_text="Fonds",
                xaxis_showgrid=False, xaxis_title="Marché Commercial",
                yaxis_showgrid=True, yaxis_gridcolor=GRIS,
                yaxis_title="AUM (EUR)", yaxis_tickformat=".2s",
                margin=dict(l=10, r=10, t=20, b=10))
            st.plotly_chart(fig_si, use_container_width=True, config={"displayModeBar": False}, key="chart_si")
            with st.expander("Tableau détaillé Marché × Fonds", expanded=False):
                pivot_si = df_mfb.pivot_table(
                    index="marche", columns="fonds", values="aum",
                    aggfunc="sum", fill_value=0)
                pivot_si["TOTAL"] = pivot_si.sum(axis=1)
                pivot_display = pivot_si.copy()
                for col in pivot_display.columns:
                    pivot_display[col] = pivot_display[col].apply(fmt_m)
                st.dataframe(pivot_display, use_container_width=True)

        st.divider()
        st.markdown("#### Next Actions — 30 Days")
        if df_na.empty:
            st.info("Aucune action planifiée dans les 30 prochains jours.")
        else:
            owners_na    = ["Tous"] + sorted(df_na["sales_owner"].unique().tolist())
            filter_owner = st.selectbox("Filtrer par commercial", owners_na)
            df_nav = df_na if filter_owner == "Tous" else df_na[df_na["sales_owner"] == filter_owner]
            today  = date.today()
            for _, row in df_nav.iterrows():
                nad = row.get("next_action_date")
                if isinstance(nad, date):
                    delta  = (nad - today).days
                    timing = "Dans {}j".format(delta) if delta >= 0 else "RETARD +{}j".format(abs(delta))
                    dot    = CIEL if delta >= 0 else ORANGE
                    nad_s  = nad.isoformat()
                else:
                    timing = "—"; dot = GRIS; nad_s = "—"
                st.markdown(
                    '<div style="display:flex;align-items:center;gap:11px;padding:7px 11px;'
                    'margin:3px 0;background:#f8fafd;border-left:3px solid {dot};">'
                    '<div style="min-width:100px;font-size:0.73rem;color:{dot};font-weight:700;">{timing}</div>'
                    '<div style="flex:1;"><span style="font-weight:600;color:{marine};font-size:0.82rem;">{client}</span>'
                    '&nbsp;<span style="color:#888;font-size:0.73rem;">{fonds} &middot; {statut}</span></div>'
                    '<div style="font-size:0.78rem;color:{ciel};font-weight:600;min-width:68px;text-align:right;">{aum}</div>'
                    '<div style="font-size:0.73rem;color:#888;min-width:85px;text-align:right;">{owner}</div>'
                    '</div>'.format(
                        dot=dot, timing=timing, marine=MARINE, ciel=CIEL,
                        client=row.get("nom_client",""), fonds=row.get("fonds",""),
                        statut=row.get("statut",""),
                        aum=("N/C" if (float(row.get("revised_aum", 0) or 0) == 0
                                       and float(row.get("target_aum_initial", 0) or 0) == 0)
                             else fmt_m(float(row.get("revised_aum", 0) or 0)
                                       if float(row.get("revised_aum", 0) or 0) > 0
                                       else float(row.get("target_aum_initial", 0) or 0))),
                        owner=row.get("sales_owner","")), unsafe_allow_html=True)


# ============================================================================
# ONGLET 5 — ACTIVITIES
# ============================================================================
with tab_activites:
    st.markdown('<div class="section-title">Activities Journal</div>',
                unsafe_allow_html=True)

    # Toolbar
    act_tb1, act_tb_spacer = st.columns([1, 7])
    with act_tb1:
        if st.button("Enregistrer une activité", key="act_tb_new_act", use_container_width=True):
            dialog_add_activity()

    tf_act = st.selectbox("Timeframe", TIMEFRAMES, key="tf_act",
                          label_visibility="collapsed")
    cutoff_act = _timeframe_cutoff(tf_act)

    df_all_act = db.get_activities()

    df_act_base = df_all_act.copy()
    if cutoff_act is not None and not df_act_base.empty:
        df_act_base = df_act_base[df_act_base["date"].apply(
            lambda d: (isinstance(d, _date_cls) and d >= cutoff_act)
                      or (isinstance(d, str) and d >= cutoff_act.isoformat())
        )]

    fa1, fa2, fa3 = st.columns([2, 2, 3])
    with fa1:
        clients_for_filter = ["Tous"] + (sorted(df_all_act["nom_client"].unique().tolist()) if not df_all_act.empty else [])
        filt_client = st.selectbox("Client", clients_for_filter, key="filt_act_client")
    with fa2:
        filt_type = st.multiselect("Type d'interaction", TYPES_INTERACTION, key="filt_act_type")
    with fa3:
        filt_search = st.text_input("Recherche dans les notes", placeholder="mot-clé…", key="filt_act_search")

    df_act_filtered = df_all_act.copy() if not df_all_act.empty else df_all_act
    if not df_act_filtered.empty:
        if filt_client != "Tous":
            df_act_filtered = df_act_filtered[df_act_filtered["nom_client"] == filt_client]
        if filt_type:
            df_act_filtered = df_act_filtered[df_act_filtered["type_interaction"].isin(filt_type)]
        if filt_search.strip():
            df_act_filtered = df_act_filtered[
                df_act_filtered["notes"].str.contains(filt_search.strip(), case=False, na=False)]

    if not df_act_base.empty:
        sc1, sc2, sc3, sc4 = st.columns(4)
        nb_mois = len(df_act_base[df_act_base["date"].apply(
            lambda d: str(d)[:7] == date.today().strftime("%Y-%m") if d else False)])
        top_type = df_act_base["type_interaction"].value_counts()
        top_type_str = top_type.index[0] if not top_type.empty else "—"
        for col, label, val in [
            (sc1, "Total Activities",  str(len(df_act_base))),
            (sc2, "This Month",        str(nb_mois)),
            (sc3, "Top Activity Type", top_type_str),
            (sc4, "Clients Reached",   str(df_act_base["nom_client"].nunique())),
        ]:
            with col:
                st.markdown(
                    '<div class="kpi-card kpi-card-static" style="padding:10px;">'
                    '<div class="kpi-label">{}</div>'
                    '<div class="kpi-value" style="font-size:1.15rem;">{}</div>'
                    '</div>'.format(label, val), unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.divider()

    if df_act_filtered.empty:
        st.info("Aucune activité trouvée.")
    else:
        nb_total  = len(df_all_act)
        nb_filtre = len(df_act_filtered)
        st.markdown(
            "<div style='font-size:0.78rem;color:#888;margin-bottom:8px;'>"
            "<b>{}</b> activité(s){}".format(
                nb_filtre,
                " sur {} au total".format(nb_total) if nb_filtre < nb_total else ""
            ) + "</div>", unsafe_allow_html=True)

        TYPE_ICONS  = {"Call":"Call","Meeting":"Meeting","Email":"Email",
                       "Roadshow":"Roadshow","Conference":"Conference","Autre":"Autre"}
        TYPE_COLORS = {"Call":CIEL,"Meeting":B_MID,"Email":B_PAL,
                       "Roadshow":"#2c7fb8","Conference":"#004f8c","Autre":"#888"}

        df_sorted = df_act_filtered.copy()
        df_sorted["_date_str"] = df_sorted["date"].apply(lambda d: str(d) if d else "")
        df_sorted = df_sorted.sort_values("_date_str", ascending=False)

        current_date_label = None
        for _, row in df_sorted.iterrows():
            date_str = str(row.get("date",""))
            typ      = str(row.get("type_interaction","Autre"))
            client   = str(row.get("nom_client",""))
            notes    = str(row.get("notes","")) or "—"
            color    = TYPE_COLORS.get(typ, "#888")
            icon     = TYPE_ICONS.get(typ, "")
            act_id   = int(row.get("id", 0))

            if date_str != current_date_label:
                current_date_label = date_str
                try:
                    d_obj  = date.fromisoformat(date_str)
                    delta  = (date.today() - d_obj).days
                    label_d = d_obj.strftime("%A %d %B %Y").capitalize()
                    if   delta == 0: suffix = " — <b style='color:#019ee1;'>Aujourd'hui</b>"
                    elif delta == 1: suffix = " — <span style='color:#888;'>Hier</span>"
                    elif delta <= 7: suffix = " — <span style='color:#888;'>Il y a {} jours</span>".format(delta)
                    else:            suffix = ""
                except Exception:
                    label_d = date_str; suffix = ""
                st.markdown(
                    "<div style='font-size:0.72rem;font-weight:700;color:#7ab8d8;"
                    "text-transform:uppercase;letter-spacing:0.8px;"
                    "padding:10px 0 4px 0;border-bottom:1px solid #e8e8e8;margin-bottom:4px;'>"
                    "{}{}</div>".format(label_d, suffix), unsafe_allow_html=True)

            st.markdown(
                "<div style='display:flex;gap:12px;align-items:flex-start;"
                "padding:9px 14px;margin:3px 0;background:#f9fbfd;"
                "border-left:3px solid {color};'>"
                "<div style='font-size:0.63rem;font-weight:700;color:{color};"
                "background:{color}15;padding:2px 5px;white-space:nowrap;"
                "margin-top:2px;'>{icon}</div>"
                "<div style='flex:1;'>"
                "<div style='font-size:0.82rem;font-weight:700;color:#001c4b;'>"
                "{client}&nbsp;"
                "<span style='font-size:0.69rem;font-weight:600;color:{color};"
                "background:{color}18;padding:1px 8px;'>{typ}</span></div>"
                "<div style='font-size:0.78rem;color:#444;margin-top:4px;line-height:1.55;'>{notes}</div>"
                "</div></div>".format(
                    color=color, icon=icon, client=client, typ=typ, notes=notes),
                unsafe_allow_html=True)

            btn_col1, btn_col2, btn_spacer = st.columns([1, 1, 8])
            with btn_col1:
                if st.button("Modifier", key="act_edit_{}".format(act_id),
                             type="tertiary", help="Modifier cette activité"):
                    dialog_edit_activity(act_id, {
                        "date": date_str, "notes": notes, "type": typ})
            with btn_col2:
                if st.button("Supprimer", key="act_del_{}".format(act_id),
                             type="tertiary", help="Supprimer cette activité"):
                    st.session_state["act_del_confirm"] = act_id

            if st.session_state.get("act_del_confirm") == act_id:
                st.warning("Supprimer cette activité ? Action irréversible.")
                dc1, dc2, _ = st.columns([1, 1, 6])
                with dc1:
                    if st.button("Confirmer", key="act_del_yes_{}".format(act_id)):
                        ok, err = db.delete_activity(act_id)
                        if ok:
                            st.session_state.pop("act_del_confirm", None)
                            st.rerun()
                        else:
                            st.error(err)
                with dc2:
                    if st.button("Annuler", key="act_del_no_{}".format(act_id)):
                        st.session_state.pop("act_del_confirm", None)
                        st.rerun()


# ============================================================================
# ONGLET 6 — SETTINGS & ADMIN
# ============================================================================
with tab_settings:
    st.markdown('<div class="section-title">Settings &amp; Admin</div>',
                unsafe_allow_html=True)

    sa_col1, sa_col2 = st.columns([1, 1], gap="large")

    with sa_col1:
        st.markdown("#### Équipe commerciale")
        st_team_disp = db.get_sales_team()
        if not st_team_disp.empty:
            st.dataframe(st_team_disp[["nom","marche"]], hide_index=True,
                         use_container_width=True,
                         height=min(300, 46 + len(st_team_disp) * 36))
        if st.button("Ajouter un commercial", key="settings_btn_add_sales", use_container_width=True):
            dialog_manage_sales()

        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("#### Actions rapides")
        qa1, qa2 = st.columns(2)
        with qa1:
            if st.button("Nouveau compte client", key="settings_btn_new_client", use_container_width=True):
                dialog_add_client()
        with qa2:
            if st.button("Enregistrer une activité", key="settings_btn_new_act", use_container_width=True):
                dialog_add_activity()

    with sa_col2:
        st.markdown("#### Import CSV / Excel — Smart Staging")
        import_type   = st.radio("Table cible", ["Clients","Pipeline"], horizontal=True,
                                  key="settings_import_type")
        uploaded_file = st.file_uploader("Fichier CSV ou Excel (.xlsx)",
                                         type=["csv","xlsx","xls"],
                                         key="settings_file_uploader")
        if import_type == "Clients":
            st.info("Colonnes : nom_client, type_client, region")
        else:
            st.info(
                "Colonnes supportees : nom_client, type_client, region, country, fonds, statut, "
                "target_aum_initial, revised_aum, funded_aum, closing_probability, "
                "raison_perte, concurrent_choisi, next_action_date, sales_owner"
            )

        if uploaded_file:
            try:
                _file_key = "staging_{}_{}".format(uploaded_file.name, uploaded_file.size)
                if st.session_state.get("_staging_file_key") != _file_key:
                    _df_raw = (pd.read_csv(uploaded_file) if uploaded_file.name.endswith(".csv")
                               else pd.read_excel(uploaded_file))
                    _df_raw.columns = [c.strip() for c in _df_raw.columns]
                    st.session_state["_staging_file_key"] = _file_key
                    st.session_state["_staging_df"] = _df_raw
                    st.session_state.pop("_staging_validated", None)

                df_staging = st.session_state.get("_staging_df")
                if df_staging is not None:
                    if import_type == "Pipeline":
                        _ALLOWED_STATUTS = STATUTS
                        _ALLOWED_FONDS   = FONDS

                        df_check = df_staging.copy()
                        _col_map_lower = {c.lower(): c for c in df_check.columns}
                        _fonds_col  = _col_map_lower.get("fonds",
                                       _col_map_lower.get("fund",
                                       _col_map_lower.get("produit", None)))
                        _statut_col = _col_map_lower.get("statut",
                                       _col_map_lower.get("status",
                                       _col_map_lower.get("stage",
                                       _col_map_lower.get("etape", None))))

                        errors_list = []
                        for idx in range(len(df_check)):
                            row_errors = []
                            if _fonds_col:
                                val = str(df_check.at[idx, _fonds_col]).strip()
                                if val and val != "nan" and val not in _ALLOWED_FONDS:
                                    row_errors.append("Fonds '{}' inconnu".format(val))
                            if _statut_col:
                                val = str(df_check.at[idx, _statut_col]).strip()
                                if val and val != "nan" and val not in _ALLOWED_STATUTS:
                                    row_errors.append("Statut '{}' inconnu".format(val))
                            errors_list.append("; ".join(row_errors))
                        df_check["Errors"] = errors_list

                        _dup_report = db.detect_import_duplicates(df_staging)
                        _dup_exact  = _dup_report.get("exact", [])
                        _dup_fuzzy  = _dup_report.get("fuzzy", [])
                        if _dup_exact:
                            _dup_msg = "**{} doublon(s) exact(s)** : ".format(len(_dup_exact))
                            for _d in _dup_exact[:5]:
                                _dup_msg += "  L.{} & L.{} : {} ".format(
                                    _d["ligne_1"], _d["ligne_2"], _d["valeur"])
                            st.warning(_dup_msg)
                        if _dup_fuzzy:
                            with st.expander("{} doublon(s) potentiel(s) (noms similaires)".format(
                                    len(_dup_fuzzy))):
                                for d in _dup_fuzzy[:10]:
                                    st.markdown("- `{}` ↔ `{}`".format(d["nom_1"], d["nom_2"]))

                        _has_errors = any(e != "" for e in errors_list)
                        _err_count  = sum(1 for e in errors_list if e != "")
                        if _has_errors:
                            st.warning("{} ligne(s) avec erreurs. Corrigez-les ci-dessous avant d'importer.".format(
                                _err_count))

                        _fonds_options = _ALLOWED_FONDS
                        _statut_options = _ALLOWED_STATUTS
                        _col_config_staging = {}
                        if _fonds_col:
                            _col_config_staging[_fonds_col] = st.column_config.SelectboxColumn(
                                "Fonds", options=_fonds_options, required=True)
                        if _statut_col:
                            _col_config_staging[_statut_col] = st.column_config.SelectboxColumn(
                                "Statut", options=_statut_options, required=True)
                        _col_config_staging["Errors"] = st.column_config.TextColumn(
                            "Errors", disabled=True)

                        edited_staging = st.data_editor(
                            df_check,
                            use_container_width=True,
                            hide_index=False,
                            num_rows="fixed",
                            column_config=_col_config_staging,
                            key="staging_editor")

                        _edited_no_err = edited_staging.copy()
                        _new_errors = []
                        for idx in range(len(_edited_no_err)):
                            row_errors = []
                            if _fonds_col:
                                val = str(_edited_no_err.at[idx, _fonds_col]).strip()
                                if val and val != "nan" and val not in _ALLOWED_FONDS:
                                    row_errors.append("Fonds '{}' inconnu".format(val))
                            if _statut_col:
                                val = str(_edited_no_err.at[idx, _statut_col]).strip()
                                if val and val != "nan" and val not in _ALLOWED_STATUTS:
                                    row_errors.append("Statut '{}' inconnu".format(val))
                            _new_errors.append("; ".join(row_errors))
                        _still_has_errors = any(e != "" for e in _new_errors)

                        st.caption("{} ligne(s) totale(s)".format(len(edited_staging)))

                        if _still_has_errors:
                            st.button("Confirmer & Importer", disabled=True,
                                      key="settings_btn_confirm_disabled",
                                      use_container_width=True,
                                      help="Corrigez toutes les erreurs avant d'importer")
                        else:
                            if st.button("Confirmer & Importer",
                                         key="settings_btn_confirm_import",
                                         use_container_width=True, type="primary"):
                                _df_to_import = edited_staging.drop(columns=["Errors"], errors="ignore")
                                st.session_state["_staging_df"] = _df_to_import
                                ins, upd = db.upsert_pipeline_from_df(_df_to_import)
                                st.success("Import : {} créé(s), {} mis à jour.".format(ins, upd))
                                st.session_state.pop("_staging_file_key", None)
                                st.session_state.pop("_staging_df", None)

                    else:
                        st.dataframe(df_staging.head(10), use_container_width=True, height=200)
                        st.caption("{} ligne(s)".format(len(df_staging)))
                        if st.button("Lancer l'import Clients",
                                     key="settings_btn_import_clients",
                                     use_container_width=True, type="primary"):
                            ins, upd = db.upsert_clients_from_df(df_staging)
                            st.success("Import : {} créé(s), {} mis à jour.".format(ins, upd))
                            st.session_state.pop("_staging_file_key", None)
                            st.session_state.pop("_staging_df", None)
            except Exception as e:
                st.error("Erreur : {}".format(e))


# ============================================================================
# ONGLET 7 — PERFORMANCE & NAV
# ============================================================================
with tab_perf:
    st.markdown('<div class="section-title">Performance et NAV</div>', unsafe_allow_html=True)

    col_up, col_info = st.columns([1.2, 1], gap="large")

    with col_info:
        st.markdown(
            '<div style="background:#001c4b07;border:1px solid #001c4b14;padding:15px 17px;">'
            '<div style="font-weight:700;color:#001c4b;font-size:0.90rem;margin-bottom:9px;">Format attendu</div>'
            '<div style="font-size:0.79rem;color:#444;line-height:1.75;">'
            '<b>Colonnes obligatoires</b><br>'
            '&bull; <code>Date</code> — YYYY-MM-DD ou DD/MM/YYYY<br>'
            '&bull; <code>Fonds</code> — Nom du fonds<br>'
            '&bull; <code>NAV</code> — Valeur liquidative numerique<br><br>'
            '<b>Calculs produits</b><br>'
            '&bull; Base 100 normalisee &bull; Perf 1M &bull; Perf YTD<br><br>'
            '<b>Export PDF</b><br>'
            'Les donnees NAV sont integrees dans le PDF si cochees.'
            '</div></div>', unsafe_allow_html=True)
        st.markdown("#### Démonstration")
        if st.button("Generer un fichier NAV de demonstration", key="perf_btn_demo_nav", use_container_width=True):
            demo_dates = pd.date_range("{}-01-01".format(date.today().year - 1),
                                       date.today(), freq="B")
            rng  = np.random.default_rng(42)
            navs = {f: 100.0 for f in FONDS}
            rows = []
            for d in demo_dates:
                for fonds, nav in navs.items():
                    nav *= (1 + rng.normal(0.0003, 0.006))
                    navs[fonds] = nav
                    rows.append({"Date": d.date().isoformat(), "Fonds": fonds, "NAV": round(nav, 4)})
            buf_demo = io.BytesIO()
            pd.DataFrame(rows).to_excel(buf_demo, index=False)
            buf_demo.seek(0)
            st.download_button("Telecharger nav_demo.xlsx", data=buf_demo,
                               file_name="nav_demo.xlsx", mime="application/vnd.ms-excel",
                               use_container_width=True)

    with col_up:
        nav_file = st.file_uploader("Charger l'historique NAV (Excel ou CSV)",
                                    type=["xlsx","xls","csv"])

    if nav_file is not None:
        try:
            df_nav = (pd.read_csv(nav_file) if nav_file.name.endswith(".csv")
                      else pd.read_excel(nav_file))
            df_nav.columns = [c.strip() for c in df_nav.columns]
            missing = [c for c in ["Date","Fonds","NAV"] if c not in df_nav.columns]
            if missing:
                st.error("Colonnes manquantes : {}".format(missing)); st.stop()
            df_nav["Date"] = pd.to_datetime(df_nav["Date"], format="mixed", dayfirst=True, errors="coerce")
            df_nav = df_nav.dropna(subset=["Date"])
            df_nav["NAV"]   = pd.to_numeric(df_nav["NAV"], errors="coerce")
            df_nav = df_nav.dropna(subset=["NAV"])
            df_nav["Fonds"] = df_nav["Fonds"].astype(str).str.strip()
            df_nav = df_nav.sort_values("Date").reset_index(drop=True)
            if df_nav.empty:
                st.error("Aucune donnée valide."); st.stop()

            fonds_list = sorted(df_nav["Fonds"].unique().tolist())
            d_min = df_nav["Date"].min(); d_max = df_nav["Date"].max()
            st.markdown(
                '<div style="background:#019ee114;border-left:3px solid #019ee1;'
                'padding:8px 14px;margin:9px 0;">'
                '{:,} points &mdash; {} fonds &mdash; {} au {}</div>'.format(
                    len(df_nav), len(fonds_list),
                    d_min.strftime('%d/%m/%Y'), d_max.strftime('%d/%m/%Y')),
                unsafe_allow_html=True)

            st.markdown("---")
            ff1, ff2, ff3 = st.columns([2, 1, 1])
            with ff1:
                fonds_sel_nav = st.multiselect("Fonds à afficher", fonds_list,
                                               default=fonds_list[:min(5, len(fonds_list))])
            with ff2: d_debut = st.date_input("Depuis", value=d_min.date())
            with ff3: d_fin   = st.date_input("Jusqu'au", value=d_max.date())

            if not fonds_sel_nav:
                st.warning("Sélectionnez au moins un fonds."); st.stop()

            mask = (df_nav["Fonds"].isin(fonds_sel_nav) &
                    (df_nav["Date"].dt.date >= d_debut) &
                    (df_nav["Date"].dt.date <= d_fin))
            df_fn = df_nav[mask].copy()
            if df_fn.empty:
                st.warning("Aucune donnée pour la période."); st.stop()

            pivot = (df_fn.pivot_table(index="Date", columns="Fonds", values="NAV", aggfunc="last")
                     .sort_index().ffill())
            pivot = pivot[[f for f in fonds_sel_nav if f in pivot.columns]]

            base100 = pivot.copy() * np.nan
            for fonds in pivot.columns:
                s = pivot[fonds].dropna()
                if s.empty: continue
                first_val = float(s.iloc[0])
                if first_val != 0 and not np.isnan(first_val):
                    base100[fonds] = pivot[fonds] / first_val * 100

            st.session_state["nav_base100"] = base100
            st.session_state["perf_fonds"]  = fonds_sel_nav

            st.markdown("#### Evolution NAV — Base 100")
            NAV_COLORS = [MARINE, CIEL, B_MID, B_PAL, B_DEP, "#2c7fb8","#004f8c","#6baed6"]
            fig_nav = go.Figure()
            for i, fonds in enumerate(pivot.columns):
                series = base100[fonds].dropna()
                if series.empty: continue
                fig_nav.add_trace(go.Scatter(
                    x=series.index.tolist(), y=series.values.tolist(),
                    mode="lines", name=fonds,
                    line=dict(color=NAV_COLORS[i % len(NAV_COLORS)], width=2),
                    hovertemplate="<b>{}</b><br>%{{y:.2f}}<extra></extra>".format(fonds)))
            fig_nav.add_hline(y=100, line_dash="dot", line_color=GRIS, line_width=1)
            fig_nav.update_layout(
                title_text="Performance NAV — Base 100", title_font_size=13,
                title_font_color=MARINE, height=380,
                paper_bgcolor=BLANC, plot_bgcolor=BLANC, font_color=MARINE,
                legend_bgcolor=BLANC, legend_bordercolor=GRIS, legend_borderwidth=1,
                legend_font_size=10, xaxis_showgrid=False,
                yaxis_showgrid=True, yaxis_gridcolor=GRIS,
                margin=dict(l=10, r=10, t=40, b=10))
            st.plotly_chart(fig_nav, use_container_width=True, config={"displayModeBar": True}, key="chart_nav")

            today_ts  = pd.Timestamp(date.today())
            one_m_ago = today_ts - pd.DateOffset(months=1)
            jan_1     = pd.Timestamp("{}-01-01".format(date.today().year))

            perf_rows = []
            for fonds in pivot.columns:
                series = pivot[fonds].dropna()
                if series.empty: continue
                nav_last = float(series.iloc[-1]); nav_first = float(series.iloc[0])
                s_1m  = series[series.index >= one_m_ago]
                p1m   = ((nav_last / float(s_1m.iloc[0]) - 1) * 100
                         if len(s_1m) > 0 and float(s_1m.iloc[0]) != 0 else float("nan"))
                s_ytd = series[series.index >= jan_1]
                pytd  = ((nav_last / float(s_ytd.iloc[0]) - 1) * 100
                         if len(s_ytd) > 0 and float(s_ytd.iloc[0]) != 0 else float("nan"))
                pp    = ((nav_last / nav_first - 1) * 100 if nav_first != 0 else float("nan"))
                b100s = base100[fonds].dropna()
                nb100 = float(b100s.iloc[-1]) if not b100s.empty else float("nan")
                perf_rows.append({
                    "Fonds": fonds, "NAV Derniere": round(nav_last, 4),
                    "Base 100 Actuel": round(nb100, 2) if not np.isnan(nb100) else None,
                    "Perf 1M (%)":     round(p1m,  2) if not np.isnan(p1m)  else None,
                    "Perf YTD (%)":    round(pytd, 2) if not np.isnan(pytd) else None,
                    "Perf Periode (%)": round(pp, 2) if not np.isnan(pp)    else None})

            if perf_rows:
                df_pt = pd.DataFrame(perf_rows)
                st.session_state["perf_data"] = df_pt
                st.markdown("#### Tableau des Performances")

                def _fp(val):
                    if val is None or (isinstance(val, float) and np.isnan(val)):
                        return '<span style="color:#999;">n.d.</span>'
                    c = CIEL if val >= 0 else "#8b2020"
                    s = "+" if val > 0 else ""
                    return '<span style="color:{};font-weight:700;">{}{:.2f}%</span>'.format(c,s,val)

                tbl_h = (
                    '<table style="width:100%;border-collapse:collapse;font-size:0.82rem;">'
                    '<thead><tr style="background:{marine};color:white;">'
                    '<th style="padding:8px 12px;text-align:left;">Fonds</th>'
                    '<th style="padding:8px 12px;text-align:right;">NAV</th>'
                    '<th style="padding:8px 12px;text-align:right;">Base 100</th>'
                    '<th style="padding:8px 12px;text-align:right;">Perf 1M</th>'
                    '<th style="padding:8px 12px;text-align:right;">Perf YTD</th>'
                    '<th style="padding:8px 12px;text-align:right;">Perf Période</th>'
                    '</tr></thead><tbody>').format(marine=MARINE)
                for i, r in enumerate(perf_rows):
                    bg = "#f0f5fa" if i % 2 == 0 else BLANC
                    tbl_h += (
                        '<tr style="background:{bg};border-bottom:1px solid {gris};">'
                        '<td style="padding:7px 12px;font-weight:600;color:{marine};">{fonds}</td>'
                        '<td style="padding:7px 12px;text-align:right;color:{marine};">{nav}</td>'
                        '<td style="padding:7px 12px;text-align:right;color:{marine};">{b100}</td>'
                        '<td style="padding:7px 12px;text-align:right;">{p1m}</td>'
                        '<td style="padding:7px 12px;text-align:right;">{pytd}</td>'
                        '<td style="padding:7px 12px;text-align:right;">{pp}</td>'
                        '</tr>').format(
                        bg=bg, gris=GRIS, marine=MARINE, fonds=r["Fonds"],
                        nav="{:.4f}".format(r["NAV Derniere"]),
                        b100="{:.2f}".format(r["Base 100 Actuel"]) if r["Base 100 Actuel"] else "n.d.",
                        p1m=_fp(r["Perf 1M (%)"]), pytd=_fp(r["Perf YTD (%)"]),
                        pp=_fp(r["Perf Periode (%)"]))
                tbl_h += "</tbody></table>"
                st.markdown(tbl_h, unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)
                st.download_button("Exporter le tableau en CSV",
                                   data=df_pt.to_csv(index=False).encode("utf-8"),
                                   file_name="performances_{}.csv".format(date.today().isoformat()),
                                   mime="text/csv")

                ytd_data = [r for r in perf_rows if r["Perf YTD (%)"] is not None]
                if ytd_data:
                    ytd_data.sort(key=lambda r: r["Perf YTD (%)"], reverse=True)
                    st.markdown("#### Comparaison Performances YTD")
                    fig_ytd = go.Figure(go.Bar(
                        x=[r["Fonds"] for r in ytd_data],
                        y=[r["Perf YTD (%)"] for r in ytd_data],
                        marker_color=[CIEL if v >= 0 else GRIS
                                      for v in [r["Perf YTD (%)"] for r in ytd_data]],
                        marker_line_color=BLANC, marker_line_width=0.4,
                        text=["{:+.2f}%".format(r["Perf YTD (%)"]) for r in ytd_data],
                        textposition="outside", textfont_size=9, textfont_color=MARINE))
                    fig_ytd.add_hline(y=0, line_color=MARINE, line_width=0.8)
                    fig_ytd.update_layout(height=300, paper_bgcolor=BLANC, plot_bgcolor=BLANC,
                                          font_color=MARINE, xaxis_tickangle=-15,
                                          xaxis_showgrid=False, yaxis_showgrid=True,
                                          yaxis_gridcolor=GRIS, margin=dict(l=10,r=10,t=36,b=10))
                    st.plotly_chart(fig_ytd, use_container_width=True,
                                    config={"displayModeBar": False}, key="chart_ytd")

        except Exception as e:
            st.error("Erreur traitement NAV : {}".format(e))
            import traceback
            with st.expander("Détails de l'erreur"):
                st.code(traceback.format_exc())
    else:
        st.markdown(
            '<div style="background:#001c4b04;border:2px dashed #001c4b1e;'
            'padding:44px;text-align:center;margin-top:12px;">'
            '<div style="font-size:0.96rem;font-weight:700;color:#001c4b;margin-bottom:5px;">'
            'Module Performance et NAV</div>'
            '<div style="color:#777;font-size:0.81rem;max-width:380px;margin:0 auto;line-height:1.65;">'
            'Chargez un fichier Excel ou CSV avec les colonnes '
            '<code>Date</code>, <code>Fonds</code>, <code>NAV</code>.<br><br>'
            'Utilisez le bouton <b>Generer un fichier NAV de demonstration</b>.'
            '</div></div>', unsafe_allow_html=True)
