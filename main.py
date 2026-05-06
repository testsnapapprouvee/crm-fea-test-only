# =============================================================================
# main.py — Meridian CRM · FastAPI backend (full feature parity)
# Serves frontend/revolut_crm_dashboard_v2.html and exposes /api/* JSON
# endpoints for full CRUD + import/export, backed by database.py (SQLite)
# =============================================================================

import io
import os
from datetime import date, datetime, timedelta
from typing import List, Optional

import numpy as np
import pandas as pd
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.responses import HTMLResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

import database as db

ROOT_DIR     = os.path.dirname(os.path.abspath(__file__))
FRONTEND_DIR = os.path.join(ROOT_DIR, "frontend")
INDEX_PATH   = os.path.join(FRONTEND_DIR, "revolut_crm_dashboard_v2.html")

app = FastAPI(
    title="Meridian CRM API",
    description="Full-feature CRM backend (clients, deals, contacts, activities, sales).",
    version="2.0.0",
)

db.init_db()

if os.path.isdir(FRONTEND_DIR):
    app.mount("/static", StaticFiles(directory=FRONTEND_DIR), name="static")


# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def _smart_aum(row) -> float:
    statut = str(row.get("statut", ""))
    if statut == "Funded":
        return float(row.get("funded_aum", 0) or 0)
    rev = float(row.get("revised_aum", 0) or 0)
    return rev if rev > 0 else float(row.get("target_aum_initial", 0) or 0)


def _safe_dict(d):
    out = {}
    for k, v in d.items():
        if isinstance(v, (np.integer,)):
            out[k] = int(v)
        elif isinstance(v, (np.floating,)):
            out[k] = float(v)
        elif isinstance(v, (np.bool_,)):
            out[k] = bool(v)
        elif isinstance(v, (pd.Timestamp,)):
            out[k] = v.isoformat()
        elif isinstance(v, (date, datetime)):
            out[k] = v.isoformat()
        else:
            out[k] = v
    return out


def _parse_date(s: Optional[str]) -> Optional[str]:
    """Normalise a date string to YYYY-MM-DD or return None."""
    if not s:
        return None
    s = str(s).strip()
    if not s:
        return None
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%Y-%m-%dT%H:%M:%S", "%Y-%m-%dT%H:%M"):
        try:
            return datetime.strptime(s, fmt).date().isoformat()
        except ValueError:
            continue
    try:
        return pd.to_datetime(s).date().isoformat()
    except Exception:
        return None


# ---------------------------------------------------------------------------
# PYDANTIC MODELS
# ---------------------------------------------------------------------------

class ClientCreate(BaseModel):
    nom_client:        str
    type_client:       str = ""
    region:            str = ""
    country:           str = ""
    parent_id:         Optional[int] = None
    tier:              str = "Tier 2"
    kyc_status:        str = "En cours"
    product_interests: str = ""


class ClientUpdate(ClientCreate):
    pass


class ContactCreate(BaseModel):
    client_id:  int
    prenom:     str = ""
    nom:        str
    role:       str = ""
    email:      str = ""
    telephone:  str = ""
    linkedin:   str = ""
    is_primary: bool = False


class ContactUpdate(BaseModel):
    prenom:     str = ""
    nom:        str
    role:       str = ""
    email:      str = ""
    telephone:  str = ""
    linkedin:   str = ""
    is_primary: bool = False


class DealCreate(BaseModel):
    client_id:           int
    fonds:               str
    statut:              str = "Prospect"
    target_aum_initial:  float = 0
    revised_aum:         float = 0
    funded_aum:          float = 0
    raison_perte:        str = ""
    concurrent_choisi:   str = ""
    next_action_date:    Optional[str] = None
    sales_owner:         str = "Non assigne"
    closing_probability: float = 50


class DealUpdate(DealCreate):
    pass


class ActivityCreate(BaseModel):
    client_id:        int
    date:             str
    notes:            str = ""
    type_interaction: str = "Call"


class ActivityUpdate(BaseModel):
    date:             str
    notes:            str = ""
    type_interaction: str = "Call"


class SalesMemberCreate(BaseModel):
    nom:    str
    marche: str = "Global"


class SalesMemberUpdate(BaseModel):
    nom:    str
    marche: str = "Global"


class BulkDeleteIds(BaseModel):
    ids: List[int]


# ---------------------------------------------------------------------------
# ROOT — serve the HTML
# ---------------------------------------------------------------------------

@app.get("/", response_class=HTMLResponse)
def root():
    if not os.path.exists(INDEX_PATH):
        raise HTTPException(404, "Frontend not found.")
    with open(INDEX_PATH, "r", encoding="utf-8") as f:
        return HTMLResponse(content=f.read())


@app.get("/health")
def health():
    return {"status": "ok", "service": "meridian-crm-api", "version": "2.0.0"}


# ---------------------------------------------------------------------------
# /api/kpis
# ---------------------------------------------------------------------------

@app.get("/api/kpis")
def api_kpis():
    k = db.get_kpis()
    return {
        "total_funded":      float(k.get("total_funded", 0)),
        "pipeline_actif":    float(k.get("pipeline_actif", 0)),
        "weighted_pipeline": float(k.get("weighted_pipeline", 0)),
        "taux_conversion":   float(k.get("taux_conversion", 0)),
        "nb_funded":         int(k.get("nb_funded", 0)),
        "nb_lost":           int(k.get("nb_lost", 0)),
        "nb_paused":         int(k.get("nb_paused", 0)),
        "nb_deals_actifs":   int(k.get("nb_deals_actifs", 0)),
        "statut_repartition": k.get("statut_repartition", {}),
    }


# ---------------------------------------------------------------------------
# DEALS — CRUD
# ---------------------------------------------------------------------------

@app.get("/api/deals")
def api_deals():
    df = db.get_pipeline_with_clients()
    if df.empty:
        return []
    items = []
    for _, r in df.iterrows():
        statut = str(r.get("statut", ""))
        aum = float(r.get("funded_aum", 0) or 0) if statut == "Funded" else _smart_aum(r)
        items.append(_safe_dict({
            "id":          int(r["id"]),
            "client_id":   int(r["client_id"]),
            "client":      str(r["nom_client"]),
            "type_client": str(r.get("type_client", "")),
            "region":      str(r.get("region", "")),
            "country":     str(r.get("country", "")),
            "fonds":       str(r.get("fonds", "")),
            "statut":      statut,
            "aum":         aum,
            "funded_aum":  float(r.get("funded_aum", 0) or 0),
            "revised_aum": float(r.get("revised_aum", 0) or 0),
            "target_aum":  float(r.get("target_aum_initial", 0) or 0),
            "sales_owner": str(r.get("sales_owner", "")),
            "next_action_date": r.get("next_action_date").isoformat()
                if isinstance(r.get("next_action_date"), date) else None,
            "closing_probability": float(r.get("closing_probability", 50) or 50),
        }))
    items.sort(key=lambda x: x["aum"], reverse=True)
    return items


@app.get("/api/deals/{deal_id}")
def api_deal_detail(deal_id: int):
    row = db.get_pipeline_row_by_id(deal_id)
    if not row:
        raise HTTPException(404, "Deal not found")
    nad = row.get("next_action_date")
    row["next_action_date"] = nad.isoformat() if isinstance(nad, date) else None
    return _safe_dict(row)


@app.post("/api/deals", status_code=201)
def api_deal_create(deal: DealCreate):
    nad = _parse_date(deal.next_action_date)
    new_id = db.add_pipeline_entry(
        client_id=deal.client_id,
        fonds=deal.fonds,
        statut=deal.statut,
        target_aum=deal.target_aum_initial,
        revised_aum=deal.revised_aum,
        funded_aum=deal.funded_aum,
        raison_perte=deal.raison_perte,
        concurrent_choisi=deal.concurrent_choisi,
        next_action_date=nad,
        sales_owner=deal.sales_owner,
        closing_probability=deal.closing_probability,
    )
    return {"id": new_id, "ok": True}


@app.put("/api/deals/{deal_id}")
def api_deal_update(deal_id: int, deal: DealUpdate):
    nad = _parse_date(deal.next_action_date)
    ok, err = db.update_pipeline_row({
        "id":                  deal_id,
        "fonds":               deal.fonds,
        "statut":              deal.statut,
        "target_aum_initial":  deal.target_aum_initial,
        "revised_aum":         deal.revised_aum,
        "funded_aum":          deal.funded_aum,
        "raison_perte":        deal.raison_perte,
        "concurrent_choisi":   deal.concurrent_choisi,
        "next_action_date":    nad,
        "sales_owner":         deal.sales_owner,
        "closing_probability": deal.closing_probability,
    })
    if not ok:
        raise HTTPException(400, err or "Update failed")
    return {"id": deal_id, "ok": True}


@app.delete("/api/deals/{deal_id}")
def api_deal_delete(deal_id: int):
    ok, err = db.delete_pipeline_row(deal_id)
    if not ok:
        raise HTTPException(400, err or "Delete failed")
    return {"id": deal_id, "ok": True}


@app.post("/api/deals/bulk-delete")
def api_deals_bulk_delete(payload: BulkDeleteIds):
    n = db.delete_pipeline_rows(payload.ids)
    return {"deleted": n}


@app.get("/api/deals/{deal_id}/audit")
def api_deal_audit(deal_id: int):
    df = db.get_audit_log(deal_id)
    if df.empty:
        return []
    return [_safe_dict(r.to_dict()) for _, r in df.iterrows()]


# ---------------------------------------------------------------------------
# CLIENTS — CRUD
# ---------------------------------------------------------------------------

@app.get("/api/clients")
def api_clients():
    df_clients = db.get_client_hierarchy()
    df_pipe    = db.get_pipeline_with_clients()
    aum_by_client = {}
    if not df_pipe.empty:
        for cid, grp in df_pipe.groupby("client_id"):
            funded = float(grp[grp["statut"] == "Funded"]["funded_aum"].sum())
            active = grp[grp["statut"].isin(
                ["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit"])]
            active_aum = float(active.apply(_smart_aum, axis=1).sum()) if not active.empty else 0.0
            aum_by_client[int(cid)] = funded + active_aum
    out = []
    for _, r in df_clients.iterrows():
        cid = int(r["id"])
        out.append(_safe_dict({
            "id":          cid,
            "nom_client":  str(r["nom_client"]),
            "type_client": str(r.get("type_client", "")),
            "region":      str(r.get("region", "")),
            "country":     str(r.get("country", "")),
            "tier":        str(r.get("tier", "")),
            "kyc_status":  str(r.get("kyc_status", "")),
            "parent_id":   int(r.get("parent_id") or 0),
            "parent_nom":  str(r.get("parent_nom", "") or ""),
            "aum":         aum_by_client.get(cid, 0.0),
        }))
    out.sort(key=lambda c: c["aum"], reverse=True)
    return out


@app.get("/api/clients/{client_id}")
def api_client_detail(client_id: int):
    df_clients = db.get_client_hierarchy()
    row = df_clients[df_clients["id"] == client_id]
    if row.empty:
        raise HTTPException(404, "Client not found")
    c = row.iloc[0].to_dict()

    contacts_df   = db.get_contacts(client_id)
    activities_df = db.get_activities(client_id=client_id)
    df_pipe       = db.get_pipeline_with_clients()
    deals = df_pipe[df_pipe["client_id"] == client_id]

    funded_aum = float(deals[deals["statut"] == "Funded"]["funded_aum"].sum()) \
        if not deals.empty else 0.0
    active = deals[deals["statut"].isin(
        ["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit"])] \
        if not deals.empty else deals
    active_aum = float(active.apply(_smart_aum, axis=1).sum()) if not active.empty else 0.0

    fonds_invested = sorted(deals[deals["statut"] == "Funded"]["fonds"].unique().tolist()) \
        if not deals.empty else []

    dominant_statut = "—"
    if not deals.empty:
        dominant_statut = str(deals["statut"].value_counts().index[0])

    primary_contact = None
    if not contacts_df.empty:
        primary = contacts_df[contacts_df["is_primary"] == 1]
        pc_row = primary.iloc[0] if not primary.empty else contacts_df.iloc[0]
        primary_contact = "{} {}".format(pc_row.get("prenom", ""), pc_row.get("nom", "")).strip()

    network = db.get_client_network(client_id)
    fund_links_clean = [_safe_dict({
        "fonds":       str(fl.get("fonds", "")),
        "client_id":   int(fl.get("client_id", 0)),
        "client_nom":  str(fl.get("nom_client", "")),
        "statut":      str(fl.get("statut", "")),
        "aum":         float(fl.get("aum", 0) or 0),
    }) for fl in network.get("fund_links", [])]
    subs_clean = [_safe_dict({
        "id":          int(s.get("id", 0)),
        "nom_client":  str(s.get("nom_client", "")),
        "type_client": str(s.get("type_client", "")),
    }) for s in network.get("subsidiaries", [])]

    contacts_clean = [_safe_dict({
        "id":         int(ct["id"]),
        "prenom":     str(ct.get("prenom", "")),
        "nom":        str(ct.get("nom", "")),
        "role":       str(ct.get("role", "")),
        "email":      str(ct.get("email", "")),
        "telephone":  str(ct.get("telephone", "")),
        "linkedin":   str(ct.get("linkedin", "")),
        "is_primary": bool(int(ct.get("is_primary", 0) or 0)),
    }) for _, ct in contacts_df.iterrows()]

    activities_clean = [_safe_dict({
        "id":               int(a["id"]),
        "date":             str(a.get("date", "")),
        "notes":            str(a.get("notes", "") or ""),
        "type_interaction": str(a.get("type_interaction", "") or ""),
    }) for _, a in activities_df.iterrows()]

    return _safe_dict({
        "id":              int(c["id"]),
        "nom_client":      str(c["nom_client"]),
        "type_client":     str(c.get("type_client", "")),
        "region":          str(c.get("region", "")),
        "country":         str(c.get("country", "")),
        "tier":            str(c.get("tier", "Tier 2")),
        "kyc_status":      str(c.get("kyc_status", "En cours")),
        "parent_id":       int(c.get("parent_id") or 0),
        "parent_nom":      str(c.get("parent_nom", "") or ""),
        "primary_contact": primary_contact,
        "aum_funded":      funded_aum,
        "aum_active":      active_aum,
        "aum_total":       funded_aum + active_aum,
        "nb_active_deals": int(len(active)),
        "nb_activities":   int(len(activities_clean)),
        "dominant_statut": dominant_statut,
        "fonds_invested":  fonds_invested,
        "created_year":    None,
        "network": {
            "root":         network.get("root"),
            "subsidiaries": subs_clean,
            "fund_links":   fund_links_clean,
        },
        "contacts":   contacts_clean,
        "activities": activities_clean,
    })


@app.post("/api/clients", status_code=201)
def api_client_create(payload: ClientCreate):
    try:
        new_id = db.add_client(
            nom_client=payload.nom_client,
            type_client=payload.type_client,
            region=payload.region,
            country=payload.country,
            parent_id=payload.parent_id,
            tier=payload.tier,
            kyc_status=payload.kyc_status,
            product_interests=payload.product_interests,
        )
        return {"id": new_id, "ok": True}
    except Exception as e:
        raise HTTPException(400, str(e))


@app.put("/api/clients/{client_id}")
def api_client_update(client_id: int, payload: ClientUpdate):
    ok, err = db.update_client(
        client_id=client_id,
        nom_client=payload.nom_client,
        type_client=payload.type_client,
        region=payload.region,
        country=payload.country,
        parent_id=payload.parent_id,
        tier=payload.tier,
        kyc_status=payload.kyc_status,
        product_interests=payload.product_interests,
    )
    if not ok:
        raise HTTPException(400, err or "Update failed")
    return {"id": client_id, "ok": True}


@app.delete("/api/clients/{client_id}")
def api_client_delete(client_id: int):
    ok, err = db.delete_client(client_id)
    if not ok:
        raise HTTPException(400, err or "Delete failed")
    return {"id": client_id, "ok": True}


# ---------------------------------------------------------------------------
# CONTACTS — CRUD
# ---------------------------------------------------------------------------

@app.get("/api/contacts")
def api_contacts(client_id: Optional[int] = None):
    if client_id is None:
        df_clients = db.get_all_clients()
        all_contacts = []
        for cid in df_clients["id"]:
            df_c = db.get_contacts(int(cid))
            for _, ct in df_c.iterrows():
                all_contacts.append(_safe_dict(ct.to_dict()))
        return all_contacts
    df = db.get_contacts(client_id)
    return [_safe_dict(ct.to_dict()) for _, ct in df.iterrows()]


@app.post("/api/contacts", status_code=201)
def api_contact_create(payload: ContactCreate):
    try:
        new_id = db.add_contact(
            client_id=payload.client_id,
            prenom=payload.prenom, nom=payload.nom,
            role=payload.role, email=payload.email,
            telephone=payload.telephone, linkedin=payload.linkedin,
            is_primary=payload.is_primary,
        )
        return {"id": new_id, "ok": True}
    except Exception as e:
        raise HTTPException(400, str(e))


@app.put("/api/contacts/{contact_id}")
def api_contact_update(contact_id: int, payload: ContactUpdate):
    ok, err = db.update_contact(
        contact_id=contact_id,
        prenom=payload.prenom, nom=payload.nom,
        role=payload.role, email=payload.email,
        telephone=payload.telephone, linkedin=payload.linkedin,
        is_primary=payload.is_primary,
    )
    if not ok:
        raise HTTPException(400, err or "Update failed")
    return {"id": contact_id, "ok": True}


@app.delete("/api/contacts/{contact_id}")
def api_contact_delete(contact_id: int):
    ok, err = db.delete_contact(contact_id)
    if not ok:
        raise HTTPException(400, err or "Delete failed")
    return {"id": contact_id, "ok": True}


# ---------------------------------------------------------------------------
# ACTIVITIES — CRUD
# ---------------------------------------------------------------------------

@app.get("/api/activities")
def api_activities(client_id: Optional[int] = None):
    df = db.get_activities(client_id=client_id)
    if df.empty:
        return []
    return [_safe_dict(a.to_dict()) for _, a in df.iterrows()]


@app.post("/api/activities", status_code=201)
def api_activity_create(payload: ActivityCreate):
    nad = _parse_date(payload.date) or date.today().isoformat()
    db.add_activity(payload.client_id, nad, payload.notes, payload.type_interaction)
    return {"ok": True}


@app.put("/api/activities/{activity_id}")
def api_activity_update(activity_id: int, payload: ActivityUpdate):
    nad = _parse_date(payload.date) or date.today().isoformat()
    ok, err = db.update_activity(activity_id, nad, payload.notes, payload.type_interaction)
    if not ok:
        raise HTTPException(400, err or "Update failed")
    return {"id": activity_id, "ok": True}


@app.delete("/api/activities/{activity_id}")
def api_activity_delete(activity_id: int):
    ok, err = db.delete_activity(activity_id)
    if not ok:
        raise HTTPException(400, err or "Delete failed")
    return {"id": activity_id, "ok": True}


# ---------------------------------------------------------------------------
# SALES TEAM — CRUD
# ---------------------------------------------------------------------------

@app.get("/api/sales")
def api_sales():
    team_df    = db.get_sales_team()
    metrics_df = db.get_sales_metrics()

    by_owner = {}
    if not metrics_df.empty:
        for _, m in metrics_df.iterrows():
            by_owner[str(m["Commercial"])] = m

    out = []
    for _, t in team_df.iterrows():
        nom = str(t["nom"])
        m = by_owner.get(nom, None)
        funded_aum   = float(m["AUM_Finance"])    if m is not None else 0.0
        pipeline_aum = float(m["Pipeline_Actif"]) if m is not None else 0.0
        nb_actifs    = int(m["Actifs"])           if m is not None else 0
        nb_funded    = int(m["Funded"])           if m is not None else 0
        nb_perdus    = int(m["Perdus"])           if m is not None else 0
        nb_total     = nb_funded + nb_perdus
        conversion   = round(nb_funded / nb_total * 100, 0) if nb_total > 0 else 0
        out.append(_safe_dict({
            "id":            int(t["id"]),
            "nom":           nom,
            "marche":        str(t.get("marche", "")),
            "funded_aum":    funded_aum,
            "pipeline_aum":  pipeline_aum,
            "nb_actifs":     nb_actifs,
            "nb_funded":     nb_funded,
            "conversion":    int(conversion),
        }))
    out.sort(key=lambda s: s["funded_aum"], reverse=True)
    return out


@app.post("/api/sales", status_code=201)
def api_sales_create(payload: SalesMemberCreate):
    ok = db.add_sales_member(payload.nom, payload.marche)
    if not ok:
        raise HTTPException(400, "Sales member already exists or invalid name")
    return {"ok": True}


@app.put("/api/sales/{sales_id}")
def api_sales_update(sales_id: int, payload: SalesMemberUpdate):
    ok, err = db.update_sales_member(sales_id, payload.nom, payload.marche)
    if not ok:
        raise HTTPException(400, err or "Update failed")
    return {"id": sales_id, "ok": True}


@app.delete("/api/sales/{sales_id}")
def api_sales_delete(sales_id: int):
    ok, err = db.delete_sales_member(sales_id)
    if not ok:
        raise HTTPException(400, err or "Delete failed")
    return {"id": sales_id, "ok": True}


# ---------------------------------------------------------------------------
# REGIONS / WHITESPACE / NEXT-ACTIONS / FUNDS / MONTHLY / MINI
# ---------------------------------------------------------------------------

@app.get("/api/regions")
def api_regions():
    region_aum = db.get_aum_by_region()
    out = [{"region": k, "aum": float(v)} for k, v in region_aum.items()]
    out.sort(key=lambda r: r["aum"], reverse=True)
    return out


# Country → ISO-3 code mapping for Plotly choropleth
_COUNTRY_ISO = {
    "United Arab Emirates": "ARE", "Saudi Arabia": "SAU", "Qatar": "QAT",
    "Kuwait": "KWT", "Bahrain": "BHR", "Oman": "OMN",
    "United Kingdom": "GBR", "France": "FRA", "Germany": "DEU",
    "Switzerland": "CHE", "Luxembourg": "LUX", "Netherlands": "NLD",
    "Italy": "ITA", "Spain": "ESP", "Belgium": "BEL", "Austria": "AUT",
    "Sweden": "SWE", "Norway": "NOR", "Denmark": "DNK", "Finland": "FIN",
    "Singapore": "SGP", "Japan": "JPN", "Hong Kong": "HKG",
    "China": "CHN", "South Korea": "KOR", "Australia": "AUS", "India": "IND",
    "United States": "USA", "Canada": "CAN", "Brazil": "BRA",
    "Mexico": "MEX", "South Africa": "ZAF", "Egypt": "EGY",
}


@app.get("/api/aum-by-country")
def api_aum_by_country():
    df = db.get_aum_by_country()
    if df is None or df.empty:
        return {"countries": []}
    out = []
    for _, r in df.iterrows():
        cty = str(r["country"])
        out.append({
            "country":      cty,
            "iso":          _COUNTRY_ISO.get(cty),
            "funded_aum":   float(r["funded_aum"]),
            "pipeline_aum": float(r["pipeline_aum"]),
            "total_aum":    float(r["total_aum"]),
            "nb_clients":   int(r["nb_clients"]),
        })
    return {"countries": out}


@app.get("/api/whitespace")
def api_whitespace():
    matrix = db.get_whitespace_matrix()
    if matrix is None or matrix.empty:
        return {"clients": [], "fonds": [], "values": []}
    totals = matrix.fillna(0).sum(axis=1).sort_values(ascending=False)
    top_clients = totals.head(5).index.tolist()
    sub = matrix.loc[top_clients]
    fonds = sub.columns.tolist()
    values = []
    for _, row in sub.iterrows():
        values.append([
            0.0 if (v is None or (isinstance(v, float) and np.isnan(v)))
            else float(v) for v in row
        ])
    return {"clients": top_clients, "fonds": fonds, "values": values}


@app.get("/api/whitespace-full")
def api_whitespace_full():
    """Full whitespace matrix: all clients × all fonds (Funded + active pipeline)."""
    df_pipe = db.get_pipeline_with_clients()
    if df_pipe.empty:
        return {"clients": [], "fonds": [], "funded": [], "pipeline": []}
    fonds = sorted(df_pipe["fonds"].dropna().unique().tolist())
    fonds = [f for f in fonds if f]
    if not fonds:
        return {"clients": [], "fonds": [], "funded": [], "pipeline": []}
    clients = sorted(df_pipe["nom_client"].dropna().unique().tolist())
    df_pipe["_aum_p"] = df_pipe.apply(_smart_aum, axis=1)
    funded = []
    pipeline = []
    for cli in clients:
        sub = df_pipe[df_pipe["nom_client"] == cli]
        f_row = []
        p_row = []
        for fund in fonds:
            sub_f = sub[sub["fonds"] == fund]
            if sub_f.empty:
                f_row.append(0.0)
                p_row.append(0.0)
            else:
                f_val = float(sub_f[sub_f["statut"] == "Funded"]["funded_aum"].sum())
                p_val = float(sub_f[sub_f["statut"].isin(
                    ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])]["_aum_p"].sum())
                f_row.append(f_val)
                p_row.append(p_val)
        funded.append(f_row)
        pipeline.append(p_row)
    # Sort clients by total AUM descending
    totals = [sum(funded[i]) + sum(pipeline[i]) for i in range(len(clients))]
    order = sorted(range(len(clients)), key=lambda i: totals[i], reverse=True)
    return {
        "clients":  [clients[i] for i in order],
        "fonds":    fonds,
        "funded":   [funded[i] for i in order],
        "pipeline": [pipeline[i] for i in order],
    }


@app.get("/api/pipeline-matrix")
def api_pipeline_matrix():
    """Cross-tab analytics for the Pipeline tab.
    Returns 4 grouped views: by sales × region, by sales × fonds,
    by region × statut, and a 'who-where' summary per sales rep."""
    df = db.get_pipeline_with_clients()
    if df.empty:
        return {
            "by_sales":     [],
            "by_region":    [],
            "by_fonds":     [],
            "who_does_what": [],
            "where":        [],
        }
    df["_aum_p"] = df.apply(_smart_aum, axis=1)
    df["funded_aum"] = pd.to_numeric(df["funded_aum"], errors="coerce").fillna(0.0)

    # By sales — total Funded + Pipeline + nb deals + regions covered
    by_sales = []
    for owner, grp in df.groupby("sales_owner"):
        funded   = float(grp[grp["statut"] == "Funded"]["funded_aum"].sum())
        active   = grp[grp["statut"].isin(
            ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])]
        pipe_aum = float(active["_aum_p"].sum())
        regions  = sorted([r for r in grp["region"].dropna().unique() if r])
        funds    = sorted([f for f in grp["fonds"].dropna().unique() if f])
        countries = sorted([c for c in grp["country"].dropna().unique() if c])
        by_sales.append({
            "sales_owner": str(owner) if owner else "Non assigné",
            "funded_aum":  funded,
            "pipeline_aum": pipe_aum,
            "total_aum":   funded + pipe_aum,
            "nb_deals":    int(len(grp)),
            "nb_funded":   int((grp["statut"] == "Funded").sum()),
            "nb_active":   int(len(active)),
            "regions":     regions,
            "fonds":       funds,
            "countries":   countries,
            "nb_clients":  int(grp["client_id"].nunique()),
        })
    by_sales.sort(key=lambda x: x["total_aum"], reverse=True)

    # By region — total Funded + Pipeline + sales reps active
    by_region = []
    for region, grp in df.groupby("region"):
        if not region:
            continue
        funded   = float(grp[grp["statut"] == "Funded"]["funded_aum"].sum())
        active   = grp[grp["statut"].isin(
            ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])]
        pipe_aum = float(active["_aum_p"].sum())
        owners   = sorted([o for o in grp["sales_owner"].dropna().unique()
                           if o and o != "Non assigne"])
        countries = sorted([c for c in grp["country"].dropna().unique() if c])
        by_region.append({
            "region":      str(region),
            "funded_aum":  funded,
            "pipeline_aum": pipe_aum,
            "total_aum":   funded + pipe_aum,
            "nb_deals":    int(len(grp)),
            "nb_clients":  int(grp["client_id"].nunique()),
            "sales_owners": owners,
            "countries":   countries,
        })
    by_region.sort(key=lambda x: x["total_aum"], reverse=True)

    # By fonds — total Funded + Pipeline + clients/sales
    by_fonds = []
    for fonds, grp in df.groupby("fonds"):
        if not fonds:
            continue
        funded   = float(grp[grp["statut"] == "Funded"]["funded_aum"].sum())
        active   = grp[grp["statut"].isin(
            ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])]
        pipe_aum = float(active["_aum_p"].sum())
        owners   = sorted([o for o in grp["sales_owner"].dropna().unique()
                           if o and o != "Non assigne"])
        regions  = sorted([r for r in grp["region"].dropna().unique() if r])
        by_fonds.append({
            "fonds":         str(fonds),
            "funded_aum":    funded,
            "pipeline_aum":  pipe_aum,
            "total_aum":     funded + pipe_aum,
            "nb_deals":      int(len(grp)),
            "nb_clients":    int(grp["client_id"].nunique()),
            "sales_owners":  owners,
            "regions":       regions,
        })
    by_fonds.sort(key=lambda x: x["total_aum"], reverse=True)

    # who-does-what: sales × fonds matrix
    who_does_what = []
    sales_list = sorted([s for s in df["sales_owner"].dropna().unique() if s])
    fonds_list = sorted([f for f in df["fonds"].dropna().unique() if f])
    for owner in sales_list:
        row = {"sales_owner": owner, "by_fonds": {}}
        for fonds in fonds_list:
            sub = df[(df["sales_owner"] == owner) & (df["fonds"] == fonds)]
            if sub.empty:
                continue
            funded = float(sub[sub["statut"] == "Funded"]["funded_aum"].sum())
            pipe   = float(sub[sub["statut"].isin(
                ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])]["_aum_p"].sum())
            row["by_fonds"][fonds] = {
                "funded_aum":  funded,
                "pipeline_aum": pipe,
                "total_aum":   funded + pipe,
                "nb_deals":    int(len(sub)),
            }
        who_does_what.append(row)

    # where: country × sales matrix (top 10 countries)
    where = []
    for country, grp in df.groupby("country"):
        if not country:
            continue
        funded   = float(grp[grp["statut"] == "Funded"]["funded_aum"].sum())
        pipe_aum = float(grp[grp["statut"].isin(
            ["Prospect","Initial Pitch","Due Diligence","Soft Commit"])]["_aum_p"].sum())
        owners = sorted([o for o in grp["sales_owner"].dropna().unique()
                         if o and o != "Non assigne"])
        clients = sorted([c for c in grp["nom_client"].dropna().unique() if c])
        where.append({
            "country":      str(country),
            "funded_aum":   funded,
            "pipeline_aum": pipe_aum,
            "total_aum":    funded + pipe_aum,
            "sales_owners": owners,
            "clients":      clients,
            "nb_deals":     int(len(grp)),
        })
    where.sort(key=lambda x: x["total_aum"], reverse=True)

    return {
        "by_sales":      by_sales,
        "by_region":     by_region,
        "by_fonds":      by_fonds,
        "who_does_what": who_does_what,
        "where":         where[:12],
    }


@app.get("/api/top-deals")
def api_top_deals(period: str = "all", group_by: str = "deal", limit: int = 10):
    """Top deals BI endpoint.
    - period: 'all' | '30d' | '90d' | '180d' | '365d' | 'ytd'
    - group_by: 'deal' | 'fonds' | 'client' | 'sales' | 'region' | 'country'
    - limit: number of top items returned (default 10)
    """
    df = db.get_pipeline_with_clients()
    if df.empty:
        return {"period": period, "group_by": group_by, "items": []}
    df = df[df["statut"] == "Funded"].copy()
    df["funded_aum"] = pd.to_numeric(df["funded_aum"], errors="coerce").fillna(0.0)
    df = df[df["funded_aum"] > 0]

    today = date.today()
    cutoff = None
    if period == "30d":  cutoff = today - timedelta(days=30)
    elif period == "90d": cutoff = today - timedelta(days=90)
    elif period == "180d": cutoff = today - timedelta(days=180)
    elif period == "365d": cutoff = today - timedelta(days=365)
    elif period == "ytd": cutoff = date(today.year, 1, 1)

    if cutoff is not None:
        # Filter on next_action_date (proxy for recency since updated_at is not in df)
        df["nad_dt"] = pd.to_datetime(df["next_action_date"], errors="coerce")
        cutoff_ts = pd.Timestamp(cutoff)
        df = df[(df["nad_dt"].isna()) | (df["nad_dt"] >= cutoff_ts)]

    if df.empty:
        return {"period": period, "group_by": group_by, "items": []}

    items = []
    if group_by == "deal":
        df = df.sort_values("funded_aum", ascending=False).head(limit)
        for _, r in df.iterrows():
            items.append(_safe_dict({
                "label":      str(r["nom_client"]),
                "sublabel":   "{} · {}".format(r.get("fonds",""), r.get("sales_owner","—")),
                "country":    str(r.get("country","")),
                "region":     str(r.get("region","")),
                "aum":        float(r["funded_aum"]),
                "deal_id":    int(r["id"]),
            }))
    elif group_by in ("fonds", "client", "sales", "region", "country"):
        col_map = {"fonds":"fonds","client":"nom_client","sales":"sales_owner",
                   "region":"region","country":"country"}
        col = col_map[group_by]
        agg = df.groupby(col)["funded_aum"].agg(["sum", "count"]).reset_index()
        agg.columns = [col, "sum", "count"]
        agg = agg.sort_values("sum", ascending=False).head(limit)
        for _, r in agg.iterrows():
            label = str(r[col]) if r[col] else "—"
            items.append({
                "label":    label,
                "sublabel": "{} deal(s)".format(int(r["count"])),
                "aum":      float(r["sum"]),
                "count":    int(r["count"]),
            })
    return {"period": period, "group_by": group_by, "items": items, "total": float(sum(i["aum"] for i in items))}


@app.get("/api/next-actions")
def api_next_actions():
    today    = date.today()
    horizon  = today + timedelta(days=30)
    overdue  = today - timedelta(days=60)
    df = db.get_pipeline_with_clients()
    if df.empty:
        return []
    df = df[df["statut"].isin(["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit"])].copy()
    out = []
    for _, r in df.iterrows():
        nad = r.get("next_action_date")
        if not isinstance(nad, date):
            continue
        if not (overdue <= nad <= horizon):
            continue
        out.append(_safe_dict({
            "id":               int(r["id"]),
            "nom_client":       str(r["nom_client"]),
            "fonds":            str(r.get("fonds", "")),
            "statut":           str(r.get("statut", "")),
            "next_action_date": nad.isoformat(),
            "sales_owner":      str(r.get("sales_owner", "")),
        }))
    out.sort(key=lambda x: x["next_action_date"])
    return out


@app.get("/api/funds-aggregate")
def api_funds_aggregate():
    df = db.get_pipeline_with_clients()
    if df.empty:
        return {}
    df_f = df[df["statut"] == "Funded"]
    return {str(k): float(v) for k, v in df_f.groupby("fonds")["funded_aum"].sum().items()}


@app.get("/api/monthly-aum")
def api_monthly_aum():
    df = db.get_historical_aum(days_back=365)
    if df is None or df.empty:
        return []
    df = df.copy()
    df["month"] = pd.to_datetime(df["date"]).dt.to_period("M")
    agg = df.groupby("month")["funded_aum"].max().reset_index()
    fr_months = ["jan", "fév", "mar", "avr", "mai", "juin",
                 "juil", "août", "sept", "oct", "nov", "déc"]
    return [{
        "label": fr_months[r["month"].month - 1],
        "month": str(r["month"]),
        "aum":   float(r["funded_aum"]),
    } for _, r in agg.tail(12).iterrows()]


@app.get("/api/mini-stats")
def api_mini_stats():
    conn = db.get_connection()
    c = conn.cursor()
    c.execute("SELECT COUNT(*) FROM clients");  nb_clients  = int(c.fetchone()[0])
    c.execute("SELECT COUNT(*) FROM contacts"); nb_contacts = int(c.fetchone()[0])
    cutoff_30d = (date.today() - timedelta(days=30)).isoformat()
    c.execute("SELECT COUNT(*) FROM activites WHERE date >= ?", (cutoff_30d,))
    nb_activities_30d = int(c.fetchone()[0])
    today_str = date.today().isoformat()
    c.execute(
        "SELECT COUNT(*) FROM pipeline"
        " WHERE next_action_date < ?"
        " AND statut NOT IN ('Lost','Funded','Redeemed')",
        (today_str,))
    nb_overdue = int(c.fetchone()[0])
    c.execute("SELECT COUNT(DISTINCT fonds) FROM pipeline WHERE fonds != ''")
    nb_fonds = int(c.fetchone()[0])
    conn.close()
    return {
        "nb_clients":        nb_clients,
        "nb_contacts":       nb_contacts,
        "nb_activities_30d": nb_activities_30d,
        "nb_overdue":        nb_overdue,
        "nb_fonds":          nb_fonds,
    }


# ---------------------------------------------------------------------------
# REFERENTIELS (for forms)
# ---------------------------------------------------------------------------

@app.get("/api/referentials")
def api_referentials():
    return {
        "fonds":             list(db.FONDS_REFERENTIEL),
        "regions":           list(db.REGIONS_REFERENTIEL),
        "tiers":             list(db.TIERS_REFERENTIEL),
        "kyc_statuts":       list(db.KYC_STATUTS),
        "product_interests": list(db.PRODUCT_INTERESTS),
        "roles_contact":     list(db.ROLES_CONTACT),
        "types_client":      ["IFA", "Wholesale", "Instit", "Family Office",
                              "Insurance", "Asset Manager", "Sovereign", "Pension"],
        "statuts":           ["Prospect", "Initial Pitch", "Due Diligence",
                              "Soft Commit", "Funded", "Paused", "Lost", "Redeemed"],
        "types_interaction": ["Call", "Meeting", "Email", "Roadshow", "Conference", "Autre"],
        "raisons_perte":     ["Pricing", "Track Record", "Macro", "Competitor", "Autre"],
        "countries":         [
            "United Arab Emirates", "Saudi Arabia", "Qatar", "Kuwait", "Bahrain", "Oman",
            "United Kingdom", "France", "Germany", "Switzerland", "Luxembourg", "Netherlands",
            "Italy", "Spain", "Belgium", "Austria", "Sweden", "Norway", "Denmark", "Finland",
            "Singapore", "Japan", "Hong Kong", "China", "South Korea", "Australia", "India",
            "United States", "Canada", "Brazil", "Mexico", "South Africa", "Egypt",
        ],
    }


# ---------------------------------------------------------------------------
# IMPORT / EXPORT / ADMIN
# ---------------------------------------------------------------------------

def _read_uploaded_file(file: UploadFile) -> pd.DataFrame:
    raw = file.file.read()
    if not raw:
        raise HTTPException(400, "Empty file")
    name = (file.filename or "").lower()
    bio = io.BytesIO(raw)
    try:
        if name.endswith(".csv"):
            return pd.read_csv(bio)
        return pd.read_excel(bio)
    except Exception as e:
        raise HTTPException(400, "Cannot parse file: {}".format(e))


@app.post("/api/import/clients")
def api_import_clients(file: UploadFile = File(...)):
    df = _read_uploaded_file(file)
    inserted, updated = db.upsert_clients_from_df(df)
    return {"inserted": inserted, "updated": updated}


@app.post("/api/import/pipeline")
def api_import_pipeline(file: UploadFile = File(...)):
    df = _read_uploaded_file(file)
    inserted, updated = db.upsert_pipeline_from_df(df)
    return {"inserted": inserted, "updated": updated}


@app.get("/api/export/excel")
def api_export_excel(anon: int = 0):
    if anon:
        # Generate anonymised version on the fly
        df_pipe = db.get_pipeline_with_clients().copy()
        df_pipe["nom_client"] = df_pipe.apply(
            lambda r: _anonymize_label(r.get("nom_client"), int(r.get("client_id", 0))), axis=1)
        # Rebuild a minimal Excel with the anonymised pipeline
        bio = io.BytesIO()
        with pd.ExcelWriter(bio, engine="openpyxl") as w:
            df_pipe.to_excel(w, sheet_name="Pipeline_Anonyme", index=False)
        data = bio.getvalue()
    else:
        data = db.get_excel_backup()
    suffix = "_anon" if anon else ""
    return Response(
        content=data,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition":
                 "attachment; filename=meridian_backup{}_{}.xlsx".format(suffix, date.today().isoformat())},
    )


def _anonymize_label(name: str, client_id: int = 0) -> str:
    if client_id:
        return "Client #{:03d}".format(int(client_id))
    h = abs(hash(str(name or ""))) % 1000
    return "Client #{:03d}".format(h)


def _fmt_aum(v) -> str:
    try:
        v = float(v or 0)
    except (TypeError, ValueError):
        return "—"
    if v >= 1e9:  return "{:.2f} Md€".format(v / 1e9)
    if v >= 1e6:  return "{:.1f} M€".format(v / 1e6)
    if v >= 1e3:  return "{:.0f} k€".format(v / 1e3)
    return "{:.0f} €".format(v)


def _make_chart_png(plot_fn, w_inch=6.5, h_inch=3.2, dpi=140):
    """Run plot_fn(ax) under a Meridian-styled matplotlib context, return PNG bytes."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({
        "font.family":         "DejaVu Sans",
        "font.size":           9,
        "axes.spines.top":     False,
        "axes.spines.right":   False,
        "axes.edgecolor":      "#E5E7EB",
        "axes.linewidth":      0.6,
        "axes.labelcolor":     "#001c4b",
        "xtick.color":         "#374151",
        "ytick.color":         "#374151",
        "figure.facecolor":    "#FFFFFF",
        "axes.facecolor":      "#FFFFFF",
        "grid.color":          "#F3F4F6",
        "grid.linewidth":      0.5,
    })
    fig, ax = plt.subplots(figsize=(w_inch, h_inch), dpi=dpi)
    plot_fn(ax)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor="#FFFFFF", pad_inches=0.05)
    plt.close(fig)
    buf.seek(0)
    return buf


def _build_pdf_report(anon: bool = False) -> bytes:
    """Multi-page executive PDF with charts."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.lib.enums import TA_LEFT, TA_RIGHT, TA_CENTER
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
        Image as RLImage, KeepTogether,
    )
    from reportlab.platypus.flowables import Flowable

    MARINE = HexColor("#001c4b")
    CIEL   = HexColor("#019ee1")
    GREY   = HexColor("#6B7280")
    LIGHT  = HexColor("#F3F4F6")
    BLANC  = HexColor("#FFFFFF")
    DARKB  = HexColor("#0F3A7A")
    SOFT   = HexColor("#7ab8d8")

    PAGE_W, PAGE_H = A4
    MARGIN_H, MARGIN_V = 2 * cm, 1.8 * cm
    USABLE_W = PAGE_W - 2 * MARGIN_H

    class ColorRect(Flowable):
        def __init__(self, w, h, color):
            super().__init__()
            self.width, self.height, self.color = w, h, color
        def draw(self):
            self.canv.setFillColor(self.color)
            self.canv.rect(0, 0, self.width, self.height, fill=1, stroke=0)

    class CoverHero(Flowable):
        """Full-bleed navy hero section for the cover."""
        def __init__(self, width, height, title, subtitle, date_str, anon_flag):
            super().__init__()
            self.width  = width; self.height = height
            self.title    = title
            self.subtitle = subtitle
            self.date_str = date_str
            self.anon     = anon_flag
        def draw(self):
            c = self.canv
            c.saveState()
            c.setFillColor(MARINE)
            c.rect(0, 0, self.width, self.height, fill=1, stroke=0)
            c.setFillColor(CIEL)
            c.rect(0, 0, self.width, 4, fill=1, stroke=0)
            # Brand mark
            c.setFillColor(SOFT)
            c.setFont("Helvetica", 8)
            c.drawString(22, self.height - 22,
                         "MERIDIAN CAPITAL  ·  ASSET MANAGEMENT")
            if self.anon:
                c.setFillColor(HexColor("#f07d00"))
                c.setFont("Helvetica-Bold", 8)
                c.drawRightString(self.width - 22, self.height - 22,
                                  "MODE ANONYME · COMEX")
            # Title
            c.setFillColor(BLANC)
            c.setFont("Helvetica-Bold", 26)
            c.drawString(22, self.height - 70, self.title)
            c.setFillColor(SOFT)
            c.setFont("Helvetica", 13)
            c.drawString(22, self.height - 92, self.subtitle)
            # Date
            c.setFillColor(SOFT)
            c.setFont("Helvetica", 10)
            c.drawString(22, 18, self.date_str.upper())
            c.restoreState()

    s_h2   = ParagraphStyle("h2",   fontName="Helvetica-Bold", fontSize=13, textColor=MARINE, spaceBefore=10, spaceAfter=6)
    s_sub  = ParagraphStyle("sub",  fontName="Helvetica-Oblique", fontSize=8.5, textColor=GREY, spaceAfter=8)
    s_body = ParagraphStyle("body", fontName="Helvetica", fontSize=9, textColor=HexColor("#374151"), leading=12)
    s_kpi_lbl = ParagraphStyle("kpi_lbl", fontName="Helvetica",      fontSize=8,  textColor=SOFT,  alignment=TA_CENTER)
    s_kpi_val = ParagraphStyle("kpi_val", fontName="Helvetica-Bold", fontSize=18, textColor=BLANC, alignment=TA_CENTER)
    s_th      = ParagraphStyle("th",   fontName="Helvetica-Bold", fontSize=8,   textColor=BLANC,  alignment=TA_LEFT)
    s_th_r    = ParagraphStyle("th_r", fontName="Helvetica-Bold", fontSize=8,   textColor=BLANC,  alignment=TA_RIGHT)
    s_td      = ParagraphStyle("td",   fontName="Helvetica",      fontSize=8.5, textColor=MARINE, alignment=TA_LEFT)
    s_td_r    = ParagraphStyle("td_r", fontName="Helvetica",      fontSize=8.5, textColor=MARINE, alignment=TA_RIGHT)
    s_disc    = ParagraphStyle("disc", fontName="Helvetica-Oblique", fontSize=7, textColor=GREY, alignment=TA_CENTER)
    s_caption = ParagraphStyle("cap",  fontName="Helvetica-Oblique", fontSize=7.5, textColor=GREY, alignment=TA_CENTER, spaceAfter=6)

    def header_footer(canvas_obj, doc):
        canvas_obj.saveState()
        if doc.page > 1:
            # Top accent line
            canvas_obj.setStrokeColor(CIEL)
            canvas_obj.setLineWidth(1.6)
            canvas_obj.line(0, PAGE_H - 0.4*cm, PAGE_W, PAGE_H - 0.4*cm)
            # Footer bar
            canvas_obj.setFillColor(MARINE)
            canvas_obj.rect(0, 0, PAGE_W, 0.85*cm, fill=1, stroke=0)
            canvas_obj.setFillColor(SOFT)
            canvas_obj.setFont("Helvetica", 7)
            label = "Meridian Capital · Executive Report"
            if anon:
                label += "  ·  MODE ANONYME"
            canvas_obj.drawString(2*cm, 0.3*cm, label)
            canvas_obj.drawRightString(PAGE_W - 2*cm, 0.3*cm,
                                        "Page {} · CONFIDENTIEL".format(doc.page))
        canvas_obj.restoreState()

    deals_all = api_deals()
    sales     = api_sales()
    regions   = api_regions()
    k         = api_kpis()
    funds_agg = api_funds_aggregate()

    if anon:
        for d in deals_all:
            d["client"] = _anonymize_label(d.get("client"), d.get("client_id", 0))
        # Sales reps stay (internal staff), but their owner names in deals stay too.

    funded = [d for d in deals_all if d["statut"] == "Funded"]
    funded_sorted = sorted(funded, key=lambda d: d["funded_aum"], reverse=True)

    # ── BUILD ─────────────────────────────────────────────────
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=MARGIN_H, rightMargin=MARGIN_H,
                            topMargin=1.4*cm, bottomMargin=1.4*cm,
                            title="Meridian Capital — Executive Report",
                            author="Meridian Capital Partners")
    story = []

    # ── PAGE 1 — COVER + KPIs ─────────────────────────────────────────────
    story.append(CoverHero(USABLE_W, 200,
                            "Executive Report",
                            "Pipeline · Performance · Distribution",
                            date.today().strftime("%d %B %Y"), anon))
    story.append(Spacer(1, 14))

    kpi_items = [
        ("AUM Financé Total",  _fmt_aum(k["total_funded"])),
        ("Pipeline Actif",     _fmt_aum(k["pipeline_actif"])),
        ("Pipeline Pondéré",   _fmt_aum(k["weighted_pipeline"])),
        ("Taux de Conversion", "{:.1f} %".format(k["taux_conversion"])),
    ]
    kpi_tbl = Table(
        [[Paragraph(lbl, s_kpi_lbl) for lbl, _ in kpi_items],
         [Paragraph(val, s_kpi_val) for _, val in kpi_items]],
        colWidths=[USABLE_W / 4] * 4,
        rowHeights=[0.7*cm, 1.2*cm],
    )
    kpi_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,-1), MARINE),
        ("LINEAFTER",     (0,0),(2,-1),  0.6, CIEL),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 6),
        ("BOTTOMPADDING", (0,0),(-1,-1), 6),
    ]))
    story.append(kpi_tbl)
    story.append(Spacer(1, 14))

    # Summary stats line
    nb_clients = api_mini_stats()["nb_clients"]
    nb_deals_active = k.get("nb_deals_actifs", 0)
    nb_funded       = k.get("nb_funded", 0)
    summary_html = (
        "<font color='#001c4b' size='10'><b>Synthèse exécutive.</b></font> "
        + "Le portefeuille compte <b>{nb_clients} client(s)</b>, "
        + "<b>{nb_funded} deal(s) Funded</b> et <b>{nb_active} deal(s) actif(s)</b> "
        + "couvrant <b>{nb_regions} région(s)</b> et <b>{nb_funds} fonds</b>. "
        + "Le pipeline pondéré (probabilité de closing) atteint <b>{wp}</b>, "
        + "soit un potentiel additionnel de <b>{wp_pct:.0f} %</b> de l'AUM financé."
    ).format(
        nb_clients=nb_clients, nb_funded=nb_funded,
        nb_active=nb_deals_active, nb_regions=len(regions), nb_funds=len(funds_agg),
        wp=_fmt_aum(k["weighted_pipeline"]),
        wp_pct=(k["weighted_pipeline"] / k["total_funded"] * 100)
                if k["total_funded"] > 0 else 0,
    )
    story.append(Paragraph(summary_html, s_body))
    story.append(Spacer(1, 14))

    # ── CHART 1 — AUM par fonds (horizontal bar) ──────────────────────────
    story.append(Paragraph("Distribution AUM par fonds", s_h2))
    story.append(Paragraph("Concentration des encours financés sur les principaux véhicules.", s_sub))
    if funds_agg:
        items = sorted(funds_agg.items(), key=lambda x: x[1], reverse=True)[:8]
        names = [n for n, _ in items]
        vals  = [v for _, v in items]
        def plot_funds(ax):
            colors = ["#001c4b","#1a5e8a","#3060A8","#6BAED6","#94A3B8","#CBD5E1","#9AB0CF","#E5E7EB"]
            bars = ax.barh(range(len(names)), [v/1e6 for v in vals],
                           color=colors[:len(names)], edgecolor="white", height=0.7)
            ax.set_yticks(range(len(names)))
            ax.set_yticklabels(names, fontsize=9, color="#001c4b")
            ax.invert_yaxis()
            ax.set_xlabel("AUM (M€)", fontsize=8, color="#6B7280")
            ax.tick_params(axis="x", labelsize=7.5, color="#94A3B8")
            ax.grid(axis="x", alpha=0.3)
            for i, b in enumerate(bars):
                w = b.get_width()
                ax.text(w + max(vals)/1e6 * 0.012, b.get_y() + b.get_height()/2,
                        _fmt_aum(vals[i]), va="center", fontsize=8,
                        color="#001c4b", fontweight="bold")
            ax.set_xlim(0, max(vals)/1e6 * 1.22)
        png = _make_chart_png(plot_funds, w_inch=6.7, h_inch=2.8)
        story.append(RLImage(png, width=USABLE_W, height=USABLE_W * 0.42))
        story.append(Paragraph("Figure 1 — AUM Financé par fonds (M€)", s_caption))
    story.append(PageBreak())

    # ── PAGE 2 — TOP 10 DEALS ─────────────────────────────────────────────
    story.append(Paragraph("Top Deals — AUM Financé", s_h2))
    story.append(Paragraph("Classement des dix plus importants deals Funded.", s_sub))
    if funded_sorted:
        ratios = [0.05, 0.30, 0.20, 0.16, 0.16, 0.13]
        col_w  = [USABLE_W * r for r in ratios]
        rows = [[Paragraph(h, s_th_r if i in (0, 5) else s_th)
                 for i, h in enumerate(["#", "Client", "Fonds", "Pays / Région", "Commercial", "AUM Financé"])]]
        for i, d in enumerate(funded_sorted[:10], 1):
            rows.append([
                Paragraph(str(i), s_td_r),
                Paragraph(str(d["client"])[:36], s_td),
                Paragraph(str(d.get("fonds", ""))[:24], s_td),
                Paragraph((d.get("country") or d.get("region") or "—")[:18], s_td),
                Paragraph(str(d.get("sales_owner", "—"))[:20], s_td),
                Paragraph(_fmt_aum(d["funded_aum"]), s_td_r),
            ])
        tbl = Table(rows, colWidths=col_w, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",     (0,0),(-1,0), MARINE),
            ("LINEBELOW",      (0,0),(-1,0), 1.4, CIEL),
            ("ROWBACKGROUNDS", (0,1),(-1,-1), [LIGHT, BLANC]),
            ("VALIGN",         (0,0),(-1,-1), "MIDDLE"),
            ("TOPPADDING",     (0,0),(-1,-1), 5),
            ("BOTTOMPADDING",  (0,0),(-1,-1), 5),
            ("LEFTPADDING",    (0,0),(-1,-1), 6),
            ("RIGHTPADDING",   (0,0),(-1,-1), 6),
        ]))
        story.append(tbl)
    else:
        story.append(Paragraph("Aucun deal Funded enregistré.", s_body))

    story.append(Spacer(1, 16))
    story.append(Paragraph("Statuts du pipeline", s_h2))
    story.append(Paragraph("Répartition des deals par étape (count).", s_sub))
    statut_rep = k.get("statut_repartition", {})
    if statut_rep:
        order = ["Prospect", "Initial Pitch", "Due Diligence", "Soft Commit", "Funded", "Lost", "Paused", "Redeemed"]
        labels_p = [s for s in order if statut_rep.get(s, 0) > 0]
        vals_p   = [statut_rep[s] for s in labels_p]
        def plot_funnel(ax):
            colors = {"Prospect":"#94A3B8","Initial Pitch":"#6BAED6",
                      "Due Diligence":"#3060A8","Soft Commit":"#1a5e8a",
                      "Funded":"#001c4b","Lost":"#DC2626","Paused":"#A1A1AA","Redeemed":"#7C3AED"}
            bars = ax.barh(range(len(labels_p)), vals_p,
                           color=[colors.get(l, "#001c4b") for l in labels_p],
                           edgecolor="white", height=0.65)
            ax.set_yticks(range(len(labels_p)))
            ax.set_yticklabels(labels_p, fontsize=9, color="#001c4b")
            ax.invert_yaxis()
            ax.set_xlabel("Nombre de deals", fontsize=8, color="#6B7280")
            ax.grid(axis="x", alpha=0.3)
            for i, b in enumerate(bars):
                ax.text(b.get_width() + 0.1, b.get_y() + b.get_height()/2,
                        str(vals_p[i]), va="center", fontsize=8,
                        color="#001c4b", fontweight="bold")
            ax.set_xlim(0, max(vals_p) * 1.18)
        png = _make_chart_png(plot_funnel, w_inch=6.7, h_inch=2.4)
        story.append(RLImage(png, width=USABLE_W, height=USABLE_W * 0.36))
        story.append(Paragraph("Figure 2 — Distribution des deals par statut", s_caption))
    story.append(PageBreak())

    # ── PAGE 3 — SALES PERFORMANCE ─────────────────────────────────────────
    story.append(Paragraph("Performance commerciale", s_h2))
    story.append(Paragraph("AUM Financé et Pipeline actif par membre de l'équipe commerciale.", s_sub))
    if sales:
        labels_s = [s["nom"][:18] for s in sales]
        funded_s = [s["funded_aum"]/1e6 for s in sales]
        pipe_s   = [s["pipeline_aum"]/1e6 for s in sales]
        def plot_sales(ax):
            import numpy as _np
            x = _np.arange(len(labels_s))
            w = 0.35
            ax.bar(x - w/2, funded_s, w, color="#001c4b", label="AUM Funded", edgecolor="white")
            ax.bar(x + w/2, pipe_s,  w, color="#019ee1", label="AUM Pipeline", edgecolor="white")
            ax.set_xticks(x); ax.set_xticklabels(labels_s, fontsize=9, color="#001c4b")
            ax.set_ylabel("AUM (M€)", fontsize=8, color="#6B7280")
            ax.legend(fontsize=8, frameon=False, labelcolor="#001c4b", loc="upper right")
            ax.grid(axis="y", alpha=0.3)
            for i, v in enumerate(funded_s):
                if v > 0: ax.text(i - w/2, v + max(funded_s+pipe_s)*0.01,
                                  "{:.0f}".format(v), ha="center", fontsize=7.5, color="#001c4b")
            for i, v in enumerate(pipe_s):
                if v > 0: ax.text(i + w/2, v + max(funded_s+pipe_s)*0.01,
                                  "{:.0f}".format(v), ha="center", fontsize=7.5, color="#019ee1")
        png = _make_chart_png(plot_sales, w_inch=6.7, h_inch=3.0)
        story.append(RLImage(png, width=USABLE_W, height=USABLE_W * 0.45))
        story.append(Paragraph("Figure 3 — AUM Funded vs Pipeline par commercial (M€)", s_caption))
        story.append(Spacer(1, 10))

        # Sales table
        ratios = [0.32, 0.18, 0.20, 0.16, 0.14]
        col_w  = [USABLE_W * r for r in ratios]
        rows = [[Paragraph(h, s_th_r if i in (1,2,4) else s_th)
                 for i, h in enumerate(["Commercial", "AUM Financé", "Pipeline Actif", "Marché", "Conversion"])]]
        for s in sales:
            rows.append([
                Paragraph(s["nom"][:30], s_td),
                Paragraph(_fmt_aum(s["funded_aum"]), s_td_r),
                Paragraph(_fmt_aum(s["pipeline_aum"]), s_td_r),
                Paragraph(s.get("marche", "")[:18], s_td),
                Paragraph("{} %".format(s["conversion"]), s_td_r),
            ])
        tbl = Table(rows, colWidths=col_w, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",     (0,0),(-1,0), MARINE),
            ("LINEBELOW",      (0,0),(-1,0), 1.4, CIEL),
            ("ROWBACKGROUNDS", (0,1),(-1,-1), [LIGHT, BLANC]),
            ("VALIGN",         (0,0),(-1,-1), "MIDDLE"),
            ("TOPPADDING",     (0,0),(-1,-1), 5),
            ("BOTTOMPADDING",  (0,0),(-1,-1), 5),
            ("LEFTPADDING",    (0,0),(-1,-1), 6),
            ("RIGHTPADDING",   (0,0),(-1,-1), 6),
        ]))
        story.append(tbl)
    else:
        story.append(Paragraph("Équipe commerciale non renseignée.", s_body))
    story.append(PageBreak())

    # ── PAGE 4 — REGIONAL DISTRIBUTION ─────────────────────────────────────
    story.append(Paragraph("Distribution géographique", s_h2))
    story.append(Paragraph("Répartition de l'AUM Financé par région.", s_sub))
    if regions:
        labels_r = [r["region"] for r in regions]
        vals_r   = [r["aum"]    for r in regions]
        def plot_regions(ax):
            colors = ["#001c4b","#1a5e8a","#3060A8","#6BAED6","#94A3B8","#CBD5E1"]
            wedges, _, autotxt = ax.pie(vals_r, labels=labels_r,
                                          colors=colors[:len(labels_r)],
                                          autopct="%1.0f%%", pctdistance=0.78,
                                          wedgeprops={"width": 0.45, "edgecolor": "white", "linewidth": 1.4},
                                          textprops={"fontsize": 8.5, "color":"#001c4b"})
            for at in autotxt:
                at.set_color("white"); at.set_fontweight("bold"); at.set_fontsize(8)
            total = sum(vals_r)
            ax.text(0, 0.05, _fmt_aum(total), ha="center", va="center",
                    fontsize=11, fontweight="bold", color="#001c4b")
            ax.text(0, -0.10, "AUM Financé total", ha="center", va="center",
                    fontsize=7.5, color="#6B7280")
        png = _make_chart_png(plot_regions, w_inch=5.6, h_inch=3.0)
        story.append(RLImage(png, width=USABLE_W * 0.7, height=USABLE_W * 0.4))
        story.append(Paragraph("Figure 4 — Répartition AUM Financé par région", s_caption))
        story.append(Spacer(1, 10))

        # Regions table
        ratios = [0.50, 0.30, 0.20]
        col_w  = [USABLE_W * r for r in ratios]
        total = sum(vals_r) or 1
        rows = [[Paragraph(h, s_th_r if i > 0 else s_th)
                 for i, h in enumerate(["Région", "AUM Financé", "Part"])]]
        for r in regions:
            pct = r["aum"] / total * 100
            rows.append([
                Paragraph(r["region"], s_td),
                Paragraph(_fmt_aum(r["aum"]), s_td_r),
                Paragraph("{:.1f} %".format(pct), s_td_r),
            ])
        tbl = Table(rows, colWidths=col_w, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",     (0,0),(-1,0), MARINE),
            ("LINEBELOW",      (0,0),(-1,0), 1.4, CIEL),
            ("ROWBACKGROUNDS", (0,1),(-1,-1), [LIGHT, BLANC]),
            ("VALIGN",         (0,0),(-1,-1), "MIDDLE"),
            ("TOPPADDING",     (0,0),(-1,-1), 5),
            ("BOTTOMPADDING",  (0,0),(-1,-1), 5),
            ("LEFTPADDING",    (0,0),(-1,-1), 6),
            ("RIGHTPADDING",   (0,0),(-1,-1), 6),
        ]))
        story.append(tbl)
    else:
        story.append(Paragraph("Aucune donnée régionale.", s_body))

    story.append(Spacer(1, 16))
    story.append(Paragraph(
        "<i>Document strictement confidentiel · usage interne exclusif. "
        "Reproduction et diffusion externe interdites. Les performances passées "
        "ne préjugent pas des performances futures.</i>", s_disc))

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)
    return buf.getvalue()


def _build_pptx_report(anon: bool = False) -> bytes:
    """Generate a 3-slide executive PPTX (Cover, KPIs, Top Deals)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(hexs):
        h = hexs.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    MARINE = rgb("#001c4b")
    CIEL   = rgb("#019ee1")
    BLANC  = rgb("#ffffff")
    LIGHT  = rgb("#f4f6fa")
    GREY   = rgb("#6B7280")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, color):
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        return sh

    def text(slide, txt, x, y, w, h, fs, bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(txt)
        r.font.size = Pt(fs); r.font.bold = bold
        if color: r.font.color.rgb = color
        return tb

    def fmt(v):
        v = float(v or 0)
        if v >= 1e9:  return "{:.2f} Md\u20ac".format(v/1e9)
        if v >= 1e6:  return "{:.1f} M\u20ac".format(v/1e6)
        if v >= 1e3:  return "{:.0f} k\u20ac".format(v/1e3)
        return "{:.0f} \u20ac".format(v)

    # ── SLIDE 1 — COVER + KPIs ─────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    rect(s1, 0, 0, 13.33, 2.6, MARINE)
    rect(s1, 0, 2.6, 13.33, 0.06, CIEL)
    text(s1, "MERIDIAN CAPITAL — EXECUTIVE REPORT" + ("  ·  MODE ANONYME" if anon else ""),
         0.5, 0.4, 12.0, 0.4, 11, color=CIEL)
    text(s1, "Pipeline & Performance",
         0.5, 0.85, 12.0, 0.9, 32, bold=True, color=BLANC)
    text(s1, date.today().strftime("%d %B %Y").upper(),
         0.5, 1.85, 6.0, 0.4, 11, color=rgb("#7ab8d8"))

    k = api_kpis()
    kpis = [
        ("AUM Financ\u00e9 Total", fmt(k["total_funded"]),    CIEL),
        ("Pipeline Actif",         fmt(k["pipeline_actif"]),   rgb("#1a5e8a")),
        ("Pipeline Pond\u00e9r\u00e9", fmt(k["weighted_pipeline"]), rgb("#f07d00")),
        ("Conversion",             "{:.1f} %".format(k["taux_conversion"]), rgb("#22a062")),
    ]
    for i, (lbl, val, accent) in enumerate(kpis):
        x = 0.5 + i * 3.15
        rect(s1, x, 3.0, 2.95, 1.7, MARINE)
        rect(s1, x, 3.0, 2.95, 0.07, accent)
        text(s1, lbl, x+0.15, 3.15, 2.7, 0.4,  9, color=rgb("#7ab8d8"))
        text(s1, val, x+0.15, 3.55, 2.7, 1.0, 22, bold=True, color=BLANC)

    text(s1, "Document confidentiel · usage interne exclusif",
         0.5, 7.05, 12.5, 0.3, 8, color=GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 — TOP DEALS ─────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    rect(s2, 0, 0, 13.33, 0.7, MARINE)
    text(s2, "Top Deals — AUM Financ\u00e9", 0.5, 0.18, 11.0, 0.4, 14, bold=True, color=BLANC)
    text(s2, date.today().strftime("%d/%m/%Y"), 11.0, 0.18, 2.0, 0.4, 11,
         color=rgb("#7ab8d8"), align=PP_ALIGN.RIGHT)

    deals = [d for d in api_deals() if d["statut"] == "Funded"][:10]
    if anon:
        deals = [dict(d, client=_anonymize_label(d.get("client"), d.get("client_id", 0))) for d in deals]
    headers = ["#", "Client", "Fonds", "Pays", "Commercial", "AUM Financ\u00e9"]
    col_w = [0.55, 3.6, 2.1, 1.7, 2.0, 2.0]
    col_x = [0.5]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    hy = 1.0
    for hi, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        rect(s2, x, hy, w, 0.45, MARINE)
        text(s2, h, x+0.06, hy+0.07, w-0.12, 0.32, 9, bold=True, color=BLANC)

    if deals:
        for ri, d in enumerate(deals):
            ry = hy + 0.45 + ri * 0.5
            bg = LIGHT if ri % 2 == 0 else BLANC
            row = [
                str(ri + 1),
                (d.get("client") or "")[:32],
                (d.get("fonds") or "")[:22],
                (d.get("country") or d.get("region") or "—")[:16],
                (d.get("sales_owner") or "—")[:20],
                fmt(d.get("funded_aum", 0)),
            ]
            for vi, (val, x, w) in enumerate(zip(row, col_x, col_w)):
                rect(s2, x, ry, w, 0.46, bg)
                clr = CIEL if vi == 5 else MARINE
                fw  = True  if vi == 5 else False
                text(s2, val, x+0.06, ry+0.10, w-0.12, 0.32, 9, bold=fw, color=clr)
    else:
        text(s2, "Aucun deal Funded enregistr\u00e9.",
             0.5, 1.7, 12.0, 0.4, 12, color=GREY)

    # ── SLIDE 3 — SALES PERFORMANCE ─────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    rect(s3, 0, 0, 13.33, 0.7, MARINE)
    text(s3, "Performance commerciale", 0.5, 0.18, 12.0, 0.4, 14, bold=True, color=BLANC)

    sales = api_sales()
    if sales:
        for i, s in enumerate(sales[:6]):
            col = i % 3
            row = i // 3
            x = 0.5 + col * 4.25
            y = 1.2 + row * 2.8
            rect(s3, x, y, 4.0, 2.4, LIGHT)
            rect(s3, x, y, 4.0, 0.08, CIEL)
            text(s3, s["nom"], x+0.2, y+0.15, 3.6, 0.4, 14, bold=True, color=MARINE)
            text(s3, s.get("marche", ""), x+0.2, y+0.55, 3.6, 0.3, 9, color=GREY)
            text(s3, "AUM Financ\u00e9", x+0.2, y+0.95, 1.7, 0.25, 8, color=GREY)
            text(s3, fmt(s["funded_aum"]), x+0.2, y+1.20, 1.7, 0.4, 14, bold=True, color=MARINE)
            text(s3, "Pipeline", x+2.05, y+0.95, 1.7, 0.25, 8, color=GREY)
            text(s3, fmt(s["pipeline_aum"]), x+2.05, y+1.20, 1.7, 0.4, 14, bold=True, color=CIEL)
            text(s3, "Conversion : {} %".format(s["conversion"]),
                 x+0.2, y+1.85, 3.6, 0.3, 9, color=MARINE)
    else:
        text(s3, "\u00c9quipe commerciale non renseign\u00e9e.",
             0.5, 1.5, 12.0, 0.4, 12, color=GREY)

    text(s3, "Document confidentiel · usage interne exclusif",
         0.5, 7.05, 12.5, 0.3, 8, color=GREY, align=PP_ALIGN.CENTER)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


@app.get("/api/export/pdf")
def api_export_pdf(anon: int = 0):
    data = _build_pdf_report(anon=bool(anon))
    suffix = "_anon" if anon else ""
    return Response(
        content=data,
        media_type="application/pdf",
        headers={"Content-Disposition":
                 "attachment; filename=meridian_executive{}_{}.pdf".format(suffix, date.today().isoformat())},
    )


@app.get("/api/export/pptx")
def api_export_pptx(anon: int = 0):
    data = _build_pptx_report(anon=bool(anon))
    suffix = "_anon" if anon else ""
    return Response(
        content=data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition":
                 "attachment; filename=meridian_executive{}_{}.pptx".format(suffix, date.today().isoformat())},
    )


@app.post("/api/admin/reset")
def api_admin_reset():
    db.reset_database()
    return {"ok": True}


@app.post("/api/admin/seed-demo")
def api_admin_seed_demo():
    """Populate the DB with demo data matching the Meridian mockup."""
    db.reset_database()
    db.add_sales_member("Sophie Martin",   "GCC")
    db.add_sales_member("Antoine Dubois",  "EMEA")
    db.add_sales_member("Liya Chen",       "APAC")

    p1   = db.add_client("Allianz Global Inv.", "Insurance",     "EMEA",          country="Germany")
    sub1 = db.add_client("Allianz GI Asset",    "Insurance",     "EMEA",          country="Germany",  parent_id=p1)
    sub2 = db.add_client("PIMCO",               "Asset Manager", "EMEA",          country="United States", parent_id=p1)
    c2   = db.add_client("Pictet AM",           "Asset Manager", "EMEA",          country="Switzerland")
    c3   = db.add_client("Norges Bank IM",      "Sovereign",     "Nordics",       country="Norway")
    c4   = db.add_client("GIC Singapore",       "Sovereign",     "APAC",          country="Singapore")
    c5   = db.add_client("Amundi Asset Mgmt",   "Asset Manager", "EMEA",          country="France")
    c6   = db.add_client("ADIA Abu Dhabi",      "Sovereign",     "GCC",           country="United Arab Emirates")
    c7   = db.add_client("CalPERS",             "Pension",       "North America", country="United States")

    db.add_pipeline_entry(p1,   "Global Value",      "Funded",         142_300_000, 142_300_000, 142_300_000, "", "", None, sales_owner="Sophie Martin")
    db.add_pipeline_entry(sub1, "Income Builder",    "Funded",          50_000_000,  50_000_000,  50_000_000, "", "", None, sales_owner="Sophie Martin")
    db.add_pipeline_entry(sub2, "Active ETFs",       "Soft Commit",     30_000_000,  25_000_000,           0, "", "", (date.today() + timedelta(days=14)).isoformat(), sales_owner="Antoine Dubois")
    db.add_pipeline_entry(c2,   "Private Debt",      "Soft Commit",     90_000_000,  87_000_000,           0, "", "", (date.today() + timedelta(days=7)).isoformat(),  sales_owner="Antoine Dubois")
    db.add_pipeline_entry(c3,   "Income Builder",    "Funded",         210_500_000, 210_500_000, 210_500_000, "", "", None, sales_owner="Sophie Martin")
    db.add_pipeline_entry(c4,   "Resilient Equity",  "Due Diligence",   65_000_000,  65_000_000,           0, "", "", (date.today() + timedelta(days=3)).isoformat(),  sales_owner="Liya Chen")
    db.add_pipeline_entry(c5,   "Active ETFs",       "Funded",         330_000_000, 330_000_000, 330_000_000, "", "", None, sales_owner="Antoine Dubois")
    db.add_pipeline_entry(c6,   "Global Value",      "Initial Pitch",  175_000_000,           0,           0, "", "", date.today().isoformat(), sales_owner="Liya Chen")
    db.add_pipeline_entry(c7,   "Income Builder",    "Funded",          95_500_000,  95_500_000,  95_500_000, "", "", None, sales_owner="Sophie Martin")
    db.add_pipeline_entry(c2,   "Global Value",      "Prospect",        40_000_000,           0,           0, "", "", (date.today() + timedelta(days=21)).isoformat(), sales_owner="Antoine Dubois")

    db.add_contact(p1, "Klaus", "Weber", role="CIO",                email="k.weber@allianz.com",   is_primary=True)
    db.add_contact(p1, "Anna",  "Bauer", role="Portfolio Manager",  email="a.bauer@allianz.com")
    db.add_contact(c2, "Claire","Morel", role="CIO",                email="c.morel@pictet.com",    is_primary=True)
    db.add_contact(c3, "Erik",  "Larsen", role="Head of Allocations", email="e.larsen@nbim.no",   is_primary=True)
    db.add_contact(c4, "Wei",   "Zhang", role="CIO",                email="w.zhang@gic.sg",        is_primary=True)

    today = date.today()
    db.add_activity(c3, today.isoformat(),                                 "Présentation Q3 — équipe très engagée sur Income Builder",       "Meeting")
    db.add_activity(c2, (today - timedelta(days=1)).isoformat(),           "Suivi due diligence Private Debt, docs KYC transmis",            "Call")
    db.add_activity(c6, (today - timedelta(days=2)).isoformat(),           "Initial pitch réalisé — intérêt confirmé Global Value",         "Roadshow")
    db.add_activity(c4, (today - timedelta(days=3)).isoformat(),           "Envoi memorandum d'investissement Resilient Equity",             "Email")
    db.add_activity(p1, (today - timedelta(days=5)).isoformat(),           "Review annuelle — renouvellement Global Value confirmé",         "Meeting")
    return {"ok": True, "message": "Demo data seeded"}


# ---------------------------------------------------------------------------
# /api/overview — aggregate boot endpoint
# ---------------------------------------------------------------------------

@app.get("/api/overview")
def api_overview():
    return {
        "kpis":              api_kpis(),
        "top_deals":         api_deals()[:8],
        "all_deals":         api_deals(),
        "statut_counts":     api_kpis()["statut_repartition"],
        "recent_activities": api_activities(),
        "whitespace":        api_whitespace(),
        "regions":           api_regions(),
        "next_actions":      api_next_actions(),
        "sales_team":        api_sales(),
        "mini_stats":        api_mini_stats(),
        "funds_aggregate":   api_funds_aggregate(),
        "monthly_aum":       api_monthly_aum(),
        "aum_by_country":    api_aum_by_country(),
    }


# ---------------------------------------------------------------------------
# Local launcher
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
